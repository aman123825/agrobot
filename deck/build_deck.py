"""Build the AgriRover GrowwXIITB final-evaluation pitch deck (.pptx).

Design goals: a clean, modern, dark "agri-tech" theme with green/lime accents,
soft shadows, illustrated art, icon cards and a native chart -- all generated
deterministically so the file always opens cleanly in PowerPoint.

Run:  python build_deck.py   (run assets.py first, or let main() do it)
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")


def asset(name):
    return os.path.join(A, name)


# ------------------------------------------------------------- colours -------
def C(hexstr):
    return RGBColor.from_string(hexstr)

INK      = C("0A130D")
PANEL    = C("122218")
PANEL2   = C("16301F")
CARD     = C("14281B")
CARD_LN  = C("2C4A34")
GREEN    = C("2E7D32")
GREEN_L  = C("66BB6A")
LIME     = C("C6FF00")
MINT     = C("22E38C")
SKY      = C("2176D2")
CLOUD    = C("E9F0EB")
WHITE    = C("FFFFFF")
MUTE     = C("9DB1A5")
AMBER    = C("FFB300")
CORAL    = C("FF6B6B")
LIGHTCARD= C("F4F8F5")
DARKTXT  = C("14281B")

FONT   = "Segoe UI"
FONT_SB= "Segoe UI Semibold"
FONT_BL= "Segoe UI Black"
FONT_LT= "Segoe UI Light"
MONO   = "Consolas"

EMU_IN = 914400
SW = Inches(13.333)
SH = Inches(7.5)


# ------------------------------------------------------------- helpers -------
def soft_shadow(shape, blur=150000, dist=55000, direction=5400000,
                alpha=45000, color="000000"):
    spPr = shape._element.spPr
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{direction}" '
        'rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )
    spPr.append(parse_xml(xml))


def no_line(shape):
    shape.line.fill.background()


def bg(slide, prs, image):
    pic = slide.shapes.add_picture(asset(image), 0, 0,
                                   width=prs.slide_width,
                                   height=prs.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return pic


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def card(slide, x, y, w, h, fill=CARD, line=CARD_LN, radius=0.06,
         shadow=True):
    sp = rect(slide, x, y, w, h, fill=fill, line=line, line_w=1.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)
    if shadow:
        soft_shadow(sp)
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of run dicts
    {t, size, color, bold, font, italic}. A string paragraph is shorthand."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [{"t": para}]
        for rd in para:
            r = p.add_run()
            r.text = rd.get("t", "")
            f = r.font
            f.size = Pt(rd.get("size", 18))
            f.name = rd.get("font", FONT)
            f.bold = rd.get("bold", False)
            f.italic = rd.get("italic", False)
            f.color.rgb = rd.get("color", CLOUD)
    return tb


def pic(slide, image, x, y, w=None, h=None):
    kw = {}
    if w is not None:
        kw["width"] = Inches(w)
    if h is not None:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(asset(image), Inches(x), Inches(y), **kw)


def icon_badge(slide, x, y, d, icon, fill=None, ring=LIME):
    """Circular badge with a centred icon image."""
    circ = rect(slide, x, y, d, d, fill=(fill or PANEL2), line=ring,
                line_w=1.5, shape=MSO_SHAPE.OVAL)
    soft_shadow(circ, blur=90000, dist=30000, alpha=40000)
    pad = d * 0.22
    pic(slide, f"ic_{icon}.png", x + pad, y + pad, w=d - 2 * pad, h=d - 2 * pad)
    return circ


def pill(slide, x, y, w, h, label, fill=PANEL2, txt_color=LIME,
         size=11, line=CARD_LN):
    sp = rect(slide, x, y, w, h, fill=fill, line=line, line_w=1.0,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(slide, x, y, w, h, [[{"t": label, "size": size, "color": txt_color,
                               "bold": True, "font": FONT_SB}]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return sp


PAGE = {"n": 0}


def header(slide, kicker, title, n=True):
    # accent tab
    rect(slide, 0.62, 0.62, 0.12, 0.62, fill=LIME,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(slide, 0.92, 0.55, 10.5, 0.4,
         [[{"t": kicker.upper(), "size": 12.5, "color": LIME, "bold": True,
            "font": FONT_SB}]])
    text(slide, 0.9, 0.86, 11.6, 0.7,
         [[{"t": title, "size": 30, "color": WHITE, "bold": True,
            "font": FONT_SB}]])
    # wordmark top-right
    text(slide, 10.4, 0.58, 2.5, 0.3,
         [[{"t": "AGRI", "size": 13, "color": WHITE, "bold": True,
            "font": FONT_BL},
           {"t": "ROVER", "size": 13, "color": LIME, "bold": True,
            "font": FONT_BL}]], align=PP_ALIGN.RIGHT)
    if n:
        footer(slide)


def footer(slide):
    PAGE["n"] += 1
    rect(slide, 0.62, 7.06, 12.1, 0.02, fill=CARD_LN)
    text(slide, 0.62, 7.12, 8, 0.3,
         [[{"t": "AgriRover  ·  GrowwXIITB Final Evaluation  ·  IIT Bombay",
            "size": 9.5, "color": MUTE}]])
    text(slide, 11.2, 7.12, 1.5, 0.3,
         [[{"t": f"{PAGE['n']:02d}", "size": 9.5, "color": LIME,
            "bold": True}]], align=PP_ALIGN.RIGHT)


def new_slide(prs, layout_bg="content_bg.png"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, prs, layout_bg)
    return slide


# ----------------------------------------------------------- content ---------
def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ============================================================ 1. TITLE ==
    s = new_slide(prs, "title_bg.png")
    pic(s, "rover.png", 6.75, 2.05, w=6.4)
    pill(s, 0.9, 1.02, 3.5, 0.42, "GROWWXIITB  ·  FINAL PITCH", size=11.5)
    text(s, 0.86, 1.75, 6.6, 2.4,
         [[{"t": "AgriRover", "size": 72, "color": WHITE, "bold": True,
            "font": FONT_BL}],
          [{"t": "The autonomous ", "size": 25, "color": CLOUD},
           {"t": "precision-farming", "size": 25, "color": LIME, "bold": True},
           {"t": " rover", "size": 25, "color": CLOUD}]],
         line_spacing=1.0, space_after=6)
    text(s, 0.9, 3.95, 6.2, 1.2,
         [[{"t": "Sense the soil. Spot the weed. Dose the exact plant.",
            "size": 16, "color": MUTE}],
          [{"t": "A dual-brain robot that brings AI-grade precision to the "
                 "small Indian farm — at a maker-scale cost.",
            "size": 15, "color": MUTE}]], line_spacing=1.08, space_after=4)
    # feature chips
    chips = ["EDGE AI VISION", "PRECISION DOSING", "GPS AUTONOMY",
             "FIELD-HARDENED"]
    cx = 0.9
    for ch in chips:
        w = 0.24 + len(ch) * 0.098
        pill(s, cx, 5.35, w, 0.4, ch, size=10.5)
        cx += w + 0.18
    text(s, 0.9, 6.35, 8, 0.5,
         [[{"t": "Hitanshu Kapadiya   ·   Vivek Kumar Gupta   ·   Shreyash "
                 "Wagh", "size": 13, "color": CLOUD, "bold": True}],
          [{"t": "B.Tech Mechanical Engineering, IIT Bombay  ·  Class of 2029",
            "size": 11, "color": MUTE}]], space_after=2)

    # ========================================================== 2. PROBLEM ==
    s = new_slide(prs)
    header(s, "The problem", "Small farms can't afford precision — so they "
           "overspend and under-yield")
    probs = [
        ("gauge", "Blind input use",
         "Fertiliser & water are applied uniformly by guesswork — over-dosed "
         "here, starved there. India is among the world's top fertiliser "
         "consumers, much of it wasted."),
        ("leaf", "Weeds & disease spread",
         "Manual scouting misses early weeds and leaf disease. By the time "
         "it's visible field-wide, yield is already lost."),
        ("rupee", "Automation is out of reach",
         "Commercial ag-robots cost lakhs to crores — impossible for the "
         "~86% of Indian holdings that are small or marginal (<2 ha)."),
        ("thermo", "Brutal field conditions",
         "40–45 °C heat, dust, dew and irrigation spray destroy fragile "
         "electronics that were never built for the field."),
    ]
    x0, y0, cw, chh, gap = 0.9, 1.95, 5.72, 2.28, 0.34
    for i, (ic, t, body) in enumerate(probs):
        cx = x0 + (i % 2) * (cw + gap)
        cy = y0 + (i // 2) * (chh + 0.28)
        card(s, cx, cy, cw, chh)
        icon_badge(s, cx + 0.32, cy + 0.34, 0.92, ic, ring=CORAL)
        text(s, cx + 1.5, cy + 0.36, cw - 1.75, 0.5,
             [[{"t": t, "size": 18.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 1.5, cy + 0.86, cw - 1.75, chh - 1.0,
             [[{"t": body, "size": 13, "color": MUTE}]], line_spacing=1.08)

    # ========================================================= 3. SOLUTION ==
    s = new_slide(prs)
    header(s, "Our solution", "One rover that senses, decides and acts — plant "
           "by plant")
    pic(s, "rover.png", 7.35, 2.35, w=5.7)
    text(s, 0.9, 1.95, 6.3, 1.2,
         [[{"t": "AgriRover", "size": 22, "color": LIME, "bold": True,
            "font": FONT_SB},
           {"t": "  is a low-cost autonomous ground robot that snake-routes a "
                 "field, reads the soil, sees crops with edge-AI vision, and "
                 "micro-doses only the plants that need it.",
            "size": 15.5, "color": CLOUD}]], line_spacing=1.12)
    steps = [
        ("gauge", "SENSE", "7-in-1 NPK probe, moisture, TDS, pH, temperature + "
         "a camera reading every frame."),
        ("chip", "DECIDE", "On-board AI (YOLOv8n + Coral TPU) classifies weeds, "
         "disease and obstacles in real time."),
        ("target", "ACT", "Aimed spot-spray, sequential micro-dosing and safe "
         "drive — triggered only where needed."),
        ("wave", "REPORT", "Live dashboard, GPS plant-map, black-box log and "
         "instant Telegram field alerts."),
    ]
    y = 3.35
    for i, (ic, t, body) in enumerate(steps):
        cy = y + i * 0.92
        icon_badge(s, 0.95, cy, 0.74, ic)
        text(s, 1.85, cy + 0.02, 1.6, 0.4,
             [[{"t": t, "size": 15, "color": LIME, "bold": True,
                "font": FONT_SB}]])
        text(s, 1.85, cy + 0.36, 5.3, 0.6,
             [[{"t": body, "size": 12.5, "color": MUTE}]], line_spacing=1.05)

    # ===================================================== 4. ARCHITECTURE ==
    s = new_slide(prs)
    header(s, "How it works", "A dual-brain architecture: real-time control "
           "meets edge AI")
    # ESP32 column
    def stack(cx, title, subtitle, ring, items, w=3.55):
        card(s, cx, 1.95, w, 4.7, fill=PANEL, line=CARD_LN)
        rect(s, cx, 1.95, w, 0.9, fill=PANEL2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        text(s, cx + 0.3, 2.06, w - 0.5, 0.4,
             [[{"t": title, "size": 17, "color": ring, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 0.3, 2.44, w - 0.5, 0.3,
             [[{"t": subtitle, "size": 10.5, "color": MUTE, "font": MONO}]])
        yy = 3.06
        for it in items:
            row = rect(s, cx + 0.26, yy, w - 0.52, 0.52, fill=CARD,
                       line=CARD_LN, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                       radius=0.2)
            text(s, cx + 0.44, yy, w - 0.7, 0.52,
                 [[{"t": it, "size": 11.5, "color": CLOUD}]],
                 anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.6
    stack(0.9, "ESP32 · FreeRTOS", "REAL-TIME  ·  DUAL-CORE", LIME, [
        "Core 1 — tank drive (BTS7960 PWM)",
        "Core 0 — sensors · dosing · MQTT",
        "NPK / moisture / TDS / GPS / IMU",
        "Fail-safe relays + E-stop on EN pin",
        "HMAC-signed command link",
    ])
    stack(4.9, "Raspberry Pi 4", "EDGE AI  ·  PYTHON", MINT, [
        "YOLOv8n obstacle + weed detection",
        "MobileNetV2 disease classifier",
        "Coral Edge-TPU acceleration",
        "Pathway stream + Streamlit dashboard",
        "GPS geo-tagging + plant database",
    ])
    # right: link + flow
    card(s, 8.95, 1.95, 3.78, 4.7, fill=PANEL, line=CARD_LN)
    text(s, 9.2, 2.1, 3.3, 0.4,
         [[{"t": "The secure link", "size": 15, "color": WHITE, "bold": True,
            "font": FONT_SB}]])
    flow = [
        ("Pi vision → STOP / RESUME", "authenticated UART"),
        ("Pi → PING heartbeat 2.5 Hz", "dead-man safety"),
        ("ESP32 → telemetry over MQTT", "TLS + topic ACLs"),
        ("Both ↔ isolated WiFi VLAN", "no open ports"),
    ]
    yy = 2.62
    for a, b in flow:
        rect(s, 9.2, yy, 3.3, 0.82, fill=CARD, line=CARD_LN,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        text(s, 9.4, yy + 0.09, 2.95, 0.4,
             [[{"t": a, "size": 12, "color": CLOUD, "bold": True}]])
        text(s, 9.4, yy + 0.44, 2.95, 0.3,
             [[{"t": b, "size": 10, "color": LIME, "font": MONO}]])
        yy += 0.96

    # ====================================================== 5. CAPABILITIES ==
    s = new_slide(prs)
    header(s, "Capabilities", "Six field jobs, one modular platform")
    caps = [
        ("gauge", "Soil intelligence", "7-in-1 RS485 NPK probe + moisture, "
         "TDS & pH — a live nutrient map of the field."),
        ("target", "Precision micro-dosing", "Sequential pump + actuator state "
         "machine doses the exact plant, not the whole row."),
        ("camera", "AI crop vision", "Weed, disease & obstacle detection at "
         "≥15 FPS on the edge — no cloud needed."),
        ("pin", "GPS snake navigation", "Autonomous coverage routing with "
         "EKF fusion and per-plant geo-tagging."),
        ("drop", "Aimed spot-spray", "Pan/tilt nozzle points at the weed from "
         "the camera box — targets weeds at any height."),
        ("robot", "Modular attachments", "Swap grass-cutter, seed-sower, "
         "weeder or sprayer on slide-and-lock rails."),
    ]
    x0, y0, cw, chh, gx, gy = 0.9, 1.95, 3.78, 2.28, 0.28, 0.28
    for i, (ic, t, body) in enumerate(caps):
        cx = x0 + (i % 3) * (cw + gx)
        cy = y0 + (i // 3) * (chh + gy)
        card(s, cx, cy, cw, chh)
        icon_badge(s, cx + 0.3, cy + 0.3, 0.86, ic)
        text(s, cx + 0.3, cy + 1.24, cw - 0.6, 0.4,
             [[{"t": t, "size": 15.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 0.3, cy + 1.64, cw - 0.55, 0.6,
             [[{"t": body, "size": 11.5, "color": MUTE}]], line_spacing=1.05)

    # =========================================================== 6. AI/CV ===
    s = new_slide(prs)
    header(s, "AI & computer vision", "Three trained models running on the "
           "edge — right on the rover")
    models = [
        ("camera", "Obstacle detection", "YOLOv8n · 7 classes",
         "person · vehicle · animal · rock · stump · fence · ditch → instant "
         "STOP / RESUME.", "≥15 FPS on-device"),
        ("leaf", "Weed detection", "YOLOv8n · DeepWeeds",
         "9 weed species (17,509 images) → drives the aimed spot-spray, "
         "cutting herbicide use.", "INT8 Edge-TPU"),
        ("gauge", "Disease classifier", "MobileNetV2 · PlantVillage",
         "38 disease classes across 14 crops (54k+ images) → logged per plant "
         "with GPS.", "Coral-compiled"),
    ]
    x0, cw, gx = 0.9, 3.78, 0.28
    for i, (ic, t, sub, body, tag) in enumerate(models):
        cx = x0 + i * (cw + gx)
        card(s, cx, 1.95, cw, 3.3)
        icon_badge(s, cx + 0.3, 2.25, 0.9, ic)
        text(s, cx + 1.35, 2.28, cw - 1.5, 0.4,
             [[{"t": t, "size": 15.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 1.35, 2.72, cw - 1.5, 0.3,
             [[{"t": sub, "size": 11, "color": LIME, "font": MONO}]])
        text(s, cx + 0.32, 3.4, cw - 0.62, 1.0,
             [[{"t": body, "size": 12.5, "color": MUTE}]], line_spacing=1.1)
        pill(s, cx + 0.32, 4.62, 2.1, 0.42, tag, size=10.5)
    # bottom band
    card(s, 0.9, 5.5, 11.53, 1.15, fill=PANEL, line=CARD_LN)
    facts = [("Colab-trained", "free T4 GPU"), ("Runs offline", "no internet"),
             ("Graceful", "degrades safely"), ("Black-box", "every event logged")]
    for i, (a, b) in enumerate(facts):
        fx = 1.2 + i * 2.9
        text(s, fx, 5.68, 2.7, 0.4,
             [[{"t": a, "size": 15, "color": LIME, "bold": True,
                "font": FONT_SB}]])
        text(s, fx, 6.08, 2.7, 0.4,
             [[{"t": b, "size": 12, "color": MUTE}]])
        if i < 3:
            rect(s, fx + 2.55, 5.66, 0.014, 0.8, fill=CARD_LN)

    # ====================================================== 7. DOSING/SOIL ===
    s = new_slide(prs)
    header(s, "Precision dosing", "Read the soil, then treat the exact plant")
    # left: sensing list
    card(s, 0.9, 1.95, 5.55, 4.7, fill=PANEL, line=CARD_LN)
    text(s, 1.2, 2.15, 5, 0.4,
         [[{"t": "What it senses", "size": 16, "color": WHITE, "bold": True,
            "font": FONT_SB}]])
    sensors = [
        ("N · P · K", "7-in-1 RS485 probe — the three macronutrients"),
        ("Moisture", "capacitive probe → irrigation decisions"),
        ("TDS / EC", "dissolved salts, calibrated to 1413 µS"),
        ("pH", "buffer-calibrated soil acidity"),
        ("Temp / humidity", "DHT22 ambient + DS18B20 soil"),
    ]
    yy = 2.66
    for a, b in sensors:
        rect(s, 1.2, yy, 4.95, 0.66, fill=CARD, line=CARD_LN,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
        text(s, 1.42, yy + 0.06, 1.7, 0.5,
             [[{"t": a, "size": 13, "color": LIME, "bold": True,
                "font": FONT_SB}]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 3.05, yy + 0.06, 3.0, 0.5,
             [[{"t": b, "size": 11.5, "color": MUTE}]],
             anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.76
    # right: dosing sequence
    card(s, 6.7, 1.95, 5.73, 4.7, fill=PANEL, line=CARD_LN)
    text(s, 7.0, 2.15, 5, 0.4,
         [[{"t": "The dosing sequence", "size": 16, "color": WHITE,
            "bold": True, "font": FONT_SB}]])
    seq = [
        ("01", "Pre-soak", "pump primes the line"),
        ("02", "Extend", "linear actuator lowers the applicator"),
        ("03", "Dose", "metered micro-dose to the target plant"),
        ("04", "Retract", "actuator returns; line clears"),
        ("05", "Drive frozen", "wheels stay locked the entire time"),
    ]
    yy = 2.66
    for n, t, b in seq:
        icon = rect(s, 7.0, yy, 0.58, 0.58, fill=PANEL2, line=LIME,
                    line_w=1.4, shape=MSO_SHAPE.OVAL)
        text(s, 7.0, yy, 0.58, 0.58,
             [[{"t": n, "size": 13, "color": LIME, "bold": True}]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 7.75, yy + 0.02, 4.5, 0.35,
             [[{"t": t, "size": 13.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, 7.75, yy + 0.32, 4.5, 0.3,
             [[{"t": b, "size": 11, "color": MUTE}]])
        yy += 0.76
    text(s, 7.0, 6.28, 5.3, 0.3,
         [[{"t": "→ Result: less chemical, healthier crop, lower cost per acre.",
            "size": 12, "color": MINT, "bold": True, "italic": True}]])

    # ====================================================== 8. NAVIGATION ====
    s = new_slide(prs)
    header(s, "Autonomy", "It knows where it is — and where every plant is")
    nav = [
        ("pin", "Snake-route coverage",
         "Zone-based boustrophedon path planning walks the whole field with no "
         "human driver."),
        ("target", "EKF sensor fusion",
         "GPS + wheel encoders + MPU6050 IMU fused for a stable pose, even when "
         "GPS drifts."),
        ("camera", "Vision geo-tagging",
         "Camera geometry turns a bounding box into real lat/long — ~10–20 cm "
         "per-plant precision from a ₹400 GPS."),
        ("shield", "Layered safety stops",
         "AI obstacle, tilt, motor-stall current and a hardware E-stop can all "
         "halt the drive independently."),
    ]
    x0, y0, cw, chh, gx, gy = 0.9, 1.95, 5.72, 2.28, 0.34, 0.28
    for i, (ic, t, body) in enumerate(nav):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (chh + gy)
        card(s, cx, cy, cw, chh)
        icon_badge(s, cx + 0.32, cy + 0.34, 0.92, ic)
        text(s, cx + 1.5, cy + 0.36, cw - 1.75, 0.5,
             [[{"t": t, "size": 17.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 1.5, cy + 0.86, cw - 1.75, chh - 1.0,
             [[{"t": body, "size": 13, "color": MUTE}]], line_spacing=1.1)

    # ==================================================== 9. FIELD-HARDENED ==
    s = new_slide(prs)
    header(s, "Engineered for the field", "10 real-world failure modes — "
           "found, solved, coded")
    rows = [
        ("Weeds grow above the pathway", "Vision-guided pan/tilt aimed spray"),
        ("45 °C heat → thermal / battery fire", "Thermal guardian + LiFePO4 path"),
        ("Analog sensors over-volt the ADC", "3.3 V power + dividers"),
        ("Relays twitch ON at boot", "Pull-ups + firmware fail-safe"),
        ("Comms link drops mid-drive", "Heartbeat dead-man auto-stop"),
        ("Dust · dew · irrigation spray", "Conformal coat + glands + gaskets"),
        ("Wheel / cutter jams (stall)", "Current-sense stall detection → stop"),
        ("L298N under-spec for soil load", "Upgraded to 2× BTS7960 (43 A)"),
    ]
    x0, y0, cw, gx, gy = 0.9, 1.95, 5.72, 0.34, 0.2
    for i, (prob, fix) in enumerate(rows):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (0.98 + gy)
        card(s, cx, cy, cw, 0.98, fill=PANEL, shadow=(i < 2))
        rect(s, cx + 0.22, cy + 0.28, 0.42, 0.42, fill=GREEN, line=None,
             shape=MSO_SHAPE.OVAL)
        text(s, cx + 0.22, cy + 0.25, 0.42, 0.42,
             [[{"t": "✓", "size": 14, "color": WHITE, "bold": True}]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 0.82, cy + 0.14, cw - 1.0, 0.4,
             [[{"t": prob, "size": 12.5, "color": CLOUD, "bold": True}]])
        text(s, cx + 0.82, cy + 0.5, cw - 1.0, 0.4,
             [[{"t": "→ " + fix, "size": 11.5, "color": LIME}]])
    text(s, 0.9, 6.72, 11.5, 0.3,
         [[{"t": "Every fix is already implemented in firmware/software — "
                 "what remains is physical assembly, not code.",
            "size": 11.5, "color": MUTE, "italic": True}]])

    # ==================================================== 10. SECURITY =======
    s = new_slide(prs)
    header(s, "Safety & security", "A robot you can trust in a real field")
    sec = [
        ("shield", "Authenticated commands", "Every drive/dose command is "
         "HMAC-SHA256 signed with a per-counter anti-replay envelope — "
         "unsigned lines are rejected."),
        ("chip", "Encrypted telemetry", "MQTT over TLS (:8883) with pinned CA, "
         "topic ACLs and no anonymous or plaintext listener."),
        ("bolt", "Hardware E-stop", "A physical button on the ESP32 EN pin "
         "halts everything, regardless of software state."),
        ("target", "Secrets stay off git", "Keys live only in gitignored "
         "secrets — production adds ESP32 Flash Encryption + Secure Boot v2."),
    ]
    x0, y0, cw, chh, gx, gy = 0.9, 1.95, 5.72, 2.28, 0.34, 0.28
    for i, (ic, t, body) in enumerate(sec):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (chh + gy)
        card(s, cx, cy, cw, chh)
        icon_badge(s, cx + 0.32, cy + 0.34, 0.92, ic, ring=MINT)
        text(s, cx + 1.5, cy + 0.36, cw - 1.75, 0.5,
             [[{"t": t, "size": 17, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 1.5, cy + 0.9, cw - 1.75, chh - 1.05,
             [[{"t": body, "size": 12.5, "color": MUTE}]], line_spacing=1.1)

    # ==================================================== 11. TECH STACK =====
    s = new_slide(prs)
    header(s, "Technology stack", "Proven, open, and free — nothing locked "
           "behind a licence")
    groups = [
        ("Firmware", LIME, ["C++ / Arduino", "FreeRTOS (dual-core)",
                            "PlatformIO", "Modbus RS485", "PubSubClient MQTT"]),
        ("Edge AI", MINT, ["Python 3.11", "YOLOv8n · TFLite", "MobileNetV2",
                          "Coral Edge-TPU", "OpenCV"]),
        ("Data & UI", SKY, ["Pathway streaming", "Streamlit dashboard",
                           "InfluxDB / CSV", "Telegram alerts", "Matplotlib"]),
        ("Ops & CI", AMBER, ["Mosquitto (TLS)", "systemd services",
                            "GitHub Actions CI", "ruff · clang", "ROS 2 nodes"]),
    ]
    x0, cw, gx = 0.9, 2.78, 0.23
    for i, (title, col, items) in enumerate(groups):
        cx = x0 + i * (cw + gx)
        card(s, cx, 1.95, cw, 4.55, fill=PANEL, line=CARD_LN)
        rect(s, cx, 1.95, cw, 0.12, fill=col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        text(s, cx + 0.28, 2.2, cw - 0.5, 0.4,
             [[{"t": title, "size": 15.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        yy = 2.78
        for it in items:
            rect(s, cx + 0.24, yy, cw - 0.48, 0.56, fill=CARD, line=CARD_LN,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
            text(s, cx + 0.42, yy, cw - 0.6, 0.56,
                 [[{"t": it, "size": 12, "color": CLOUD}]],
                 anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.68
    text(s, 0.9, 6.66, 11.5, 0.3,
         [[{"t": "Validated by continuous integration: C++ syntax + host "
                 "tests, Python compile, ruff lint and notebook checks on "
                 "every push.", "size": 11.5, "color": MUTE, "italic": True}]])

    # ==================================================== 12. COST/CHART =====
    s = new_slide(prs)
    header(s, "Affordable by design", "Precision agriculture for the price of "
           "a smartphone")
    # chart card
    card(s, 0.9, 1.95, 6.6, 4.7, fill=LIGHTCARD, line=None)
    text(s, 1.2, 2.12, 6, 0.4,
         [[{"t": "Build cost by tier  (₹ '000)", "size": 14,
            "color": DARKTXT, "bold": True, "font": FONT_SB}]])
    chart_data = CategoryChartData()
    chart_data.categories = ["Core\n(demo)", "Core + Nav", "Full AI\n(complete)"]
    chart_data.add_series("Cost", (11.0, 16.0, 38.0))
    gframe = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.55),
        Inches(6.3), Inches(3.9), chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 80
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '₹#,##0"k"'
    dl.number_format_is_linked = False
    dl.font.size = Pt(13)
    dl.font.bold = True
    dl.font.color.rgb = DARKTXT
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    ser = plot.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = GREEN
    # colour the "Full AI" bar lime for emphasis
    from pptx.oxml.ns import qn as _qn
    pts_colors = [GREEN, GREEN_L, LIME]
    for idx, col in enumerate(pts_colors):
        pt = ser.points[idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = col
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(11)
    cat_ax.tick_labels.font.color.rgb = DARKTXT
    cat_ax.tick_labels.font.bold = True
    cat_ax.format.line.color.rgb = C("C8D6CC")
    val_ax = chart.value_axis
    val_ax.visible = False
    val_ax.has_major_gridlines = False
    val_ax.minimum_scale = 0
    val_ax.maximum_scale = 45
    # right: stat callouts
    stats = [
        ("₹28k–49k", "complete Full-AI build, incl. Coral TPU & AI camera"),
        ("₹0", "all software — FreeRTOS, TFLite, YOLO, Mosquitto, Streamlit"),
        ("10–100×", "cheaper than commercial ag-robots (lakhs to crores)"),
    ]
    yy = 1.95
    for big, small in stats:
        card(s, 7.75, yy, 4.68, 1.42, fill=PANEL, line=CARD_LN)
        text(s, 8.05, yy + 0.16, 4.2, 0.6,
             [[{"t": big, "size": 30, "color": LIME, "bold": True,
                "font": FONT_BL}]])
        text(s, 8.05, yy + 0.85, 4.2, 0.5,
             [[{"t": small, "size": 12, "color": MUTE}]], line_spacing=1.02)
        yy += 1.56

    # ==================================================== 13. IMPACT/MARKET ==
    s = new_slide(prs)
    header(s, "Impact & opportunity", "Built for the farms that need it most")
    big = [
        ("~55%", "of India's workforce depends on agriculture"),
        ("~86%", "of farm holdings are small / marginal (<2 ha)"),
        ("10–100×", "lower cost opens automation to those farms"),
    ]
    for i, (a, b) in enumerate(big):
        cx = 0.9 + i * 3.95
        card(s, cx, 1.95, 3.7, 1.85, fill=PANEL, line=CARD_LN)
        text(s, cx + 0.3, 2.14, 3.2, 0.8,
             [[{"t": a, "size": 40, "color": LIME, "bold": True,
                "font": FONT_BL}]])
        text(s, cx + 0.3, 3.05, 3.2, 0.7,
             [[{"t": b, "size": 12.5, "color": CLOUD}]], line_spacing=1.05)
    # impact bullets
    card(s, 0.9, 4.15, 11.53, 2.35, fill=PANEL, line=CARD_LN)
    text(s, 1.2, 4.32, 11, 0.4,
         [[{"t": "Why it matters", "size": 15, "color": WHITE, "bold": True,
            "font": FONT_SB}]])
    imp = [
        ("drop", "Less input, more yield", "Targeted dosing cuts fertiliser & "
         "water waste and lifts productivity per acre."),
        ("leaf", "Healthier crops", "Early weed & disease detection stops "
         "field-wide loss before it starts."),
        ("rupee", "Accessible economics", "Maker-scale BOM + free software puts "
         "precision farming within reach."),
        ("robot", "Less drudgery", "Autonomy removes hours of manual scouting, "
         "spraying and weeding."),
    ]
    for i, (ic, t, b) in enumerate(imp):
        cx = 1.2 + (i % 2) * 5.75
        cy = 4.78 + (i // 2) * 0.82
        icon_badge(s, cx, cy, 0.6, ic)
        text(s, cx + 0.76, cy - 0.02, 4.9, 0.35,
             [[{"t": t, "size": 13, "color": LIME, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 0.76, cy + 0.3, 4.9, 0.4,
             [[{"t": b, "size": 11, "color": MUTE}]], line_spacing=1.0)

    # ==================================================== 14. EDGE / MOAT ====
    s = new_slide(prs)
    header(s, "Why AgriRover wins", "The combination is the moat")
    edge = [
        ("chip", "Two brains, done right", "Hard real-time control (FreeRTOS) "
         "AND heavy edge-AI (Coral) — most hobby builds pick only one."),
        ("target", "Per-plant, not per-field", "We act on the individual plant "
         "— aimed spray + micro-dosing — where others blanket-spray."),
        ("shield", "Field-hardened & secure", "10 documented failure modes "
         "solved, plus an authenticated, encrypted command chain."),
        ("rupee", "Radically affordable", "A complete AI build under ₹50k, "
         "entirely on free and open software."),
    ]
    x0, y0, cw, chh, gx, gy = 0.9, 1.95, 5.72, 2.28, 0.34, 0.28
    for i, (ic, t, body) in enumerate(edge):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (chh + gy)
        card(s, cx, cy, cw, chh)
        icon_badge(s, cx + 0.32, cy + 0.34, 0.92, ic)
        text(s, cx + 1.5, cy + 0.36, cw - 1.75, 0.5,
             [[{"t": t, "size": 17, "color": WHITE, "bold": True,
                "font": FONT_SB}]])
        text(s, cx + 1.5, cy + 0.9, cw - 1.75, chh - 1.05,
             [[{"t": body, "size": 12.5, "color": MUTE}]], line_spacing=1.1)

    # ==================================================== 15. STATUS/ROADMAP =
    s = new_slide(prs)
    header(s, "Status & roadmap", "From working software to a field-proven "
           "product")
    # status pills row
    now = [
        ("DONE", GREEN, "Firmware + Pi software, dual-core, secure link"),
        ("DONE", GREEN, "3 AI models trained & Edge-TPU compiled"),
        ("DONE", GREEN, "CI green: C++, Python, lint, notebooks"),
        ("DONE", GREEN, "All 10 field-challenge fixes coded"),
    ]
    text(s, 0.9, 1.9, 6, 0.4,
         [[{"t": "Where we are", "size": 16, "color": LIME, "bold": True,
            "font": FONT_SB}]])
    yy = 2.4
    for tag, col, body in now:
        card(s, 0.9, yy, 5.55, 0.78, fill=PANEL, shadow=False)
        pill(s, 1.1, yy + 0.22, 0.95, 0.36, tag, fill=col, txt_color=WHITE,
             size=10, line=None)
        text(s, 2.2, yy, 4.15, 0.78,
             [[{"t": body, "size": 12, "color": CLOUD}]],
             anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.9
    # next steps
    text(s, 6.85, 1.9, 6, 0.4,
         [[{"t": "What's next", "size": 16, "color": MINT, "bold": True,
            "font": FONT_SB}]])
    nxt = [
        ("Q1", "Physical assembly & wiring of the v2 hardware"),
        ("Q2", "Full field trial on a partner farm + data collection"),
        ("Q3", "LiFePO4 + RTK-GPS upgrade for cm precision"),
        ("Q4", "Pilot units with smallholder farmer groups"),
    ]
    yy = 2.4
    for q, body in nxt:
        card(s, 6.85, yy, 5.55, 0.78, fill=PANEL, shadow=False)
        icon = rect(s, 7.05, yy + 0.19, 0.52, 0.42, fill=PANEL2, line=MINT,
                    line_w=1.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
        text(s, 7.05, yy + 0.17, 0.52, 0.42,
             [[{"t": q, "size": 12, "color": MINT, "bold": True}]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 7.75, yy, 4.5, 0.78,
             [[{"t": body, "size": 12, "color": CLOUD}]],
             anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.9

    # ========================================================= 16. TEAM ======
    s = new_slide(prs)
    header(s, "The team", "Three mechanical engineers, one field robot")
    team = [
        ("Hitanshu Kapadiya", "Mechanical Engineering", "IIT Bombay · 2029",
         "Mechanics, chassis & drivetrain"),
        ("Vivek Kumar Gupta", "Mechanical Engineering", "IIT Bombay · 2029",
         "Systems, firmware & integration"),
        ("Shreyash Wagh", "Mechanical Engineering", "IIT Bombay · 2029",
         "Electronics, sensing & AI"),
    ]
    x0, cw, gx = 0.9, 3.78, 0.28
    for i, (name, deg, inst, role) in enumerate(team):
        cx = x0 + i * (cw + gx)
        card(s, cx, 2.1, cw, 3.9)
        # avatar circle with initials
        av = rect(s, cx + cw / 2 - 0.7, 2.45, 1.4, 1.4, fill=PANEL2,
                  line=LIME, line_w=1.6, shape=MSO_SHAPE.OVAL)
        soft_shadow(av, blur=90000, dist=30000, alpha=40000)
        initials = "".join(p[0] for p in name.split()[:2]).upper()
        text(s, cx + cw / 2 - 0.7, 2.45, 1.4, 1.4,
             [[{"t": initials, "size": 34, "color": LIME, "bold": True,
                "font": FONT_BL}]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 0.2, 4.05, cw - 0.4, 0.4,
             [[{"t": name, "size": 16.5, "color": WHITE, "bold": True,
                "font": FONT_SB}]], align=PP_ALIGN.CENTER)
        text(s, cx + 0.2, 4.5, cw - 0.4, 0.3,
             [[{"t": deg, "size": 12, "color": CLOUD}]],
             align=PP_ALIGN.CENTER)
        text(s, cx + 0.2, 4.8, cw - 0.4, 0.3,
             [[{"t": inst, "size": 11, "color": LIME, "font": MONO}]],
             align=PP_ALIGN.CENTER)
        rect(s, cx + 0.9, 5.2, cw - 1.8, 0.02, fill=CARD_LN)
        text(s, cx + 0.2, 5.34, cw - 0.4, 0.5,
             [[{"t": role, "size": 11.5, "color": MUTE}]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # ========================================================= 17. CLOSE =====
    s = new_slide(prs, "close_bg.png")
    pic(s, "rover.png", 7.4, 3.0, w=5.7)
    pill(s, 0.9, 1.35, 2.2, 0.42, "THANK YOU", size=12)
    text(s, 0.86, 2.05, 7.2, 2.2,
         [[{"t": "Let's make precision", "size": 40, "color": WHITE,
            "bold": True, "font": FONT_BL}],
          [{"t": "farming ", "size": 40, "color": WHITE, "bold": True,
            "font": FONT_BL},
           {"t": "affordable.", "size": 40, "color": LIME, "bold": True,
            "font": FONT_BL}]], line_spacing=1.0, space_after=2)
    text(s, 0.9, 3.95, 6.4, 0.8,
         [[{"t": "AgriRover — an autonomous, AI-guided, field-hardened rover "
                 "that the smallest farm can actually own.",
            "size": 15, "color": MUTE}]], line_spacing=1.12)
    # contact card
    card(s, 0.9, 5.0, 6.2, 1.55, fill=PANEL, line=CARD_LN)
    text(s, 1.2, 5.16, 5.6, 0.4,
         [[{"t": "Team AgriRover · IIT Bombay", "size": 14, "color": WHITE,
            "bold": True, "font": FONT_SB}]])
    text(s, 1.2, 5.56, 5.7, 0.9,
         [[{"t": "Hitanshu Kapadiya  ·  Vivek Kumar Gupta  ·  Shreyash Wagh",
            "size": 11.5, "color": CLOUD}],
          [{"t": "25b2269@iitb.ac.in   ·   GrowwXIITB Final Evaluation 2026",
            "size": 11, "color": LIME, "font": MONO}]], space_after=3)

    out = os.path.join(os.path.dirname(HERE), "AgriRover_GrowwXIITB.pptx")
    prs.save(out)
    return out, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    import assets
    assets.main()
    path, n = build()
    print(f"OK — {n} slides -> {path}")
