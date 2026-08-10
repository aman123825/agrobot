# -*- coding: utf-8 -*-
"""
AgriRover - IDEAS kickoff deck generator
========================================

Builds `deck/AgriRover_IDEAS_Kickoff.pptx` (16:9) from facts that are verifiable
in THIS repository only. Every number on a slide traces to a file listed in
`SOURCES` below. No market claim, accuracy figure or performance result is
invented: README.md explicitly states no field evidence exists yet, so the deck
carries architectural facts plus labelled validation gates instead.

Run:
    python deck/build_ideas_kickoff.py

Requires: python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------- #
#  Provenance - each slide's numbers come from these repository files
# --------------------------------------------------------------------------- #
SOURCES = {
    "arch":   "README.md",
    "pins":   "firmware/include/pins.h",
    "cfg":    "firmware/include/config.h",
    "servo":  "firmware/src/servo.cpp",
    "mech":   "docs/mechanical-layout.md",
    "cad":    "cad/AgriRoverPrototype/AgriRoverPrototype.py",
    "market": "research/AGRIROVER_MARKET_ADOPTION_RESEARCH.md",
    "team":   "IDEAS_team_members.csv",
}

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "AgriRover_IDEAS_Kickoff.pptx")

# --------------------------------------------------------------------------- #
#  Palette - 5 colours total (2 neutrals, 1 surface, 1 primary, 1 caution)
# --------------------------------------------------------------------------- #
BG      = RGBColor(0x0A, 0x0F, 0x24)   # deep navy canvas
SURFACE = RGBColor(0x14, 0x1D, 0x3D)   # raised panel
BLUE    = RGBColor(0x3D, 0x8B, 0xFF)   # electric blue primary
AMBER   = RGBColor(0xF5, 0xA6, 0x23)   # caution / unproven
INK     = RGBColor(0xF2, 0xF5, 0xFF)   # near-white text
MUTED   = RGBColor(0x8E, 0x9B, 0xC4)   # secondary text

FONT_UI   = "Segoe UI"
FONT_MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.72)          # page margin


# --------------------------------------------------------------------------- #
#  Primitives
# --------------------------------------------------------------------------- #
def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1), shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = lw
    return s


def text(slide, x, y, w, h, runs, size=14, color=INK, bold=False, font=FONT_UI,
         align=PP_ALIGN.LEFT, spacing=1.35, anchor=MSO_ANCHOR.TOP, caps=False):
    """runs: str or list of (str, dict-overrides) paragraphs."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        body, over = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("spacing", spacing)
        if over.get("space_before"):
            p.space_before = over["space_before"]
        r = p.add_run()
        r.text = body.upper() if over.get("caps", caps) else body
        f = r.font
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.name = over.get("font", font)
        f.color.rgb = over.get("color", color)
    return box


def bg(slide):
    """Flat navy canvas. The bundled *_bg.png assets are built for a light theme,
    so they are intentionally not used - a solid field keeps the data legible."""
    rect(slide, 0, 0, W, H, fill=BG)
    return slide


def icon(slide, name, x, y, size=Inches(0.34)):
    p = os.path.join(ASSETS, name)
    if os.path.exists(p):
        return slide.shapes.add_picture(p, x, y, width=size, height=size)
    return None


def slide_head(slide, kicker, title, rule=True):
    """Consistent header: small blue kicker, large title, hairline rule."""
    text(slide, M, Inches(0.62), Inches(11.0), Inches(0.3), kicker,
         size=11, color=BLUE, bold=True, caps=True, spacing=1.0)
    text(slide, M, Inches(0.95), Inches(11.4), Inches(0.75), title,
         size=30, color=INK, bold=True, spacing=1.05)
    if rule:
        rect(slide, M, Inches(1.78), Inches(1.5), Pt(3), fill=BLUE)


FOOTER_Y = Inches(6.88)


def source_note(slide, *keys):
    """Footer proving where the slide's numbers came from."""
    files = "   ".join(SOURCES[k] for k in keys)
    box = text(slide, M, FOOTER_Y, Inches(11.9), Inches(0.28),
               "SOURCE   " + files, size=8.5, color=MUTED, font=FONT_MONO, spacing=1.0)
    box.name = "footer-source"
    return box


def page_num(slide, n, total):
    box = text(slide, Inches(12.1), FOOTER_Y, Inches(0.55), Inches(0.28),
               f"{n:02d}/{total:02d}", size=9, color=MUTED, font=FONT_MONO,
               align=PP_ALIGN.RIGHT, spacing=1.0)
    box.name = "footer-page"
    return box


def metric(slide, x, y, w, value, label, unit="", color=BLUE):
    """Data callout: big mono number over a quiet label."""
    h = Inches(1.16)
    rect(slide, x, y, w, h, fill=SURFACE)
    rect(slide, x, y, Pt(3), h, fill=color)
    text(slide, x + Inches(0.24), y + Inches(0.17), w - Inches(0.4), Inches(0.45),
         value + unit, size=23, color=color, bold=True, font=FONT_MONO, spacing=1.0)
    text(slide, x + Inches(0.24), y + Inches(0.68), w - Inches(0.4), Inches(0.36),
         label, size=9.5, color=MUTED, caps=True, spacing=1.15)


def panel(slide, x, y, w, h, title, rows, accent=BLUE, mono_left=True, row_gap=Inches(0.365)):
    """Surface card with a titled key/value table."""
    rect(slide, x, y, w, h, fill=SURFACE)
    rect(slide, x, y, w, Pt(3), fill=accent)
    text(slide, x + Inches(0.28), y + Inches(0.26), w - Inches(0.56), Inches(0.3),
         title, size=11, color=accent, bold=True, caps=True, spacing=1.0)
    cy = y + Inches(0.72)
    for k, v in rows:
        text(slide, x + Inches(0.28), cy, Inches(1.85), Inches(0.3), k,
             size=10, color=MUTED, font=FONT_MONO if mono_left else FONT_UI, spacing=1.15)
        text(slide, x + Inches(2.22), cy, w - Inches(2.5), Inches(0.32), v,
             size=10.5, color=INK, spacing=1.2)
        cy += row_gap


def bullets(slide, x, y, w, items, size=12.5, gap=Inches(0.5), color=INK, marker=BLUE):
    cy = y
    for it in items:
        rect(slide, x, cy + Inches(0.07), Inches(0.075), Inches(0.075), fill=marker,
             shape=MSO_SHAPE.OVAL)
        text(slide, x + Inches(0.26), cy, w - Inches(0.26), gap, it,
             size=size, color=color, spacing=1.35)
        cy += gap
    return cy


# --------------------------------------------------------------------------- #
#  Slides
# --------------------------------------------------------------------------- #
TOTAL = 10


def s01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)

    # rover render anchored right; it is the one signature element
    p = os.path.join(ASSETS, "rover.png")
    if os.path.exists(p):
        s.shapes.add_picture(p, Inches(7.15), Inches(1.62), width=Inches(5.6))

    rect(s, M, Inches(1.62), Inches(1.5), Pt(3), fill=BLUE)
    text(s, M, Inches(1.95), Inches(6.4), Inches(0.34),
         "IDEAS kickoff  ·  IIT Bombay", size=12, color=BLUE, bold=True, caps=True, spacing=1.0)
    text(s, M, Inches(2.42), Inches(6.3), Inches(1.5), "AgriRover",
         size=66, color=INK, bold=True, spacing=0.95)
    text(s, M, Inches(3.5), Inches(6.0), Inches(1.3),
         "A production-intent agricultural scouting rover: dual-controller "
         "firmware, on-device AI inference, and a fully parametric CAD build.",
         size=15, color=MUTED, spacing=1.45)

    for i, (v, l) in enumerate([("2", "controllers"), ("42", "CAD parts"), ("12,137", "lines of code")]):
        metric(s, M + i * Inches(2.06), Inches(4.9), Inches(1.9), v, l)

    text(s, M, Inches(6.42), Inches(7.0), Inches(0.3),
         "Vivek Kumar Gupta   ·   Shreyash Wagh   ·   Hitanshu Kapadiya   ·   Pritish Nandy",
         size=10.5, color=MUTED, font=FONT_MONO, spacing=1.2)
    page_num(s, 1, TOTAL)


def s02_scope(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "what exists today", "A built system, not a concept sketch")

    icon(s, "icw_robot.png", M, Inches(2.15))
    text(s, M + Inches(0.5), Inches(2.16), Inches(5.5), Inches(0.3),
         "Shipped in this repository", size=12.5, color=INK, bold=True, spacing=1.0)
    bullets(s, M, Inches(2.72), Inches(5.5), [
        "ESP32 firmware (PlatformIO + FreeRTOS) for drive, sensing, dosing and MQTT",
        "71 Raspberry Pi Python modules: AI inference, EKF pose fusion, path planning",
        "A ROS 2 package with sensor, drive, AI and mission nodes",
        "Fusion 360 script that generates the whole rover to scale",
        "11 pytest modules covering AI interfaces, navigation, data and stop logic",
    ])

    panel(s, Inches(7.05), Inches(2.15), Inches(5.55), Inches(4.05),
          "repository scale", [
              ("files", "116 C++ / Python source files"),
              ("lines", "12,137 across firmware, pi, ros2, cad"),
              ("docs", "11 design documents in docs/"),
              ("BOM", "110 components + 15 gap-audit items"),
              ("training", "3 Colab notebooks + dataset builder"),
              ("deploy", "systemd, MQTT TLS, Pi hardening"),
          ])
    source_note(s, "arch")
    page_num(s, 2, TOTAL)


def s03_service(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "commercial hypothesis", "Tomato crop-protection, Narayangaon/Junnar")

    text(s, M, Inches(2.12), Inches(11.9), Inches(0.9),
         "Weekly repeatable sentinel sample and trap read  →  count and geotag suspected "
         "problems  →  compare against an economic threshold  →  trigger a full-row scan "
         "only when risk justifies it  →  agronomist-reviewed action  →  record treatment "
         "and pre-harvest interval  →  revisit and verify.",
         size=14, color=INK, spacing=1.5)

    rect(s, M, Inches(3.32), Inches(11.9), Pt(1), fill=SURFACE)

    text(s, M, Inches(3.62), Inches(5.4), Inches(0.3), "Why this cluster",
         size=12.5, color=INK, bold=True, spacing=1.0)
    bullets(s, M, Inches(4.1), Inches(5.4), [
        "2023-24 tomato area: Nashik 22,040 ha, Pune 4,772 ha, Ahmednagar 4,324 ha, Solapur 4,072 ha",
        "KVK Narayangaon provides local agronomy; a tomato FPO aggregates the route",
        "Concentration dense enough to test operating economics, not just technology",
    ], size=11.5, gap=Inches(0.72))

    panel(s, Inches(7.05), Inches(3.62), Inches(5.55), Inches(2.6),
          "scope limits stated in the repo", [
              ("pilot v1", "pesticide application stays manual"),
              ("NPK probe", "not used to prescribe fertilizer"),
              ("500 ml tank", "spot treatment only, not full-field"),
              ("category", "not a universal farm robot"),
          ], accent=AMBER)
    source_note(s, "market", "arch")
    page_num(s, 3, TOTAL)


def s04_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "system architecture", "Two controllers, split by timing guarantee")

    # left controller
    panel(s, M, Inches(2.15), Inches(5.75), Inches(3.05),
          "real-time tier", [
              ("board", "ESP32 DevKit V1 (WROOM-32)"),
              ("runtime", "FreeRTOS, both cores used"),
              ("core 1", "driveTask - motor PWM, obstacle reaction"),
              ("core 0", "sensorTask - NPK, dosing, MQTT publish"),
              ("signalling", "FreeRTOS EventGroup, not volatile bool"),
          ])
    icon(s, "icw_chip.png", Inches(11.9), Inches(2.36))
    panel(s, Inches(7.05), Inches(2.15), Inches(5.55), Inches(3.05),
          "inference tier", [
              ("board", "Raspberry Pi 5 + Hailo-8 AI HAT+"),
              ("throughput", "26 TOPS accelerator"),
              ("fallback", "Pi 4 + Coral USB Edge TPU"),
              ("workload", "obstacle, weed, leaf-disease paths"),
              ("services", "evidence capture, dashboard, OTA"),
          ])

    # link between the tiers
    rect(s, M, Inches(5.5), Inches(11.55), Inches(0.72), fill=SURFACE)
    rect(s, M, Inches(5.5), Pt(3), Inches(0.72), fill=BLUE)
    text(s, M + Inches(0.28), Inches(5.66), Inches(11.0), Inches(0.42),
         "Command link is HMAC-authenticated and encrypted   ·   heartbeat timeout 1500 ms "
         "halts the rover   ·   10 s lockout after repeated auth failures",
         size=11, color=INK, spacing=1.25)
    source_note(s, "arch", "cfg")
    page_num(s, 4, TOTAL)


def s05_drive(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "firmware · actuation", "Drive and servo PWM on separate timers")

    panel(s, M, Inches(2.15), Inches(5.75), Inches(2.5),
          "tank drive - 2x BTS7960 (IBT-2)", [
              ("GPIO", "19 / 21 left, 22 / 23 right"),
              ("direction", "PWM on RPWM fwd, LPWM rev"),
              ("enable", "R_EN + L_EN tied high to 3.3 V"),
              ("PWM", "LEDC 1 kHz, 8-bit, duty 0-255"),
          ])
    panel(s, Inches(7.05), Inches(2.15), Inches(5.55), Inches(2.5),
          "SG90 sweep mount", [
              ("GPIO", "27, LEDC channel 4 → timer 2"),
              ("frame", "50 Hz, 16-bit resolution"),
              ("range", "duty 1638 (0°) → 8192 (180°)"),
              ("startup", "centred at 90°"),
          ])

    text(s, M, Inches(4.95), Inches(11.9), Inches(0.44),
         "Channel 4 is deliberately on timer 2 so the servo's 20 ms frame can never collide "
         "with the 1 kHz drive PWM on timers 0 and 1. Moving from L298N to BTS7960 freed GPIO 32/33.",
         size=11.5, color=MUTED, spacing=1.4)

    # dosing sequence - strictly ordered, never simultaneous
    text(s, M, Inches(5.62), Inches(11.9), Inches(0.3),
         "Sequential dosing state machine — motors held stopped while the probe is in soil",
         size=11, color=BLUE, bold=True, caps=True, spacing=1.0)
    steps = [("pre-soak", "1500 ms"), ("extend", "4000 ms"), ("dwell", "800 ms"),
             ("micro-dose", "1500 ms"), ("retract", "4000 ms")]
    bw = Inches(2.24)
    for i, (name, ms) in enumerate(steps):
        x = M + i * Inches(2.32)
        rect(s, x, Inches(6.05), bw, Inches(0.66), fill=SURFACE)
        rect(s, x, Inches(6.05), Pt(3), Inches(0.66), fill=BLUE)
        text(s, x + Inches(0.18), Inches(6.16), bw - Inches(0.3), Inches(0.22), name,
             size=10.5, color=INK, bold=True, spacing=1.0)
        text(s, x + Inches(0.18), Inches(6.42), bw - Inches(0.3), Inches(0.22), ms,
             size=10, color=MUTED, font=FONT_MONO, spacing=1.0)
    source_note(s, "pins", "servo", "cfg")
    page_num(s, 5, TOTAL)


def s06_mechanical(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "mechanical design", "Double-decker chassis, 320 × 450 mm")

    # deck stack elevation, drawn from the real Z values in the CAD script
    ox, oy = M, Inches(2.25)
    stack_h = Inches(3.9)
    rect(s, ox, oy, Inches(5.4), stack_h, fill=SURFACE)
    # top-to-bottom, matching the Z constants in the Fusion generator
    decks = [
        ("Top & mast — GPS, camera, antenna", "Z 257+", 0.08),
        ("Sun canopy", "Z 254-257", 0.24),
        ("Upper deck — compute, sensors, UI", "Z 158-162", 0.46),
        ("Lower deck — power, drive, fluid", "Z 64-68", 0.66),
        ("Underside — ultrasonics, actuator, NPK probe", "Z -30", 0.86),
    ]
    for label, z, frac in decks:
        y = oy + Emu(int(stack_h * frac))
        rect(s, ox + Inches(0.3), y, Inches(4.8), Pt(2.5), fill=BLUE)
        text(s, ox + Inches(0.3), y + Inches(0.08), Inches(3.6), Inches(0.24), label,
             size=9.5, color=INK, spacing=1.0)
        text(s, ox + Inches(3.95), y + Inches(0.08), Inches(1.15), Inches(0.24), z,
             size=9, color=MUTED, font=FONT_MONO, align=PP_ALIGN.RIGHT, spacing=1.0)

    panel(s, Inches(7.05), Inches(2.25), Inches(5.55), Inches(2.35),
          "Fusion 360 generator", [
              ("input", "docs/mechanical-layout.md is source of truth"),
              ("output", "42 named, colour-coded components"),
              ("tree", "6 assembly groups, browser mirrors the build"),
              ("origin", "chassis front-left corner, +Y toward rear"),
          ])

    text(s, Inches(7.05), Inches(4.9), Inches(5.55), Inches(0.3),
         "Layout rules that drove the stack", size=11.5, color=INK, bold=True, spacing=1.0)
    bullets(s, Inches(7.05), Inches(5.3), Inches(5.55), [
        "Fluid sits below electronics so leaks drain away",
        "MPU6050 at centre of gravity; GPS and LoRa on the mast, clear of BTS7960 RF noise",
        "Battery and tank low and centred for stability on uneven ground",
    ], size=11, gap=Inches(0.52))
    source_note(s, "mech", "cad")
    page_num(s, 6, TOTAL)


def s07_sensing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "sensing & safety", "Layered stop logic and a verified pin map")

    panel(s, M, Inches(2.15), Inches(5.75), Inches(3.2),
          "sensing buses", [
              ("RS485", "NPK probe via MAX485, GPIO 17/16/4"),
              ("analog", "moisture 34, VBAT 35, TDS 36 (ADC1)"),
              ("1-wire", "DHT22 ambient on GPIO 14"),
              ("UART1", "GPS Neo-6M, RX 39 / TX 15 for RTCM"),
              ("relays", "pump 26, actuator 13, polarity 2"),
              ("I2C", "MPU6050, ADS1115, PCF8574, VL53L1X"),
          ])
    panel(s, Inches(7.05), Inches(2.15), Inches(5.55), Inches(3.2),
          "stop layers (design controls)", [
              ("proximity", "VL53L1X time-of-flight"),
              ("vision", "obstacle inference on device"),
              ("attitude", "IMU tilt threshold"),
              ("link", "heartbeat watchdog, 1500 ms"),
              ("motor", "stall detection via ACS712"),
              ("manual", "E-STOP on the upper deck + Telegram /stop"),
          ])
    text(s, M, Inches(5.62), Inches(11.9), Inches(0.5),
         "GPIO 34/35/36/39 are input-only, so every analog sensor is placed on ADC1. "
         "GPIO 15 is a strapping pin whose required boot level matches an idle UART TX line, "
         "which is why GPS TX can safely live there.",
         size=11.5, color=MUTED, spacing=1.4)
    source_note(s, "pins", "arch")
    page_num(s, 7, TOTAL)


def s08_status(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "honest status", "What is built vs. what is proven")

    rect(s, M, Inches(2.15), Inches(5.75), Inches(4.0), fill=SURFACE)
    rect(s, M, Inches(2.15), Inches(5.75), Pt(3), fill=BLUE)
    text(s, M + Inches(0.28), Inches(2.42), Inches(5.2), Inches(0.3),
         "verified in repository", size=11, color=BLUE, bold=True, caps=True, spacing=1.0)
    bullets(s, M + Inches(0.28), Inches(2.88), Inches(5.2), [
        "Firmware, Pi services, ROS 2 nodes, CAD, wiring and BOMs",
        "Simulation that runs with no hardware attached",
        "Automated test suite over AI interfaces and stop logic",
        "Encrypted, authenticated Pi ↔ ESP32 command link",
        "Model OTA and active-learning frame capture",
    ], size=11.5, gap=Inches(0.6))

    rect(s, Inches(7.05), Inches(2.15), Inches(5.55), Inches(4.0), fill=SURFACE)
    rect(s, Inches(7.05), Inches(2.15), Inches(5.55), Pt(3), fill=AMBER)
    text(s, Inches(7.33), Inches(2.42), Inches(5.0), Inches(0.3),
         "not yet established", size=11, color=AMBER, bold=True, caps=True, spacing=1.0)
    bullets(s, Inches(7.33), Inches(2.88), Inches(5.0), [
        "Crop-detection accuracy on local tomato classes",
        "Productive acres per day and route density",
        "Field reliability and ruggedness",
        "Chemical reduction, yield protection, rupee savings",
        "Willingness to pay and payback period",
    ], size=11.5, gap=Inches(0.6), marker=AMBER)

    text(s, M, Inches(6.42), Inches(11.9), Inches(0.34),
         "The README states this explicitly. Nothing in this deck claims a field result.",
         size=11, color=MUTED, spacing=1.2)
    source_note(s, "arch", "market")
    page_num(s, 8, TOTAL)


def s09_gates(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "next milestones", "Validation gates before any paid launch")

    gates = [
        ("01", "Supervised traversal", "Physical drive tests with the layered stop logic active in a real row."),
        ("02", "Local tomato data", "Collect and label local classes; held-out per-class validation."),
        ("03", "Hailo export path", "Compile and benchmark the model chain on the Hailo-8 accelerator."),
        ("04", "Row navigation", "Camera-based row following and boundary behaviour beyond GPS geotags."),
        ("05", "12-grower pilot", "Paired pilot with KVK / agronomist review and calibrated measurement."),
    ]
    cw = Inches(2.24)
    for i, (n, title, body) in enumerate(gates):
        x = M + i * Inches(2.32)
        rect(s, x, Inches(2.3), cw, Inches(3.5), fill=SURFACE)
        rect(s, x, Inches(2.3), cw, Pt(3), fill=BLUE)
        text(s, x + Inches(0.22), Inches(2.6), cw - Inches(0.44), Inches(0.3), n,
             size=15, color=BLUE, bold=True, font=FONT_MONO, spacing=1.0)
        text(s, x + Inches(0.22), Inches(3.05), cw - Inches(0.44), Inches(0.6), title,
             size=12.5, color=INK, bold=True, spacing=1.15)
        text(s, x + Inches(0.22), Inches(3.85), cw - Inches(0.44), Inches(1.7), body,
             size=10.5, color=MUTED, spacing=1.35)

    text(s, M, Inches(6.15), Inches(11.9), Inches(0.4),
         "Repo analysis also flags open engineering items to close first: a 45× flow-configuration "
         "conflict, and a savings tracker that currently mixes incompatible units.",
         size=11, color=AMBER, spacing=1.35)
    source_note(s, "market", "arch")
    page_num(s, 9, TOTAL)


def s10_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    slide_head(s, "team", "IIT Bombay · Mechanical Engineering")

    members = [
        ("Vivek Kumar Gupta", "25B2269", "B.Tech 2029 · Uttar Pradesh"),
        ("Shreyash Wagh", "25B2227", "B.Tech 2029 · Nashik, Maharashtra"),
        ("Hitanshu Kapadiya", "25B2228", "B.Tech 2029 · Gujarat"),
        ("Pritish Nandy", "—", "Team member"),
    ]
    cw = Inches(2.85)
    for i, (name, roll, meta) in enumerate(members):
        x = M + i * Inches(2.95)
        rect(s, x, Inches(2.4), cw, Inches(3.15), fill=SURFACE)
        # circular photo placeholder
        ph = rect(s, x + Inches(0.92), Inches(2.72), Inches(1.0), Inches(1.0),
                  fill=BG, line=BLUE, lw=Pt(1.5), shape=MSO_SHAPE.OVAL)
        text(s, x + Inches(0.18), Inches(3.95), cw - Inches(0.36), Inches(0.5), name,
             size=13, color=INK, bold=True, align=PP_ALIGN.CENTER, spacing=1.15)
        text(s, x + Inches(0.18), Inches(4.52), cw - Inches(0.36), Inches(0.26), roll,
             size=10.5, color=BLUE, font=FONT_MONO, align=PP_ALIGN.CENTER, spacing=1.0)
        text(s, x + Inches(0.18), Inches(4.85), cw - Inches(0.36), Inches(0.5), meta,
             size=10, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.25)

    text(s, M, Inches(5.95), Inches(11.9), Inches(0.34),
         "github.com/aman123825/agrobot", size=12, color=BLUE, bold=True,
         font=FONT_MONO, spacing=1.0)
    text(s, M, Inches(6.34), Inches(11.9), Inches(0.34),
         "Roll numbers and details for the first three are from IDEAS_team_members.csv; "
         "the repository holds no record for Pritish Nandy, so no details are asserted.",
         size=9.5, color=MUTED, spacing=1.25)
    page_num(s, 10, TOTAL)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for fn in (s01_title, s02_scope, s03_service, s04_architecture, s05_drive,
               s06_mechanical, s07_sensing, s08_status, s09_gates, s10_team):
        fn(prs)

    prs.save(OUT)
    print(f"[deck] wrote {OUT}")
    print(f"[deck] {len(prs.slides._sldIdLst)} slides, 16:9")


if __name__ == "__main__":
    main()
