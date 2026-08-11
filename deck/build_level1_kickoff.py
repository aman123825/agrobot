from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "AgriRover_IDEAS_Level_1_Kickoff_14_Aug_2026.pptx"
HERO = ROOT / "_ppt_build/assets/v3/s1_Image_0.png"

# Apple-keynote direction: large editorial type, calm whitespace, one photographic moment.
INK = RGBColor(19, 22, 20)
FOREST = RGBColor(24, 61, 43)
GREEN = RGBColor(45, 132, 78)
LIME = RGBColor(188, 224, 91)
MIST = RGBColor(242, 245, 240)
LINE = RGBColor(214, 220, 214)
MUTED = RGBColor(92, 101, 95)
WHITE = RGBColor(255, 255, 255)
DISPLAY = "Aptos Display"
BODY = "Aptos"

# "Angry Birds" is the team; "AgriRover" is the project the team is working on.
# Keeping the two in separate constants stops the deck from presenting the
# product name where the template asks for the team name. Slides set both in
# caps; the document metadata carries the proper mixed-case spellings, because
# .title() would flatten the internal capital in "AgriRover".
TEAM = "ANGRY BIRDS"
PROJECT = "AGRIROVER"
TEAM_MIXED = "Angry Birds"
PROJECT_MIXED = "AgriRover"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
prs.core_properties.title = f"{PROJECT_MIXED} — IDEAS IIT Bombay Level 1 Kickoff"
prs.core_properties.subject = "Customer discovery: customer, pain, payer and value hypothesis"
prs.core_properties.author = f"Team {TEAM_MIXED}, IIT Bombay"
prs.core_properties.comments = "14 August 2026. Secondary evidence is identified; no prototype or field validation claim is made."
blank = prs.slide_layouts[6]


def shape(slide, x, y, w, h, fill, radius=False, line=None, transparency=0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.fill.transparency = transparency
    s.line.color.rgb = line or fill; s.line.width = Pt(.8)
    if radius:
        try: s.adjustments[0] = .08
        except Exception: pass
    return s


def txt(slide, x, y, w, h, value, size=16, color=INK, bold=False, font=BODY,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(0)
    r = p.add_run(); r.text = value
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return box


def transition(slide, duration=650):
    el = parse_xml(f'<p:transition {nsdecls("p")} spd="med" advClick="1" dur="{duration}"><p:fade/></p:transition>')
    slide._element.insert(2, el)


def base(slide, dark=False):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = INK if dark else WHITE
    transition(slide)


def rule(slide, x, y, w, color=LINE, thick=.012):
    shape(slide, x, y, w, thick, color)


def header(slide, number, title, dark=False):
    color = WHITE if dark else INK
    txt(slide, .62, .32, 4.2, .3, "IDEAS L1C19 — KICKOFF", 11, color, True)
    txt(slide, 8.5, .32, 4.2, .3, title.upper(), 11, color, True, align=PP_ALIGN.RIGHT)


def source(slide, value, dark=False):
    rule(slide, .62, 6.82, 12.08, RGBColor(62, 68, 64) if dark else LINE)
    # The team credit is wider than the project name was, so the source note
    # gives up the matching amount of width rather than being overprinted.
    txt(slide, .62, 6.93, 9.7, .28, value, 9, RGBColor(173, 180, 175) if dark else MUTED)
    txt(slide, 10.45, 6.93, 2.25, .28, f"TEAM: {TEAM}", 9, WHITE if dark else INK, True, align=PP_ALIGN.RIGHT)


def tag(slide, x, y, label, dark=False, width=1.55):
    fill = LIME if dark else FOREST
    fg = INK if dark else WHITE
    shape(slide, x, y, width, .34, fill, True)
    txt(slide, x, y+.01, width, .23, label, 8.5, fg, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def bullet(slide, x, y, w, title, body, dark=False):
    fg = WHITE if dark else INK; sub = RGBColor(188,195,190) if dark else MUTED
    shape(slide, x, y+.06, .11, .11, LIME if dark else GREEN, True)
    txt(slide, x+.28, y, w-.28, .32, title, 15, fg, True)
    txt(slide, x+.28, y+.36, w-.28, .58, body, 13, sub)

# 1 — cinematic cover
s = prs.slides.add_slide(blank); base(s, True)
s.shapes.add_picture(str(HERO), Inches(7.55), Inches(0), width=Inches(5.78), height=Inches(7.5))
# dark veil makes the image feel photographic rather than illustrative
shape(s, 7.55, 0, 5.78, 7.5, INK, transparency=42)
shape(s, 7.53, 0, .04, 7.5, LIME)
header(s, 1, "Level 1 kickoff", True)
tag(s, .64, 1.02, "14 AUGUST 2026", True, 1.72)
txt(s, .62, 1.72, 6.3, .72, PROJECT, 46, WHITE, True, DISPLAY)
txt(s, .62, 2.55, 6.25, 1.55, "Before we build,\nwe need to know.", 38, WHITE, True, DISPLAY)
txt(s, .65, 4.42, 5.95, .9, "Is there a painful, frequent and paid problem in agricultural compliance evidence?", 18, RGBColor(221,226,222), True)
txt(s, .65, 5.55, 5.95, .34, f"TEAM NAME: {TEAM}", 14, WHITE, True)
txt(s, .65, 5.96, 5.95, .34, "DATE: 14 AUGUST 2026", 14, LIME, True)
txt(s, .65, 6.37, 5.95, .34, "LEVEL 1 · CUSTOMER DISCOVERY", 13, WHITE, True)

# 2 — problem, an argument not a dashboard
s = prs.slides.add_slide(blank); base(s); header(s, 2, "Problem & Customer Segment")
tag(s, .62, .86, "OUR STARTING THESIS", False, 1.75)
txt(s, .62, 1.38, 9.5, .68, "The record may matter more than the robot.", 31, INK, True, DISPLAY)
txt(s, .62, 2.08, 10.9, .58, "Export-linked agriculture requires credible treatment and PHI evidence. Today, that evidence may be fragmented, late or expensive to verify.", 16, MUTED)
# main customer and consequence
shape(s, .62, 2.92, 4.05, 2.87, MIST, True, LINE)
txt(s, .94, 3.24, 3.4, .25, "BEACHHEAD CUSTOMER HYPOTHESIS", 8.5, GREEN, True)
txt(s, .94, 3.72, 3.35, .85, "Exporter / pack-house\nquality team", 23, INK, True, DISPLAY)
txt(s, .94, 4.82, 3.35, .55, "Maharashtra grape and pomegranate clusters", 12, MUTED)
shape(s, 4.97, 2.92, 3.42, 2.87, FOREST, True, FOREST)
txt(s, 5.28, 3.24, 2.78, .25, "THE CONSEQUENCE TO TEST", 8.5, LIME, True)
txt(s, 5.28, 3.72, 2.72, .85, "Audit friction.\nRework. Rejection risk.", 22, WHITE, True, DISPLAY)
txt(s, 5.28, 4.84, 2.68, .58, "We have secondary signals—no first-hand cost baseline yet.", 11, RGBColor(201,211,205))
# four decisive questions
shape(s, 8.7, 2.92, 4.0, 2.87, WHITE, True, LINE)
txt(s, 9.02, 3.24, 3.35, .25, "FOUR QUESTIONS DECIDE THE THESIS", 8.5, GREEN, True)
for i, line in enumerate(["Is the problem frequent and costly?", "Who owns the loss and budget?", "What evidence is trusted today?", "Will the payer name a ₹ value?"]):
    y=3.72+i*.47
    txt(s, 9.02, y, .28, .28, f"{i+1}", 12, GREEN, True)
    txt(s, 9.42, y-.02, 2.85, .38, line, 13.5, INK, True)
source(s, "Sources: APEDA GrapeNet / Residue Monitoring workflow; repository evidence [E05–E06], [F17], [G03]. Status: secondary evidence only; 0 AgriRover interviews completed.")

# 3 — solution as hypothesis with explicit boundary
s = prs.slides.add_slide(blank); base(s); header(s, 3, "Proposed solution")
tag(s, .62, .86, "HYPOTHESIS · NOT BUILT", False, 1.78)
txt(s, .62, 1.38, 10.7, .68, "One trusted evidence record.", 31, INK, True, DISPLAY)
txt(s, .62, 2.05, 11.4, .58, "A geotagged, time-stamped treatment + PHI record, reviewed by an agronomist and usable by the buyer’s QA workflow.", 16, MUTED)
# keynote sequence
sequence=[("01", "Observe", "Capture what happened in the field."), ("02", "Verify", "Attach place, time and accountable review."), ("03", "Record", "Map evidence to the buyer’s workflow."), ("04", "Decide", "Test whether the buyer will pay.")]
for i,(n,h,b) in enumerate(sequence):
    x=.62+i*3.03
    txt(s,x,3.0,.5,.3,n,12,GREEN,True)
    rule(s,x,3.42,2.65,FOREST if i==3 else LINE,.025)
    txt(s,x,3.68,2.56,.4,h,19,INK,True,DISPLAY)
    txt(s,x,4.18,2.55,.72,b,14,MUTED)
# boundary statement
shape(s, .62, 5.17, 12.08, 1.15, INK, True, INK)
txt(s, .94, 5.48, 2.05, .26, "LEVEL 1 BOUNDARY", 9, LIME, True)
txt(s, 3.1, 5.4, 9.1, .48, "Interview first. Quantify pain. Identify payer. Seek real advancement. Only then decide go, pivot or stop.", 14, WHITE, True)
source(s, "No prototype, accuracy, coverage, savings or payback claim is made. The solution is an early value hypothesis to be tested through H1–H8 interviews.")

# 4 — why team / why now; photo-ready without fake portraits
s = prs.slides.add_slide(blank); base(s); header(s, 4, "Why us · why now")
tag(s, .62, .86, "TEAM + TIMING", False, 1.45)
txt(s, .62, 1.38, 11.2, .98, "A team with agricultural context—and access to people who can challenge our assumptions.", 27, INK, True, DISPLAY)
# members in clean editorial bands; photo boxes remain editable for later replacement
people=[
    ("PRITISH NANDY", "Team Head · coordinates Level 1 customer-discovery execution"),
    ("HITANSHU KAPADIYA", "Customer discovery · wants to reduce field-to-audit friction"),
    ("VIVEK KUMAR GUPTA", "From Jaunpur, Uttar Pradesh · agriculture-family background"),
    ("SHREYASH WAGH", "Nashik context · connected to real horticulture workflows"),
]
for i,(name,role) in enumerate(people):
    y=2.72+i*.82
    ph=shape(s,.62,y,.62,.62,MIST,True,GREEN)
    ph.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    txt(s,.62,y+.18,.62,.22,"PHOTO",7,GREEN,True,align=PP_ALIGN.CENTER)
    txt(s,1.46,y,4.15,.3,name,14,INK,True)
    txt(s,1.46,y+.32,4.25,.34,role,11.5,MUTED)
# right-side discovery-access case
shape(s, 6.02, 2.56, 6.68, 3.62, FOREST, True, FOREST)
txt(s, 6.38, 2.91, 5.8, .3, "CUSTOMER-DISCOVERY ACCESS", 11, LIME, True)
bullet(s, 6.38, 3.42, 5.72, "Tinkerers’ Lab, IIT Bombay", "We are in contact with Ashwini Ma’am to seek guidance and relevant discovery connections; this is access, not validation.", True)
bullet(s, 6.38, 4.42, 5.72, "Agricultural context in the team", "Vivek’s Jaunpur agriculture background and Shreyash’s Nashik context can help recruit grounded conversations.", True)
bullet(s, 6.38, 5.42, 5.72, "Online evidence guides—not replaces—interviews", "APEDA traceability, PHI and market evidence will shape questions; customer statements will test the thesis.", True)
source(s, f"Team/access details: supplied by Team {TEAM_MIXED}. Public discovery context: APEDA traceability and PHI guidance; repository evidence register. No endorsement or validation is implied.")

# 5 — preserve the template's required closing slide
s = prs.slides.add_slide(blank); base(s, True); header(s, 5, "Thank You", True)
shape(s, .62, 1.18, .12, 4.85, LIME)
txt(s, 1.08, 1.48, 10.9, 1.08, "THANK YOU", 52, WHITE, True, DISPLAY)
txt(s, 1.1, 2.78, 9.8, .58, "Level 1 begins with listening—not building.", 24, LIME, True, DISPLAY)
txt(s, 1.1, 3.72, 10.6, .55, "Next: test the customer, pain, payer and value hypotheses through structured discovery.", 17, RGBColor(210,218,212))
txt(s, 1.1, 5.2, 5.0, .35, f"TEAM NAME: {TEAM}", 15, WHITE, True)
txt(s, 1.1, 5.65, 5.0, .35, "14 AUGUST 2026", 15, RGBColor(183,192,186))
source(s, "IDEAS L1C19 — Kickoff · Desai Sethi School of Entrepreneurship, IIT Bombay", True)

prs.save(OUT)
print(OUT)
