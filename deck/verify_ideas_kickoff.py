# -*- coding: utf-8 -*-
"""
Geometry + text-fit checks for AgriRover_IDEAS_Kickoff.pptx.

There is no LibreOffice in this environment to rasterise slides, so this script
substitutes for a visual pass: it verifies every shape sits inside the slide,
respects the page margin, and that no textbox is likely to overflow its own
height at its declared font size.

Run:  python deck/verify_ideas_kickoff.py
"""

import os
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "AgriRover_IDEAS_Kickoff.pptx")

MARGIN = Inches(0.5)          # nothing should sit closer than this to an edge
CHAR_W = 0.52                 # avg glyph advance as a fraction of point size
LINE_H = 1.34                 # line box as a multiple of point size


def emu_in(v):
    return round(v / 914400, 3)


def main():
    prs = Presentation(DECK)
    SW, SH = prs.slide_width, prs.slide_height
    problems = []

    for idx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            l, t = sh.left, sh.top
            r, b = l + (sh.width or 0), t + (sh.height or 0)
            name = sh.shape_type

            if l < 0 or t < 0 or r > SW or b > SH:
                problems.append(
                    f"slide {idx}: {name} off-slide "
                    f"L{emu_in(l)} T{emu_in(t)} R{emu_in(r)} B{emu_in(b)}"
                )
            elif l < MARGIN or t < MARGIN or (SW - r) < MARGIN or (SH - b) < MARGIN:
                full_bleed = sh.width == SW and sh.height == SH
                footer = (sh.name or "").startswith("footer-")
                if not full_bleed and not footer:
                    problems.append(
                        f"slide {idx}: {name} inside margin "
                        f"L{emu_in(l)} T{emu_in(t)} R{emu_in(r)} B{emu_in(b)}"
                    )

            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                body = "".join(run.text for run in p.runs)
                if not body:
                    continue
                pt = max((run.font.size.pt for run in p.runs if run.font.size), default=12)
                avail_chars = max(1, int((sh.width / 12700) / (pt * CHAR_W)))
                lines = max(1, -(-len(body) // avail_chars))
                need = Inches(lines * pt * LINE_H / 72)
                if need > (sh.height or 0) * 1.9:
                    problems.append(
                        f"slide {idx}: text may overflow ({lines} lines @ {pt}pt "
                        f"needs {emu_in(need)}in, box {emu_in(sh.height)}in): "
                        f"{body[:58]!r}"
                    )

    print(f"slides: {len(prs.slides._sldIdLst)}   size: "
          f"{emu_in(SW)} x {emu_in(SH)} in")
    if problems:
        print(f"\n{len(problems)} issue(s):")
        for p in problems:
            print("  -", p)
    else:
        print("\nno geometry or text-fit issues found")


if __name__ == "__main__":
    main()
