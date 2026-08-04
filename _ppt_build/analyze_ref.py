"""Analyze AgriRover_Pitch.pptx: dump text, fonts, colors, layout, images, theme."""
import os
import zipfile
import re
from pptx import Presentation
from pptx.util import Emu

SRC = "AgriRover_Pitch.pptx"
OUTDIR = "_ppt_build/ref_imgs"
os.makedirs(OUTDIR, exist_ok=True)


def col(run):
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            return str(c.rgb)
    except Exception:
        pass
    return "-"


def shp_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and f.fore_color and f.fore_color.rgb is not None:
            return str(f.fore_color.rgb)
    except Exception:
        pass
    return "-"


p = Presentation(SRC)
print("=== AgriRover_Pitch.pptx ===")
print("slides:", len(p.slides._sldIdLst),
      "| size:", round(Emu(p.slide_width).inches, 2), "x",
      round(Emu(p.slide_height).inches, 2))
print("=" * 80)

for i, s in enumerate(p.slides, 1):
    print(f"\n---------- SLIDE {i}  (layout: {s.slide_layout.name}) ----------")
    for sh in s.shapes:
        kind = str(sh.shape_type)
        pos = ""
        try:
            pos = f"@({round(Emu(sh.left).inches,2)},{round(Emu(sh.top).inches,2)}) " \
                  f"{round(Emu(sh.width).inches,2)}x{round(Emu(sh.height).inches,2)}"
        except Exception:
            pass
        fill = shp_fill(sh)
        tag = f"[{sh.shape_type}] {pos} fill={fill}"
        if sh.has_text_frame and sh.text_frame.text.strip():
            print(f"  TEXT {tag}")
            for para in sh.text_frame.paragraphs:
                if not para.runs:
                    continue
                seg = []
                for r in para.runs:
                    sz = r.font.size.pt if r.font.size else "?"
                    b = "B" if r.font.bold else ""
                    nm = r.font.name or "?"
                    seg.append(f'"{r.text}"[{sz}{b} {nm} {col(r)}]')
                print("      ", " ".join(seg))
        elif sh.shape_type == 13:
            print(f"  PICTURE {tag}")
        elif sh.has_table:
            print(f"  TABLE {tag} rows={len(sh.table.rows)} cols={len(sh.table.columns)}")
        else:
            print(f"  SHAPE {tag} name={sh.name}")

# theme colors + fonts
print("\n" + "=" * 80)
print("THEME:")
z = zipfile.ZipFile(SRC)
theme = [n for n in z.namelist() if n.startswith("ppt/theme/theme")][0]
xml = z.read(theme).decode("utf-8", "ignore")
for m in re.findall(r'<a:(dk1|lt1|dk2|lt2|accent[1-6]|hlink)>.*?(?:srgbClr val="([0-9A-Fa-f]{6})"|sysClr[^>]*lastClr="([0-9A-Fa-f]{6})")', xml):
    print("  ", m[0], "#" + (m[1] or m[2]))
fonts = re.findall(r'<a:(?:majorFont|minorFont)>\s*<a:latin typeface="([^"]+)"', xml)
print("  fonts:", fonts)

# images
imgs = [n for n in z.namelist() if n.startswith("ppt/media/")]
print("\nEMBEDDED MEDIA:", len(imgs))
for n in imgs:
    print("  ", n, z.getinfo(n).file_size, "bytes")
    with open(os.path.join(OUTDIR, os.path.basename(n)), "wb") as fh:
        fh.write(z.read(n))
