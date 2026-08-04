# Build AgriRover_Investor_5Slides.pptx — a NEW, clean 5-slide investor deck.
# Brief (from GrowwxIITB Track A interviewer):
#   * 4-5 slides that look like a real investor deck
#   * no grants / subsidies anywhere
#   * problem in simple form, solution in simple terms
#   * 1-2 concrete examples, mechanism in simple words, uses,
#     and how it differs from drones & tractors
# Reuses the illustration art extracted from the Gamma deck (_ppt_build/assets/v3/).
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = "_ppt_build/assets/v3"
OUT = "AgriRover_Investor_5Slides.pptx"

# ---- design system -------------------------------------------------------
INK    = RGBColor(0x1F, 0x2A, 0x22)   # near-black green (headlines)
MUT    = RGBColor(0x5C, 0x6B, 0x5E)   # muted gray-green (body)
ACC    = RGBColor(0x2E, 0x7D, 0x32)   # green accent
ACC_D  = RGBColor(0x1E, 0x4D, 0x36)   # deep green (ask card, numbers)
CARD   = RGBColor(0xF2, 0xF7, 0xF2)   # card fill
CARD2  = RGBColor(0xE7, 0xF1, 0xE7)   # stronger tint (pills, highlight col)
LINE   = RGBColor(0xDD, 0xE8, 0xDD)   # card border
BG     = RGBColor(0xFC, 0xFB, 0xF7)   # warm white background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAYX  = RGBColor(0xA8, 0xB5, 0xA8)   # the "no" dash in the table
FONT   = "Calibri"

prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333 in — same 16:9 as the old deck
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(slide, x, y, w, h, fill=None, line=None, round_=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        shp.adjustments[0] = round_
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of dicts {runs:[(text,size,bold,color[,italic])], align, sb, sa, ls}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if "sb" in p: para.space_before = Pt(p["sb"])
        if "sa" in p: para.space_after = Pt(p["sa"])
        if "ls" in p: para.line_spacing = p["ls"]
        for r in p["runs"]:
            text, size, bold, color = r[0], r[1], r[2], r[3]
            run = para.add_run()
            run.text = text
            f = run.font
            f.name, f.size, f.bold = FONT, Pt(size), bold
            f.color.rgb = color
            if len(r) > 4 and r[4]:
                f.italic = True
    return tb


def kicker(slide, x, y, label, w):
    pill = box(slide, x, y, w, 0.34, fill=CARD2, round_=0.5)
    tf = pill.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    f = run.font
    f.name, f.size, f.bold, f.color.rgb = FONT, Pt(10.5), True, ACC_D


def footer(slide, n):
    txt(slide, 0.7, 7.08, 6.0, 0.3,
        [{"runs": [("AgriRover · IITB–Groww INV.ENT · Track A", 8.5, False, GRAYX)]}])
    txt(slide, 11.9, 7.08, 0.75, 0.3,
        [{"runs": [(f"{n} / 5", 8.5, False, GRAYX)], "align": PP_ALIGN.RIGHT}])


def side_img(slide, fname):
    slide.shapes.add_picture(f"{ASSETS}/{fname}", Inches(8.08), 0, height=Inches(7.5))


def wordmark(slide):
    box(slide, 0.7, 0.58, 0.26, 0.26, fill=ACC, round_=0.35)
    txt(slide, 1.06, 0.56, 3.0, 0.32,
        [{"runs": [("AgriRover", 14, True, INK)]}])


# ===========================================================================
# SLIDE 1 — TITLE
# ===========================================================================
s = new_slide()
side_img(s, "s1_Image_0_sm.png")
wordmark(s)
txt(s, 0.7, 1.95, 7.0, 1.6, [
    {"runs": [("Precision farming", 40, True, INK)], "ls": 1.02},
    {"runs": [("for every small farm.", 40, True, ACC_D)], "ls": 1.02},
])
txt(s, 0.7, 3.5, 6.9, 1.1, [
    {"runs": [("A small self-driving rover that checks the soil, looks at every plant, "
               "and gives each one exactly what it needs — nothing more, nothing wasted.",
               13.5, False, MUT)], "ls": 1.15},
])
pill = box(s, 0.7, 4.85, 2.95, 0.44, fill=ACC_D, round_=0.5)
tf = pill.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = tf.margin_right = Inches(0.1)
tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "IITB–Groww INV.ENT · Track A"
r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(11), True, WHITE
txt(s, 0.7, 6.35, 7.0, 0.7, [
    {"runs": [("Hitanshu Kapadiya · Vivek Gupta · Shreyash Wagh · Pritish Nandy", 11.5, True, INK)], "sa": 2},
    {"runs": [("Team of 4 · IIT Bombay", 10.5, False, MUT)]},
])

# ===========================================================================
# SLIDE 2 — THE PROBLEM
# ===========================================================================
s = new_slide()
side_img(s, "s2_Image_0_sm.png")
kicker(s, 0.7, 0.6, "THE PROBLEM", 1.55)
txt(s, 0.7, 1.08, 7.1, 1.0, [
    {"runs": [("Farmers spend more every year —", 27, True, INK)], "ls": 1.02},
    {"runs": [("and waste much of it.", 27, True, ACC_D)], "ls": 1.02},
])
problems = [
    ("Farming is done blind.",
     "Fertilizer and water are spread evenly by hand, because there is no way to know "
     "what each plant needs. A big share is simply wasted."),
    ("Waste costs real money.",
     "Fertilizer, chemicals and labor get costlier every season — and over-use quietly "
     "damages the soil for years to come."),
    ("Existing machines don't fit.",
     "86% of India's farms are under 2 hectares — too small for tractors and big precision "
     "equipment. A lab soil test takes 2–3 weeks, so almost nobody does one."),
]
y = 2.28
for title, body in problems:
    box(s, 0.7, y, 6.95, 1.28, fill=CARD, line=LINE, round_=0.08)
    txt(s, 0.95, y + 0.13, 6.45, 0.3, [{"runs": [(title, 13.5, True, INK)]}])
    txt(s, 0.95, y + 0.46, 6.45, 0.72, [{"runs": [(body, 11, False, MUT)], "ls": 1.1}])
    y += 1.44
txt(s, 0.7, y + 0.04, 6.95, 0.35, [
    {"runs": [("From our field visits: fertilizer is still thrown by hand — evenly, "
               "everywhere, every season.", 10.5, False, ACC_D, True)]},
])
footer(s, 2)

# ===========================================================================
# SLIDE 3 — THE SOLUTION (simple mechanism included)
# ===========================================================================
s = new_slide()
side_img(s, "s6_Image_0_sm.png")
kicker(s, 0.7, 0.6, "OUR SOLUTION", 1.65)
txt(s, 0.7, 1.08, 7.1, 0.55, [
    {"runs": [("A small rover that farms ", 27, True, INK),
              ("plant by plant.", 27, True, ACC_D)]},
])
txt(s, 0.7, 1.68, 7.0, 0.45, [
    {"runs": [("It drives the crop rows on its own and treats every plant individually — "
               "like a careful farmhand that never gets tired.", 12, False, MUT)], "ls": 1.1},
])
steps = [
    ("1", "CHECK", "A probe dips into the soil and reads its health in seconds — no lab, no waiting."),
    ("2", "SEE", "An AI camera spots weeds and sick plants as the rover drives the rows."),
    ("3", "TREAT", "Drops fertilizer or a tiny spray only where it is needed, plant by plant."),
    ("4", "REPORT", "Sends a simple map to the farmer's phone: what it found, what it did."),
]
positions = [(0.7, 2.35), (4.32, 2.35), (0.7, 3.88), (4.32, 3.88)]
for (num, title, body), (x, y) in zip(steps, positions):
    box(s, x, y, 3.42, 1.38, fill=CARD, line=LINE, round_=0.08)
    chip = box(s, x + 0.18, y + 0.16, 0.34, 0.34, fill=ACC, round_=0.5)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num
    cr.font.name, cr.font.size, cr.font.bold, cr.font.color.rgb = FONT, Pt(12), True, WHITE
    txt(s, x + 0.64, y + 0.2, 2.6, 0.3, [{"runs": [(title, 12.5, True, INK)]}])
    txt(s, x + 0.18, y + 0.62, 3.06, 0.7, [{"runs": [(body, 10.5, False, MUT)], "ls": 1.08}])
box(s, 0.7, 5.48, 7.04, 1.3, fill=CARD2, line=None, round_=0.08)
txt(s, 0.95, 5.62, 6.55, 0.3, [{"runs": [("HOW IT WORKS INSIDE — TWO BRAINS, ALWAYS SAFE", 10.5, True, ACC_D)]}])
txt(s, 0.95, 5.94, 6.55, 0.75, [
    {"runs": [("A thinking brain", 11, True, INK),
              (" sees and decides; an ", 11, False, MUT),
              ("acting brain", 11, True, INK),
              (" drives and doses. If anything fails — camera, AI or signal — the rover "
               "simply stops. It cannot run away or overdose a plant.", 11, False, MUT)], "ls": 1.12},
])
footer(s, 3)

# ===========================================================================
# SLIDE 4 — EXAMPLES + HOW IT'S DIFFERENT (tractor / drone)
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "IN PRACTICE", 1.5)
txt(s, 0.7, 1.0, 11.9, 0.55, [
    {"runs": [("What this means on a ", 27, True, INK),
              ("real farm.", 27, True, ACC_D)]},
])
# ---- left: two examples ----
ex1 = [
    ("Today: ", "a full bag of urea is spread evenly by hand — much of it is wasted."),
    ("With AgriRover: ", "the probe finds the one patch that is truly nitrogen-poor and feeds only there."),
    ("Target: ", "30–50% less fertilizer, same yield — what our first pilots will prove."),
]
ex2 = [
    ("Today: ", "the entire field is sprayed with weedkiller."),
    ("With AgriRover: ", "the camera finds each weed and sprays that spot only — far less chemical, and none of it on the food."),
]
box(s, 0.7, 1.9, 5.7, 2.25, fill=CARD, line=LINE, round_=0.07)
txt(s, 0.95, 2.04, 5.2, 0.3, [{"runs": [("EXAMPLE 1 · FERTILIZER", 11, True, ACC_D)]}])
paras = [{"runs": [(lead, 10.5, True, INK), (rest, 10.5, False, MUT)], "sa": 4, "ls": 1.08}
         for lead, rest in ex1]
txt(s, 0.95, 2.38, 5.2, 1.7, paras)
box(s, 0.7, 4.35, 5.7, 1.85, fill=CARD, line=LINE, round_=0.07)
txt(s, 0.95, 4.49, 5.2, 0.3, [{"runs": [("EXAMPLE 2 · WEEDS", 11, True, ACC_D)]}])
paras = [{"runs": [(lead, 10.5, True, INK), (rest, 10.5, False, MUT)], "sa": 4, "ls": 1.08}
         for lead, rest in ex2]
txt(s, 0.95, 4.83, 5.2, 1.3, paras)
# ---- right: comparison table ----
txt(s, 6.7, 1.9, 5.9, 0.35, [{"runs": [("Not a tractor. Not a drone.", 16, True, INK)]}])
tx, ty = 6.7, 2.42
w_label, w_col, row_h, head_h = 2.45, 1.16, 0.42, 0.38
cols = ["TRACTOR", "DRONE", "AGRIROVER"]
rows = [
    ("Sees every plant",          ["–", "–", "✓"]),
    ("Tests the soil",            ["–", "–", "✓"]),
    ("Treats plant-by-plant",     ["–", "–", "✓"]),
    ("Fits a 1-acre farm",        ["–", "✓", "✓"]),
    ("Works without an operator", ["–", "–", "✓"]),
]
table_h = head_h + row_h * len(rows)
# highlight the AgriRover column first (rows stripe over the other columns only)
box(s, tx + w_label + 2 * w_col, ty, w_col, table_h, fill=CARD2, round_=None)
for i in range(len(rows)):
    if i % 2 == 0:
        box(s, tx, ty + head_h + i * row_h, w_label + 2 * w_col, row_h, fill=CARD)
# header
for j, c in enumerate(cols):
    color = ACC_D if j == 2 else MUT
    txt(s, tx + w_label + j * w_col, ty + 0.06, w_col, 0.3,
        [{"runs": [(c, 9.5, True, color)], "align": PP_ALIGN.CENTER}])
# rows
for i, (label, marks) in enumerate(rows):
    ry = ty + head_h + i * row_h
    txt(s, tx + 0.12, ry, w_label - 0.12, row_h,
        [{"runs": [(label, 10.5, False, INK)]}], anchor=MSO_ANCHOR.MIDDLE)
    for j, m in enumerate(marks):
        color = ACC if m == "✓" else GRAYX
        txt(s, tx + w_label + j * w_col, ry, w_col, row_h,
            [{"runs": [(m, 13, True, color)], "align": PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
txt(s, 6.7, ty + table_h + 0.18, 5.9, 0.6, [
    {"runs": [("Tractor = muscle. Drone = a quick look from the sky. "
               "AgriRover = eyes and hands at every plant.", 11, False, MUT, True)], "ls": 1.1},
])
cost = box(s, 6.7, ty + table_h + 0.82, 5.9, 0.5, fill=CARD2, round_=0.15)
ctf = cost.text_frame
ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
ctf.margin_left = ctf.margin_right = Inches(0.1)
ctf.margin_top = ctf.margin_bottom = 0
cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
cr = cp.add_run()
cr.text = "Tractor ₹3–8 lakh  ·  agri drone ₹3–6 lakh  ·  AgriRover under ₹50,000 to build"
cr.font.name, cr.font.size, cr.font.bold, cr.font.color.rgb = FONT, Pt(10.5), True, ACC_D
footer(s, 4)

# ===========================================================================
# SLIDE 5 — MARKET, MODEL & THE ASK
# ===========================================================================
s = new_slide()
side_img(s, "s13_Image_0_sm.png")
kicker(s, 0.7, 0.55, "MARKET, MODEL & THE ASK", 2.7)
txt(s, 0.7, 1.0, 7.1, 0.55, [
    {"runs": [("A big market. ", 27, True, INK), ("A simple business.", 27, True, ACC_D)]},
])
stats = [
    ("146M", "farm holdings in India"),
    ("86%", "of farms under 2 hectares — our users"),
    ("$2.5Bn", "India agritech market by 2034"),
]
for i, (num, label) in enumerate(stats):
    x = 0.7 + i * 2.42
    txt(s, x, 1.8, 2.3, 0.45, [{"runs": [(num, 24, True, ACC_D)]}])
    txt(s, x, 2.28, 2.25, 0.55, [{"runs": [(label, 10, False, MUT)], "ls": 1.05}])
box(s, 0.7, 3.05, 7.1, 1.42, fill=CARD, line=LINE, round_=0.07)
txt(s, 0.95, 3.18, 6.6, 0.28, [{"runs": [("HOW WE EARN", 10.5, True, ACC_D)]}])
earn = [
    "1.  Sell the rover with a yearly service plan.",
    "2.  Rent it per acre through farmer groups (FPOs) — no upfront cost for the farmer.",
    "3.  Monthly subscription for soil-health and crop reports.",
]
txt(s, 0.95, 3.5, 6.6, 0.95,
    [{"runs": [(e, 11, False, INK)], "sa": 3} for e in earn])
box(s, 0.7, 4.62, 7.1, 1.05, fill=CARD, line=LINE, round_=0.07)
txt(s, 0.95, 4.74, 6.6, 0.28, [{"runs": [("WHERE WE ARE TODAY", 10.5, True, ACC_D)]}])
txt(s, 0.95, 5.05, 6.6, 0.55, [
    {"runs": [("Full software stack — driving, AI vision, navigation, farmer dashboard — built and "
               "tested in simulation. Hardware fully designed and costed: ₹27,000–48,000 per rover.",
               10.5, False, MUT)], "ls": 1.1},
])
ask = box(s, 0.7, 5.85, 7.1, 0.95, fill=ACC_D, round_=0.09)
atf = ask.text_frame
atf.vertical_anchor = MSO_ANCHOR.MIDDLE
atf.margin_left = atf.margin_right = Inches(0.25)
atf.margin_top = atf.margin_bottom = Inches(0.08)
ap = atf.paragraphs[0]
ar = ap.add_run(); ar.text = "The ask:  "
ar.font.name, ar.font.size, ar.font.bold, ar.font.color.rgb = FONT, Pt(12.5), True, WHITE
ar2 = ap.add_run()
ar2.text = ("seed support to build the first prototype, plus a pilot with one farmer "
            "group to prove 30–50% input savings on real farms.")
ar2.font.name, ar2.font.size, ar2.font.bold, ar2.font.color.rgb = FONT, Pt(12.5), False, WHITE
ap.line_spacing = 1.12
txt(s, 0.7, 7.08, 9.0, 0.3, [
    {"runs": [("Hitanshu Kapadiya · Vivek Gupta · Shreyash Wagh · Pritish Nandy · Team of 4, IIT Bombay · Track A",
               8.5, False, GRAYX)]},
])
txt(s, 11.9, 7.08, 0.75, 0.3,
    [{"runs": [("5 / 5", 8.5, False, GRAYX)], "align": PP_ALIGN.RIGHT}])

prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides)} slides")
