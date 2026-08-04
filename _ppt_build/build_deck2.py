"""Build AgriRover_GrowwxIITB.pptx — BUSINESS-FIRST pitch for IITB-Groww INV.ENT
(Track A). Reuses the helper/theme library in build_deck.py, follows the real
pitch's investor narrative + researched numbers, keeps the strong custom
diagrams (inline + engineering appendix). Team of 4, correct event branding.
"""
import os
import build_deck as B
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = os.path.join(B.REPO, "AgriRover_GrowwxIITB.pptx")
EVENT = "IITB–Groww INV.ENT  ·  Track A"

C = B.C
INK, FOREST, FOREST2 = B.INK, B.FOREST, B.FOREST2
GREEN, LEAF, AMBER, SKY, PLUM = B.GREEN, B.LEAF, B.AMBER, B.SKY, B.PLUM
CLOUD, MIST, PANEL, GRAY, LGRAY, WHITE, RED = (B.CLOUD, B.MIST, B.PANEL, B.GRAY,
                                              B.LGRAY, B.WHITE, B.RED)
HF, HF2, BF = B.HEAD_FONT, B.HEAD_FONT2, B.BODY_FONT
MX, CT, CB, CW = B.MX, B.CT, B.CB, B.CW


# ---------------- extra helpers ----------------
def stat_big(s, x, y, w, h, big, desc, color):
    p = B.rect(s, x, y, w, h, fill=PANEL, line=MIST, line_w=1,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    p.adjustments[0] = 0.06
    B.rect(s, x, y, 0.13, h, fill=color)
    _, tf = B.textbox(s, x + 0.34, y + 0.14, w - 0.5, 0.55)
    B.para(tf, big, size=25, color=color, bold=True, font=HF, first=True,
           space_after=0)
    _, tf = B.textbox(s, x + 0.34, y + 0.72, w - 0.55, h - 0.82)
    B.para(tf, desc, size=11.5, color=GRAY, first=True, space_after=0,
           line_spacing=1.12)


def metric_bar(s, x, y, w, label, big, pct, color):
    p = B.rect(s, x, y, w, 1.5, fill=WHITE, line=MIST, line_w=1,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    p.adjustments[0] = 0.07
    B.rect(s, x, y, 0.11, 1.5, fill=color)
    _, tf = B.textbox(s, x + 0.26, y + 0.15, w - 0.45, 0.35)
    B.para(tf, label, size=12.5, color=INK, bold=True, font=HF, first=True,
           space_after=0)
    _, tf = B.textbox(s, x + 0.26, y + 0.46, w - 0.45, 0.55)
    B.para(tf, big, size=27, color=color, bold=True, font=HF, first=True,
           space_after=0)
    B.rect(s, x + 0.28, y + 1.12, w - 0.55, 0.16, fill=MIST,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    B.rect(s, x + 0.28, y + 1.12, (w - 0.55) * pct, 0.16, fill=color,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def side_panel(s, x, y, w, h, sections, fill=PANEL, hcolor=GREEN):
    p = B.rect(s, x, y, w, h, fill=fill, line=MIST, line_w=1,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    p.adjustments[0] = 0.04
    _, tf = B.textbox(s, x + 0.3, y + 0.24, w - 0.55, h - 0.4)
    first = True
    for head, lines in sections:
        B.para(tf, head, size=13.5, color=hcolor, bold=True, font=HF,
               first=first, space_after=3, space_before=(0 if first else 8))
        first = False
        for ln in lines:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(3); p2.line_spacing = 1.06
            r = p2.add_run(); r.text = "▪  "
            r.font.size = Pt(11.5); r.font.color.rgb = C(hcolor); r.font.bold = True
            r = p2.add_run(); r.text = ln
            r.font.size = Pt(11.5); r.font.color.rgb = C(INK); r.font.name = BF


def source_note(s, text):
    _, tf = B.textbox(s, MX, CB - 0.28, CW, 0.3)
    B.para(tf, text, size=9, color=LGRAY, italic=True, first=True, space_after=0)


def banner(s, y, parts, fill=PANEL, h=0.66):
    b = B.rect(s, MX, y, CW, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    b.adjustments[0] = 0.4
    _, tf = B.textbox(s, MX + 0.3, y, CW - 0.6, h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    B.runs(p, parts)


# ---------------- slides ----------------
def t_title():
    s = B.slide()
    bg = B.rect(s, 0, 0, 13.333, 7.5, fill=FOREST); B.grad(bg, FOREST, FOREST2, 55)
    B.rect(s, 0, 0, 0.3, 7.5, fill=LEAF)
    c1 = B.rect(s, 9.7, -2.2, 5.6, 5.6, fill=GREEN, shape=MSO_SHAPE.OVAL); B._translucent(c1, 20000)
    c2 = B.rect(s, 11.6, 3.6, 4.2, 4.2, fill=LEAF, shape=MSO_SHAPE.OVAL); B._translucent(c2, 13000)
    c3 = B.rect(s, 8.7, 4.9, 2.2, 2.2, fill=AMBER, shape=MSO_SHAPE.OVAL); B._translucent(c3, 12000)
    B.brand_wordmark(s, MX, 0.62, h=0.52, on_dark=True)
    _, tf = B.textbox(s, MX, 2.0, 10.5, 0.4)
    B.para(tf, "PRECISION FARMING FOR EVERY SMALL FARM", size=15, color=AMBER,
           bold=True, font=HF, first=True, space_after=0)
    _, tf = B.textbox(s, MX - 0.03, 2.38, 10.8, 1.5)
    p = tf.paragraphs[0]
    B.runs(p, [("Agri", 66, WHITE, True, HF), ("Rover", 66, LEAF, True, HF)])
    _, tf = B.textbox(s, MX, 3.78, 11.0, 0.7)
    B.para(tf, "Affordable autonomous rover — soil-NPK sensing · AI weed & disease "
               "detection · per-plant micro-dosing", size=17, color="D9EAD9",
           bold=True, font=HF2, first=True, space_after=0, line_spacing=1.15)
    B.chip(s, MX, 4.85, EVENT, color=AMBER, tcolor=INK, w=4.5, size=12.5)
    B.rect(s, MX, 5.62, 11.5, 0.02, fill="27503C")
    _, tf = B.textbox(s, MX, 5.78, 11.6, 0.5)
    p = tf.paragraphs[0]
    B.runs(p, [("Team of 4 · IIT Bombay:  ", 12.5, LEAF, True),
               ("Hitanshu Kapadiya · Vivek Gupta · Shreyash Wagh · Pritish Nandy",
                12.5, WHITE, True)])
    _, tf = B.textbox(s, MX, 6.16, 11.6, 0.4)
    B.para(tf, "Mechanical Engineering ×3  ·  Aerospace Engineering ×1",
           size=11.5, color="9FC2AC", first=True, space_after=0)
    B.fade(s)


def t_problem(n):
    s = B.slide(); B.header(s, "Problem & Validation", "Farming runs on guesswork", n, RED)
    _, tf = B.textbox(s, MX, CT, CW, 0.55)
    B.para(tf, "India's small farms can't use big machinery — and fertilise blind. "
               "That wastes money and degrades soil.", size=15, color=GRAY,
           first=True, line_spacing=1.12)
    cy = 2.62; cw2 = (CW - 0.4) / 2; ch = 1.5
    B.card(s, MX, cy, cw2, ch, "Farms too small for machinery",
           ["≈86% of India's holdings are under 2 ha", "(Agriculture Census 2015–16)."],
           accent=RED, icon="▪", title_size=14, body_size=12)
    B.card(s, MX + cw2 + 0.4, cy, cw2, ch, "Fertilise by guesswork",
           ["No affordable way to sense soil NPK or", "dose per plant."],
           accent=AMBER, icon="▪", title_size=14, body_size=12)
    B.card(s, MX, cy + ch + 0.28, cw2, ch, "NPK imbalance vs 4:2:1 (FAI)",
           ["Excess urea wastes money and", "degrades the soil."],
           accent="B4632B", icon="▪", title_size=14, body_size=12)
    B.card(s, MX + cw2 + 0.4, cy + ch + 0.28, cw2, ch, "Soil Health Cards underused",
           ["Lab tests take 2–3 weeks — so old", "habits win."],
           accent=PLUM, icon="▪", title_size=14, body_size=12)
    banner(s, cy + 2 * ch + 0.5,
           [("Validated in the field:  ", 13.5, GREEN, True, HF),
            ("talks with local farmers & input dealers confirm uniform hand-"
             "broadcasting and guesswork dosing is the norm.", 12.5, INK, False)],
           fill="E7F3EC")
    source_note(s, "Sources: Agriculture Census 2015–16; FAI. Indicative.")
    B.fade(s)


def t_solution(n):
    s = B.slide(); B.header(s, "Solution & Innovation", "One low-cost autonomous rover", n)
    _, tf = B.textbox(s, MX, CT, CW, 0.55)
    B.para(tf, "AgriRover senses, decides and acts per plant — precision farming at "
               "a fraction of big-ag machinery cost.", size=15, color=GRAY,
           first=True, line_spacing=1.12)
    cy = 2.6; w4 = (CW - 3 * 0.3) / 4; ch = 2.9
    data = [
        ("SENSE", GREEN, "1", ["7-in-1 NPK probe", "moisture · TDS · pH", "GPS + IMU + encoders"]),
        ("DECIDE", SKY, "2", ["YOLOv8n vision", "weed · disease", "obstacle detection"]),
        ("ACT", AMBER, "3", ["micro-dosing", "aimed spot-spray", "cutter · seeder · weeder"]),
        ("LEARN", PLUM, "4", ["geo-tagged logs", "prescription maps", "dashboard + alerts"]),
    ]
    for i, (t, col, num, lines) in enumerate(data):
        B.card(s, MX + i * (w4 + 0.3), cy, w4, ch, t, lines, accent=col, num=num,
               title_size=16, body_size=12.5)
    banner(s, cy + ch + 0.32,
           [("Why it's different:  ", 13.5, GREEN, True, HF),
            ("the only affordable, autonomous, per-plant solution purpose-built "
             "for India's small farms.", 12.5, INK, False)], fill="E7F3EC")
    B.fade(s)


def t_why(n):
    s = B.slide(); B.header(s, "Why AgriRover?", "Vs the alternatives", n, SKY)
    headers = ["Capability", "AgriRover", "Big-ag\nmachinery", "Manual\nlabour", "Ag-drones"]
    rows = [
        ("Affordable for <2 ha farms", "✓", "✗", "✓*", "~"),
        ("Senses soil NPK per spot", "✓", "✗", "✗", "✗"),
        ("Per-plant micro-dosing", "✓", "✗", "~", "✗"),
        ("AI weed & disease ID", "✓", "✗", "✗", "~"),
        ("Autonomous operation", "✓", "✓", "✗", "~"),
        ("One-time low-cost asset", "✓", "✗", "✗", "~"),
    ]
    x, y, w, h = MX, CT + 0.05, CW, 3.75
    nrow = len(rows) + 1
    tbl_shape = s.shapes.add_table(nrow, 5, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = tbl_shape.table
    tbl.first_row = False; tbl.horz_banding = False
    widths = [4.2, 1.9, 1.8, 1.75, 1.85]
    for j, wd in enumerate(widths):
        tbl.columns[j].width = Inches(wd)
    sym_color = {"✓": GREEN, "✗": "C0473B", "~": AMBER, "✓*": GREEN}
    for i in range(nrow):
        tbl.rows[i].height = Inches(h / nrow)
        row = headers if i == 0 else rows[i - 1]
        for j in range(5):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            val = row[j]
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C(FOREST)
            elif j == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = C("E7F3EC")
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = C(WHITE if i % 2 else PANEL)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0) else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            if i == 0:
                r.font.name = HF; r.font.size = Pt(12); r.font.bold = True
                r.font.color.rgb = C(WHITE)
            elif j == 0:
                r.font.name = BF; r.font.size = Pt(12.5); r.font.bold = True
                r.font.color.rgb = C(INK)
            else:
                r.font.name = HF; r.font.size = Pt(15); r.font.bold = True
                r.font.color.rgb = C(sym_color.get(val, GRAY))
    banner(s, y + h + 0.18,
           [("Our edge:  ", 13, GREEN, True, HF),
            ("affordable + autonomous + per-plant, in one asset that also senses "
             "soil.   ", 12, INK, False),
            ("*manual labour is cheap per use but recurring.", 10.5, GRAY, False)],
           fill="E7F3EC", h=0.6)
    B.fade(s)


def t_arch(n):
    s = B.slide(); B.header(s, "How It Works", "Two brains, one mission", n)
    _, tf = B.textbox(s, MX, CT - 0.02, CW, 0.55)
    B.para(tf, "A real-time ESP32 (FreeRTOS) runs drive, sensing & dosing; a "
               "Raspberry Pi 4 runs AI vision & analytics — over an HMAC-signed link.",
           size=13.5, color=GRAY, first=True, line_spacing=1.1)
    B.place_image(s, "architecture.png", MX, 2.45, CW, 4.35,
                  caption="Dual-controller architecture — real-time control + on-board AI.")
    B.fade(s)


def t_sense_see(n):
    s = B.slide(); B.header(s, "Sense & See", "Know the soil, see the crop", n, GREEN)
    B.bullets(s, MX, CT + 0.05, 6.4, 3.4, [
        ("7-in-1 soil probe", "N·P·K·pH·EC·moisture·temp over RS485 Modbus"),
        ("Calibrated & robust", "oversampled ADC, multi-point moisture, temp-comp TDS"),
        ("Three AI models", "obstacle · weed (DeepWeeds) · disease (PlantVillage)"),
        ("Runs on Coral Edge TPU", "~10× faster than CPU; fails safe if a model is absent"),
    ], size=14, mcolor=GREEN)
    B.place_image(s, "ai_performance.png", MX + 6.7, CT + 0.15, CW - 6.7, 3.3)
    ky = 5.5; kw = (CW - 3 * 0.3) / 4
    B.kpi(s, MX, ky, kw, 1.2, "7-in-1", "soil probe", GREEN)
    B.kpi(s, MX + (kw + 0.3), ky, kw, 1.2, "38", "disease classes", SKY)
    B.kpi(s, MX + 2 * (kw + 0.3), ky, kw, 1.2, "~30", "FPS on Coral", AMBER)
    B.kpi(s, MX + 3 * (kw + 0.3), ky, kw, 1.2, "~10–20cm", "plant geo-tag", PLUM)
    B.fade(s)


def t_act(n):
    s = B.slide(); B.header(s, "Act & Autonomy", "Treat the plant, drive the field", n, AMBER)
    B.bullets(s, MX, CT + 0.05, 6.7, 4.3, [
        ("Sequential micro-dosing", "pre-soak → probe → dose → retract; drive frozen while dosing"),
        ("Vision-aimed spot-spray", "pan/tilt nozzle targets each detected weed (FC-01)"),
        ("Autonomous snake-route", "GPS + SBAS + EKF fusion → sub-meter coverage"),
        ("Closed-loop velocity PID", "encoder feedback → straight rows, repeatable spacing"),
        ("Fail-safe by design", "E-stop, dead-man heartbeat, thermal guardian, stall stop"),
    ], size=14, mcolor=AMBER)
    B.place_image(s, "savings.png", MX + 6.9, CT + 0.15, CW - 6.9, 4.3)
    B.fade(s)


def t_market(n):
    s = B.slide(); B.header(s, "Market Opportunity & Scope", "A large, growing market", n, SKY)
    x = MX; w = 5.3; h = 1.4
    stat_big(s, x, CT + 0.05, w, h, "$2.5 Bn",
             "India agritech market by 2034 (from ~$0.97 Bn in 2025, ~10.6% CAGR).", GREEN)
    stat_big(s, x, CT + 0.05 + h + 0.22, w, h, "$739 Mn",
             "India precision-agriculture market by 2034 (~9.2% CAGR).", SKY)
    stat_big(s, x, CT + 0.05 + 2 * (h + 0.22), w, h, "~146 M",
             "farm holdings in India — 86% small/marginal (<2 ha): our core market.", AMBER)
    side_panel(s, MX + w + 0.4, CT + 0.05, CW - w - 0.4, 4.55, [
        ("Scope — what it does today",
         ["Soil NPK sensing, AI weed/disease ID, targeted micro-dosing, autonomous nav."]),
        ("Scope — where it grows",
         ["Crops: vegetables, cotton, cereals, orchards, polyhouses.",
          "Geography: India's small-farm belts via FPOs / co-ops.",
          "Add-ons: seeding, spraying, weeding, scouting, yield mapping.",
          "Data services: soil & crop analytics + prescription maps (recurring)."]),
    ])
    source_note(s, "Sources: IMARC (agritech & precision-ag), Agriculture Census 2015–16. Indicative.")
    B.fade(s)


def t_govt(n):
    s = B.slide(); B.header(s, "Tailwinds & Business Model", "Policy is on our side", n)
    B.bullets(s, MX, CT + 0.08, 6.35, 4.4, [
        ("SMAM", "50–80% subsidy on farm machinery; funds Custom Hiring Centres"),
        ("Agri Infrastructure Fund", "interest-subvented loans for FPOs/SHGs to set up CHCs"),
        ("Namo Drone Didi", "15,000 drones to women SHGs — momentum for ag-robotics"),
        ("Soil Health Card scheme", "national soil-testing push — aligns with our NPK sensing"),
        ("10,000 FPOs + Digital Agri", "aggregation & farm-data tailwinds"),
    ], size=13, mcolor=GREEN)
    side_panel(s, MX + 6.65, CT + 0.05, CW - 6.65, 4.5, [
        ("How we leverage it",
         ["Rover qualifies as agri-machinery → SMAM subsidy cuts farmer cost.",
          "Deploy via FPOs / CHCs financed by AIF → shared-asset model."]),
        ("Revenue streams",
         ["Rover sale + annual service contract.",
          "Rover-as-a-Service: per-acre fee via FPO/CHC.",
          "Subscription: soil & crop analytics + prescription maps."]),
    ], fill="E7F3EC")
    source_note(s, "Scheme details per Govt of India / PIB — verify exact terms before claims.")
    B.fade(s)


def t_built(n):
    s = B.slide(); B.header(s, "What We've Built", "A working POC / MVP", n, PLUM)
    cards = [
        ("Firmware · ESP32/FreeRTOS", GREEN, ["Dual-core drive + RS485 NPK", "Dosing state machine · secure link"]),
        ("AI Vision · Raspberry Pi", SKY, ["YOLOv8n weed/disease/obstacle", "Coral TPU · training notebooks"]),
        ("Navigation & Systems", PLUM, ["EKF + GPS + vision geo-tag", "Snake planner · ROS2 · simulator"]),
        ("Software Tooling", AMBER, ["Live dashboard + Telegram", "CI pipeline · black-box logging"]),
        ("Hardware Design", "EF6C00", ["Verified circuit + wiring v2", "110-part costed BOM (+15 audit)"]),
        ("Field Reliability", "2E7D5B", ["10 real field risks engineered", "Heat, stall, comms-loss, sealing"]),
    ]
    w3 = (CW - 2 * 0.3) / 3; ch = 1.95
    for i, (t, col, lines) in enumerate(cards):
        cx = MX + (i % 3) * (w3 + 0.3)
        cyy = CT + 0.15 + (i // 3) * (ch + 0.3)
        B.card(s, cx, cyy, w3, ch, t, lines, accent=col, icon="●",
               title_size=13, body_size=11.5)
    _, tf = B.textbox(s, MX, CB - 0.32, CW, 0.3)
    B.para(tf, "Built end-to-end by our 4-member team — software validated in "
               "simulation; hardware integration is the next step.", size=11,
           color=LGRAY, italic=True, first=True, space_after=0)
    B.fade(s)


def t_metrics(n):
    s = B.slide(); B.header(s, "PoC Hypothesis & Metrics", "What we will prove", n)
    banner(s, CT + 0.1,
           [("Hypothesis:  ", 14, GREEN, True, HF),
            ("per-plant sensing + targeted dosing cuts fertiliser/chemical use "
             "sharply — with no yield loss.", 13.5, INK, False)],
           fill="E7F3EC", h=0.8)
    B.bullets(s, MX, CT + 1.25, CW, 2.5, [
        ("Input reduction", "fertiliser & pesticide saved vs hand-broadcasting"),
        ("Detection accuracy", "weed/disease precision & recall; dosing placement (cm)"),
        ("Coverage & economics", "area per charge; cost saving per acre per season"),
    ], size=15)
    banner(s, CT + 3.55,
           [("Targets to prove in pilot:  ", 14, AMBER, True, HF),
            ("30–50% input reduction  ·  >85% detection accuracy.", 14, INK, True)],
           fill=PANEL, h=0.75)
    B.fade(s)


def t_impact(n):
    s = B.slide(); B.header(s, "Projected Impact & Scale", "The numbers we're chasing", n, GREEN)
    B.place_image(s, "impact_charts.png", MX, CT + 0.05, CW, 3.0, frame=False, shadow=False)
    my = CT + 3.2; mw = (CW - 2 * 0.3) / 3
    metric_bar(s, MX, my, mw, "AI detection accuracy", "88%", 0.88, SKY)
    metric_bar(s, MX + (mw + 0.3), my, mw, "Fertiliser reduction", "45%", 0.45, GREEN)
    metric_bar(s, MX + 2 * (mw + 0.3), my, mw, "Cost saving / acre", "40%", 0.40, AMBER)
    source_note(s, "Projected from our simulation model & targets — to be validated in field pilots.")
    B.fade(s)


def t_exec(n):
    s = B.slide(); B.header(s, "Execution & Go-to-Market", "From prototype to fleet", n, SKY)
    ms = [
        ("BUILD", GREEN, "Hardware prototype", "BTS7960 drive, NPK probe, pan/tilt sprayer, LiFePO4 — ~3 months"),
        ("TEST", SKY, "Pilot on 5–10 farms", "measure input savings, accuracy & uptime"),
        ("ITERATE", AMBER, "Harden from field data", "tune dosing + AI; seal for dust & 40–45 °C heat"),
        ("SCALE", PLUM, "FPO / CHC roll-out", "SMAM/AIF support → Rover-as-a-Service + data subscription"),
    ]
    y = 3.5
    B.rect(s, MX + 0.2, y, CW - 0.4, 0.03, fill=MIST)
    step = (CW - 0.4) / (len(ms) - 1)
    for i, (tag, col, title, sub) in enumerate(ms):
        cx = MX + 0.2 + i * step
        B.rect(s, cx - 0.17, y - 0.17, 0.36, 0.36, fill=col, shape=MSO_SHAPE.OVAL, shadow=True)
        above = (i % 2 == 0)
        bx = cx - 1.4
        if i == 0:
            bx = cx - 0.2
        if i == len(ms) - 1:
            bx = cx - 2.6
        bx = max(MX, bx)
        by = y - 1.95 if above else y + 0.45
        B.chip(s, bx, by, tag, color=col, w=1.15, size=12)
        _, tf = B.textbox(s, bx, by + 0.46, 3.0, 1.3)
        B.para(tf, title, size=13.5, color=INK, bold=True, font=HF, first=True,
               space_after=2, line_spacing=1.0)
        B.para(tf, sub, size=11, color=GRAY, space_after=0, line_spacing=1.05)
    B.fade(s)


def t_team(n):
    s = B.slide(); B.header(s, "Team of 4 & Impact", "Built end-to-end at IIT Bombay", n)
    members = [
        ("Hitanshu Kapadiya", "Mechanical Engineering", GREEN),
        ("Vivek Gupta", "Mechanical Engineering", SKY),
        ("Shreyash Wagh", "Mechanical Engineering", PLUM),
        ("Pritish Nandy", "Aerospace Engineering", AMBER),
    ]
    w4 = (CW - 3 * 0.34) / 4; y = CT + 0.1; h = 2.3
    for i, (name, dept, col) in enumerate(members):
        x = MX + i * (w4 + 0.34)
        pnl = B.rect(s, x, y, w4, h, fill=WHITE, line=MIST, line_w=1,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        pnl.adjustments[0] = 0.06
        B.rect(s, x, y, w4, 0.6, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        B.rect(s, x, y + 0.32, w4, 0.28, fill=col)
        av = B.rect(s, x + w4 / 2 - 0.5, y + 0.38, 1.0, 1.0, fill=col,
                    shape=MSO_SHAPE.OVAL, shadow=True)
        _, tf = B.textbox(s, x + w4 / 2 - 0.5, y + 0.38, 1.0, 1.0, anchor=MSO_ANCHOR.MIDDLE)
        ini = "".join([q[0] for q in name.split()[:2]]).upper()
        B.para(tf, ini, size=26, color=WHITE, bold=True, font=HF,
               align=PP_ALIGN.CENTER, first=True, space_after=0)
        _, tf = B.textbox(s, x + 0.12, y + 1.5, w4 - 0.24, 0.75, anchor=MSO_ANCHOR.TOP)
        B.para(tf, name, size=13.5, color=INK, bold=True, font=HF,
               align=PP_ALIGN.CENTER, first=True, space_after=3, line_spacing=0.98)
        B.para(tf, dept, size=11, color=col, bold=True, align=PP_ALIGN.CENTER,
               space_after=0, line_spacing=1.0)
    side_panel(s, MX, y + h + 0.28, CW, 1.75, [
        ("Why us",
         ["An IIT Bombay team of 4 that built the full software stack, circuit "
          "design and mechanical fabrication — end to end."]),
        ("Impact",
         ["Lower input cost, less runoff, healthier soil and higher net income — "
          "affordable precision farming for every smallholder."]),
    ])
    B.fade(s)


def t_budget(n):
    s = B.slide(); B.header(s, "Budget & Build Tiers", "Affordable, modular, costed", n, AMBER)
    B.place_image(s, "cost_tiers.png", MX, CT + 0.15, 7.2, 4.0, frame=True)
    side_panel(s, MX + 7.5, CT + 0.1, CW - 7.5, 4.3, [
        ("Build tiers (India prices)",
         ["Core / demo: ₹8,000–14,000",
          "Core + Navigation: ₹12,000–20,000",
          "Full AI build: ₹27,000–48,000"]),
        ("Notes",
         ["Off-the-shelf parts; software is free.",
          "3D-printed parts are PETG, not bought.",
          "Verify prices before ordering."]),
    ])
    B.fade(s)


def t_thanks(n):
    s = B.slide()
    bg = B.rect(s, 0, 0, 13.333, 7.5, fill=FOREST); B.grad(bg, FOREST, FOREST2, 55)
    B.rect(s, 0, 0, 0.3, 7.5, fill=LEAF)
    c1 = B.rect(s, 9.9, 3.0, 5.6, 5.6, fill=GREEN, shape=MSO_SHAPE.OVAL); B._translucent(c1, 18000)
    c2 = B.rect(s, -1.6, -1.6, 3.6, 3.6, fill=LEAF, shape=MSO_SHAPE.OVAL); B._translucent(c2, 12000)
    B.brand_wordmark(s, MX, 0.7, h=0.5, on_dark=True)
    _, tf = B.textbox(s, MX, 2.2, 11.0, 1.2)
    B.para(tf, "Thank you.", size=54, color=WHITE, bold=True, font=HF, first=True, space_after=0)
    _, tf = B.textbox(s, MX, 3.45, 11.2, 0.6)
    B.para(tf, "Precision farming for every small farm.", size=20, color="D9EAD9",
           first=True, space_after=0)
    _, tf = B.textbox(s, MX, 4.2, 11.2, 0.4)
    B.para(tf, "OUR ASK", size=13, color=AMBER, bold=True, font=HF, first=True, space_after=0)
    B.bullets(s, MX, 4.55, 11.2, 1.4, [
        ("Mentorship & a pilot", "field access via an FPO / Custom Hiring Centre"),
        ("Seed support", "to build the hardware prototype and run the first pilots"),
    ], size=14, mcolor=LEAF, tcolor="EAF3EA")
    B.rect(s, MX, 5.95, 11.5, 0.02, fill="27503C")
    _, tf = B.textbox(s, MX, 6.08, 11.6, 0.9)
    B.para(tf, "Hitanshu Kapadiya · Vivek Gupta · Shreyash Wagh · Pritish Nandy",
           size=13, color=WHITE, bold=True, first=True, space_after=4)
    B.para(tf, "IIT Bombay  ·  " + EVENT, size=12, color="9FC2AC", space_after=0)
    B.fade(s)


def app_divider():
    s = B.slide()
    bg = B.rect(s, 0, 0, 13.333, 7.5, fill=FOREST); B.grad(bg, FOREST, FOREST2, 60)
    c1 = B.rect(s, 10.4, -1.6, 4.4, 4.4, fill=GREEN, shape=MSO_SHAPE.OVAL); B._translucent(c1, 20000)
    B.rect(s, 0, 0, 0.28, 7.5, fill=LEAF)
    _, tf = B.textbox(s, MX, 3.0, 10.5, 0.4)
    B.para(tf, "APPENDIX", size=15, color=AMBER, bold=True, font=HF, first=True, space_after=0)
    _, tf = B.textbox(s, MX, 3.35, 11.2, 1.2)
    B.para(tf, "Engineering Depth", size=42, color=WHITE, bold=True, font=HF, first=True, space_after=0)
    _, tf = B.textbox(s, MX, 4.5, 10.8, 0.6)
    B.para(tf, "Verified circuit, drive, mechanical layout & field-risk engineering.",
           size=16, color="CFE6D8", first=True, space_after=0)
    B.fade(s)


def t_app_wiring(n):
    s = B.slide(); B.header(s, "Appendix · Electronics", "Every wire accounted for", n, SKY)
    _, tf = B.textbox(s, MX, CT - 0.02, CW, 0.45)
    B.para(tf, "110 BOM components + 15 gap-audit items — every bus, divider and "
               "safety part placed.", size=13, color=GRAY, first=True, line_spacing=1.05)
    B.place_image(s, "wiring.png", MX, 2.4, CW, 4.45,
                  caption="Consolidated v2 wiring schematic — red = safety-critical paths.")
    B.fade(s)


def t_app_drive(n):
    s = B.slide(); B.header(s, "Appendix · Drive System", "MOSFET power for soil loads", n, SKY)
    B.place_image(s, "drive_schematic.png", MX, CT, 8.5, 4.85,
                  caption="Pin-level 2× BTS7960 (IBT-2) drive schematic.")
    B.bullets(s, MX + 8.8, CT + 0.25, CW - 8.8, 4.3, [
        ("2× BTS7960", "replaced the L298N (FC-10)"),
        ("~43 A peak/side", "handles soil-loaded motors"),
        ("Dual-PWM", "forward / reverse / coast"),
        ("1N5819 flyback", "clamps back-EMF"),
        ("Current-sense", "IS → stall detection"),
    ], size=13, mcolor=SKY)
    B.fade(s)


def t_app_chassis(n):
    s = B.slide(); B.header(s, "Appendix · Mechanical", "Stable, service-friendly", n)
    B.place_image(s, "chassis_layout.png", MX, CT, 8.4, 4.8,
                  caption="To-scale layout from verified component coordinates.")
    B.bullets(s, MX + 8.7, CT + 0.25, CW - 8.7, 4.3, [
        ("Double-decker", "320 × 450 mm, two layers"),
        ("Low CoG", "battery + tank low, centred"),
        ("Fluids below", "leaks drain from the Pi"),
        ("Heat apart", "drivers at the rear + fan"),
        ("Field-hardened", "grommets, glands, coating"),
    ], size=13, mcolor=GREEN)
    B.fade(s)


def t_app_fc(n):
    s = B.slide(); B.header(s, "Appendix · Field-Ready", "Ten field risks, engineered", n, AMBER)
    rows = [
        ("ID", "Field challenge", "Engineered solution", "Status"),
        ("FC-01", "Weeds above the pathway", "Vision-guided pan/tilt aimed spray", "Code-ready"),
        ("FC-02", "Field heat → battery fire", "Thermal guardian + NTC + LiFePO4", "Code-ready"),
        ("FC-03", "Actuator retract type", "Dual compile-flag branches", "Code-ready"),
        ("FC-04", "Per-plant GPS precision", "SBAS + averaging + EKF + vision", "Implemented"),
        ("FC-05", "Sensors over-volt the ADC", "Powered from the 3.3 V rail", "Implemented"),
        ("FC-06", "Relays twitch ON at boot", "10 kΩ pull-ups + fail-safe OFF", "Implemented"),
        ("FC-07", "Comms link drops", "HMAC heartbeat dead-man halt", "Implemented"),
        ("FC-08", "Dust, dew & spray", "Conformal coat + glands + gaskets", "Implemented"),
        ("FC-09", "Wheel / cutter jam", "ADS1115 current-sense stall stop", "Implemented"),
        ("FC-10", "Driver under-spec (L298N)", "2× BTS7960 MOSFET drivers", "Code-ready"),
    ]
    x, y, w, h = MX, CT + 0.02, CW, 4.85
    ts = s.shapes.add_table(len(rows), 4, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = ts.table; tbl.first_row = False; tbl.horz_banding = False
    for j, wd in enumerate([1.1, 4.2, 4.65, 1.5]):
        tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(h / len(rows))
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(FOREST if i == 0 else (WHITE if i % 2 else PANEL))
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in (0, 3) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = HF if (i == 0 or j == 0) else BF
            r.font.size = Pt(10.5 if i == 0 else 10)
            r.font.bold = (i == 0 or j == 0 or j == 3)
            if i == 0:
                r.font.color.rgb = C(WHITE)
            elif j == 0:
                r.font.color.rgb = C(GREEN)
            elif j == 3:
                r.font.color.rgb = C(GREEN if val == "Implemented" else AMBER)
            else:
                r.font.color.rgb = C(INK)
    B.fade(s)


def main():
    t_title()
    t_problem(2)
    t_solution(3)
    t_why(4)
    t_arch(5)
    t_sense_see(6)
    t_act(7)
    t_market(8)
    t_govt(9)
    t_built(10)
    t_metrics(11)
    t_impact(12)
    t_exec(13)
    t_team(14)
    t_budget(15)
    t_thanks(16)
    app_divider()
    t_app_wiring(18)
    t_app_drive(19)
    t_app_chassis(20)
    t_app_fc(21)
    B.prs.save(OUT)
    print("Saved:", OUT)
    print("Slides:", len(B.prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
