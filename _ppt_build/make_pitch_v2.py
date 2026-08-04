# Build AgriRover_GrowwxIITB_TrackA_v2.pptx from the Gamma-designed reference deck.
# Interviewer feedback applied:
#   - remove ALL grant/subsidy mentions (SMAM, AIF, CHC, Namo Drone Didi, Soil Health Card)
#   - problems & solution in simple language
#   - add a real-farm Examples slide (2 examples)
#   - mechanism in simple terms
#   - uses slide
#   - differentiation vs tractors & drones
import sys, copy
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

SRC = "Precision-farming-for-every-small-farm (3).pptx"
OUT = "AgriRover_GrowwxIITB_TrackA_v2.pptx"
prs = Presentation(SRC)


def shape_by_id(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(f"shape id {sid} not on slide")


def set_text(slide, sid, paras):
    """paras: list of lists-of-run-texts, one inner list per paragraph."""
    sh = shape_by_id(slide, sid)
    tf = sh.text_frame
    assert len(tf.paragraphs) == len(paras), (
        f"para count mismatch on shape {sid}: {len(tf.paragraphs)} vs {len(paras)}"
    )
    for p, texts in zip(tf.paragraphs, paras):
        runs = p.runs
        assert len(runs) >= len(texts), f"run count issue on shape {sid}"
        for i, t in enumerate(texts):
            runs[i].text = t
        for r in runs[len(texts):]:
            r._r.getparent().remove(r._r)


def duplicate_slide(prs, index):
    """Copy a slide (shapes, background, image rels) so the design stays consistent."""
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    for sh in list(dest.shapes):  # drop auto-added placeholders
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dest.shapes._spTree.append(copy.deepcopy(sh._element))
    src_cSld = src._element.find(qn("p:cSld"))
    bg = src_cSld.find(qn("p:bg")) if src_cSld is not None else None
    if bg is not None:
        dest._element.find(qn("p:cSld")).insert(0, copy.deepcopy(bg))
    rid_map = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith("/notesSlide"):
            continue
        if rel.is_external:
            rid_map[rId] = dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            rid_map[rId] = dest.part.relate_to(rel.target_part, rel.reltype)
    for el in dest._element.iter():
        for attr in list(el.attrib):
            if attr.endswith("}embed") or attr.endswith("}link") or attr.endswith("}id"):
                v = el.attrib[attr]
                if v in rid_map:
                    el.attrib[attr] = rid_map[v]
    return dest


# ---------- 0. duplicate slide 8 (idx 7) -> Examples slide (appended at end) ----------
ex = duplicate_slide(prs, 7)

# ---------- 1. SLIDE 2 : PROBLEM, simple ----------
s = prs.slides[1]
set_text(s, 4, [["THE PROBLEM TODAY"]])
shape_by_id(s, 3).width = Inches(1.45)  # tag pill
set_text(s, 6, [["Fertilizer and water go everywhere because farmers can't see what each plant actually needs — so much of it is simply wasted."]])
set_text(s, 9, [["86% of India's farms are under 2 hectares — big machines don't fit and cost too much."]])
set_text(s, 11, [["Farming Blind"]])
set_text(s, 12, [["No easy way to know what the soil lacks — so farmers guess, and over- or under-feed."]])
set_text(s, 14, [["Rising Costs"]])
set_text(s, 15, [["Fertilizer, chemicals and labor cost more every year — and much of the spend is wasted."]])
set_text(s, 18, [["A lab soil test takes 2–3 weeks — the season moves on, so most farmers simply skip it."]])
set_text(s, 19, [["Our field visits confirm it: fertilizer is still spread by hand — evenly, everywhere, every season."]])

# ---------- 2. SLIDE 3 : SOLUTION, simple ----------
s = prs.slides[2]
shape_by_id(s, 4).width = Inches(11.5)
set_text(s, 4, [["A Small Rover That Farms Plant by Plant"]])
set_text(s, 5, [["AgriRover drives the field on its own, ", "checks the soil and every plant", ", and gives each one only what it needs."]])
set_text(s, 7, [["1. CHECK"]])
set_text(s, 8, [["A probe reads soil nutrients & moisture on the spot."]])
set_text(s, 10, [["2. SEE"]])
set_text(s, 11, [["An AI camera spots weeds, diseases and obstacles."]])
set_text(s, 13, [["3. TREAT"]])
set_text(s, 14, [["Doses fertilizer & sprays only where it's needed."]])
set_text(s, 16, [["4. REPORT"]])
set_text(s, 17, [["Sends a simple soil & crop map to the farmer's phone."]])
set_text(s, 20, [["Why it's different:", " it works plant by plant, tests the soil itself, and costs a fraction of big machinery."]])

# ---------- 3. SLIDE 4 : MECHANISM, simple ----------
s = prs.slides[3]
set_text(s, 3, [["MECHANISM"]])
shape_by_id(s, 2).width = Inches(1.1)
set_text(s, 5, [["The rover has two brains — one that acts, one that thinks — so it always stays in control."]])
set_text(s, 8, [["Acting brain: drives, doses, stops"]])
set_text(s, 9, [["Secure Link"]])
set_text(s, 10, [["A tamper-proof wire between them"]])
set_text(s, 12, [["Thinking brain: sees & decides"]])
set_text(s, 13, [["The acting brain never waits for the thinking brain. If the AI hangs or the camera fails, the rover simply slows down and stops — it can never run away, crash, or overdose a plant."]])

# ---------- 4. SLIDE 5 : SENSE & SEE, simple ----------
s = prs.slides[4]
set_text(s, 7, [["One dip reads 7 things — N, P, K, pH, salts, moisture & temperature — in seconds, right in the field. No lab, no waiting."]])
set_text(s, 8, [["The AI Camera"]])
set_text(s, 9, [["The camera recognises weeds, 38 crop diseases and obstacles in real time, on the rover itself — like a tireless agronomist walking every row."]])
set_text(s, 12, [["One dip, full soil report"]])
set_text(s, 14, [["Diseases Spotted"]])
set_text(s, 15, [["Learned from real crop photos"]])
set_text(s, 17, [["Photos a Second"]])
set_text(s, 18, [["Sees in real time, on the rover"]])
set_text(s, 20, [["Plant Memory"]])
set_text(s, 21, [["Remembers each plant's exact spot"]])

# ---------- 5. SLIDE 6 : USES ----------
s = prs.slides[5]
set_text(s, 4, [["WHAT IT CAN DO"]])
set_text(s, 7, [["Feeds Each Plant"]])
set_text(s, 8, [["Probes the soil, then drops the exact dose of fertilizer at that plant — no more, no less."]])
set_text(s, 10, [["Sprays Only the Weeds"]])
set_text(s, 11, [["The camera finds each weed and aims the nozzle at just that spot — the crop stays clean."]])
set_text(s, 13, [["Drives Itself, Row by Row"]])
set_text(s, 14, [["Covers the field on GPS — the farmer just places it at the field edge."]])
set_text(s, 16, [["One Rover, Many Jobs"]])
set_text(s, 17, [["Swap tools — seeder, weeder, cutter, sprayer. And if anything goes wrong, it simply stops — safely."]])

# ---------- 6. SLIDE 7 : MARKET + BUSINESS MODEL (no subsidies) ----------
s = prs.slides[6]
set_text(s, 3, [["MARKET & BUSINESS MODEL"]])
shape_by_id(s, 2).width = Inches(2.6)
set_text(s, 16, [["How We Earn"]])
set_text(s, 17, [
    ["Sale:", " rover + annual service contract"],
    ["RaaS:", " per-acre fee via FPOs & farmer co-ops"],
    ["Data:", " soil & crop analytics subscription"],
    ["Growth:", " vegetables → cotton → orchards → polyhouses"],
])

# ---------- 7. SLIDE 8 : DIFFERENTIATION (replaces tailwinds/subsidies) ----------
s = prs.slides[7]
set_text(s, 3, [["HOW WE'RE DIFFERENT"]])
shape_by_id(s, 2).width = Inches(2.1)
shape_by_id(s, 4).width = Inches(9.0)
set_text(s, 4, [["Not a Tractor. Not a Drone."]])
set_text(s, 5, [["What Others Do"]])
set_text(s, 6, [
    ["Tractor:", " ploughs and pulls, but can't see plants — too big & costly for small plots."],
    ["Drone:", " a quick look and spray from the sky — it can't touch, test or treat the soil."],
    ["Lab tests:", " one soil report in 2–3 weeks — no detail on which patch needs what."],
    ["Hired labor:", " scarce and pricier every year — hand-spreading is even, not precise."],
])
set_text(s, 7, [["What AgriRover Does"]])
set_text(s, 8, [["Drives on the ground, tests the soil it touches, sees every plant, and treats each one — no pilot, no lab, no guesswork."]])
set_text(s, 9, [["In One Line"]])
set_text(s, 10, [
    ["Tractor = muscle. Drone = a look from the sky."],
    ["AgriRover = eyes + hands at every plant."],
    ["It replaces guesswork — not the tractor."],
])

# ---------- 8. EXAMPLES slide (the duplicate, currently last) ----------
s = ex
set_text(s, 3, [["REAL-FARM EXAMPLES"]])
shape_by_id(s, 2).width = Inches(2.0)
set_text(s, 4, [["Two Everyday Examples"]])
set_text(s, 5, [["Example 1 — Fertilizer"]])
set_text(s, 6, [
    ["Today:", " a full bag of urea, spread evenly by hand — much of it wasted."],
    ["Rover:", " probes the field and finds only one patch is nitrogen-poor."],
    ["Action:", " doses fertilizer only there, plant by plant."],
    ["Result:", " 30–50% less fertilizer — the target we'll prove in pilots."],
])
set_text(s, 7, [["Example 2 — Weeds"]])
set_text(s, 8, [["Instead of spraying the whole field, the camera finds each weed and sprays only those spots — far less chemical, none on the crop."]])
set_text(s, 9, [["What It Means"]])
set_text(s, 10, [
    ["A smaller input bill every season."],
    ["Healthier soil, year after year."],
    ["Less dependence on scarce farm labor."],
])

# ---------- 9. SLIDE 11 : timeline, drop SMAM/AIF/CHC ----------
s = prs.slides[10]
set_text(s, 24, [["Roll out via FPOs & farmer co-ops → Rover-as-a-Service + data."]])

# ---------- 10. SLIDE 13 : closing ask, drop Custom Hiring Centre ----------
s = prs.slides[12]
sh = shape_by_id(s, 5)
sh.text_frame.paragraphs[0].runs[1].text = (
    " mentorship + a field pilot with a farmer group (FPO), and seed support "
    "to build the prototype and run the first pilots."
)

# ---------- 11. SLIDE 12 : cost notes tidy ----------
s = prs.slides[11]
for sh in s.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith("Notes:"):
        p = sh.text_frame.paragraphs[0]
        p.runs[0].text = "Notes: standard off-the-shelf parts; our software adds no licence cost; body is 3D-printed (PETG)."
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)

# ---------- 12. reorder: Examples after Solution; Differentiation after Uses ----------
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
ex_id = ids[15]   # duplicated Examples slide (appended last)
dif_id = ids[7]   # slide 8 = differentiation
sldIdLst.remove(ex_id)
sldIdLst.insert(3, ex_id)      # after Solution
sldIdLst.remove(dif_id)
sldIdLst.insert(7, dif_id)     # after Uses, before Market

prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides)} slides")
