# Build AgriRover_Investor_Full.pptx — 10-slide deck for IITB-Groww INV.ENT.
# Changes vs v3 (per Vivek):
#   * slides 1-3 unchanged (title / problem / solution)
#   * slide 4  -> competition vs XMachines (robots) & Marut Drones (AG365)
#   * slide 5  -> personal story (Vivek, farmer family, Jaunpur UP)
#   * slide 6  -> old slide 5 (market, model & ask)
#   * slides 7-10 -> everything else the official template demands:
#       PoC Hypothesis & Metrics · Execution & Business Plan ·
#       Team, Budget & Impact · Budget Sheet
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = "_ppt_build/assets/v3"
OUT = "AgriRover_Investor_Full.pptx"
N_SLIDES = 10

# ---- design system (identical to v3) --------------------------------------
INK    = RGBColor(0x1F, 0x2A, 0x22)
MUT    = RGBColor(0x5C, 0x6B, 0x5E)
ACC    = RGBColor(0x2E, 0x7D, 0x32)
ACC_D  = RGBColor(0x1E, 0x4D, 0x36)
CARD   = RGBColor(0xF2, 0xF7, 0xF2)
CARD2  = RGBColor(0xE7, 0xF1, 0xE7)
LINE   = RGBColor(0xDD, 0xE8, 0xDD)
BG     = RGBColor(0xFC, 0xFB, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAYX  = RGBColor(0xA8, 0xB5, 0xA8)
FONT   = "Calibri"

prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333in x 7.5in (16:9)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(slide, x, y, w, h, fill=None, line=None, round_=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        shp.adjustments[0] = round_
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if "sb" in p: para.space_before = Pt(p["sb"])
        if "sa" in p: para.space_after = Pt(p["sa"])
        if "ls" in p: para.line_spacing = p["ls"]
        for r in p["runs"]:
            text, size, bold, color = r[0], r[1], r[2], r[3]
            run = para.add_run()
            run.text = text
            f = run.font
            f.name, f.size, f.bold = FONT, Pt(size), bold
            f.color.rgb = color
            if len(r) > 4 and r[4]:
                f.italic = True
    return tb


def kicker(slide, x, y, label, w):
    pill = box(slide, x, y, w, 0.34, fill=CARD2, round_=0.5)
    tf = pill.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    f = run.font
    f.name, f.size, f.bold, f.color.rgb = FONT, Pt(10.5), True, ACC_D


def footer(slide, n):
    txt(slide, 0.7, 7.08, 6.0, 0.3,
        [{"runs": [("AgriRover · IITB–Groww INV.ENT · Track A", 8.5, False, GRAYX)]}])
    txt(slide, 11.75, 7.08, 0.9, 0.3,
        [{"runs": [(f"{n} / {N_SLIDES}", 8.5, False, GRAYX)], "align": PP_ALIGN.RIGHT}])


def side_img(slide, fname):
    slide.shapes.add_picture(f"{ASSETS}/{fname}", Inches(8.08), 0, height=Inches(7.5))


def wordmark(slide):
    box(slide, 0.7, 0.58, 0.26, 0.26, fill=ACC, round_=0.35)
    txt(slide, 1.06, 0.56, 3.0, 0.32,
        [{"runs": [("AgriRover", 14, True, INK)]}])


def card_head(slide, x, y, w, label):
    txt(slide, x + 0.25, y + 0.13, w - 0.5, 0.28,
        [{"runs": [(label, 10.5, True, ACC_D)]}])


# ===========================================================================
# SLIDE 1 — TITLE  (unchanged)
# ===========================================================================
s = new_slide()
side_img(s, "s1_Image_0_sm.png")
wordmark(s)
txt(s, 0.7, 1.95, 7.0, 1.6, [
    {"runs": [("Precision farming", 40, True, INK)], "ls": 1.02},
    {"runs": [("for every small farm.", 40, True, ACC_D)], "ls": 1.02},
])
txt(s, 0.7, 3.5, 6.9, 1.1, [
    {"runs": [("A small self-driving rover that checks the soil, looks at every plant, "
               "and gives each one exactly what it needs — nothing more, nothing wasted.",
               13.5, False, MUT)], "ls": 1.15},
])
pill = box(s, 0.7, 4.85, 2.95, 0.44, fill=ACC_D, round_=0.5)
tf = pill.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = tf.margin_right = Inches(0.1)
tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "IITB–Groww INV.ENT · Track A"
r.font.name, r.font.size, r.font.bold, r.font.color.rgb = FONT, Pt(11), True, WHITE
txt(s, 0.7, 6.35, 7.0, 0.7, [
    {"runs": [("Hitanshu Kapadiya · Vivek Gupta · Shreyash Wagh · Pritish Nandy", 11.5, True, INK)], "sa": 2},
    {"runs": [("Team of 4 · IIT Bombay", 10.5, False, MUT)]},
])

# ===========================================================================
# SLIDE 2 — PROBLEM & VALIDATION  (unchanged)
# ===========================================================================
s = new_slide()
side_img(s, "s2_Image_0_sm.png")
kicker(s, 0.7, 0.6, "THE PROBLEM", 1.55)
txt(s, 0.7, 1.08, 7.1, 1.0, [
    {"runs": [("Farmers spend more every year —", 27, True, INK)], "ls": 1.02},
    {"runs": [("and waste much of it.", 27, True, ACC_D)], "ls": 1.02},
])
problems = [
    ("Farming is done blind.",
     "Fertilizer and water are spread evenly by hand, because there is no way to know "
     "what each plant needs. A big share is simply wasted."),
    ("Waste costs real money.",
     "Fertilizer, chemicals and labor get costlier every season — and over-use quietly "
     "damages the soil for years to come."),
    ("Existing machines don't fit.",
     "86% of India's farms are under 2 hectares — too small for tractors and big precision "
     "equipment. A lab soil test takes 2–3 weeks, so almost nobody does one."),
]
y = 2.28
for title, body in problems:
    box(s, 0.7, y, 6.95, 1.28, fill=CARD, line=LINE, round_=0.08)
    txt(s, 0.95, y + 0.13, 6.45, 0.3, [{"runs": [(title, 13.5, True, INK)]}])
    txt(s, 0.95, y + 0.46, 6.45, 0.72, [{"runs": [(body, 11, False, MUT)], "ls": 1.1}])
    y += 1.44
txt(s, 0.7, y + 0.04, 6.95, 0.35, [
    {"runs": [("From our field visits: fertilizer is still thrown by hand — evenly, "
               "everywhere, every season.", 10.5, False, ACC_D, True)]},
])
footer(s, 2)

# ===========================================================================
# SLIDE 3 — SOLUTION & INNOVATION  (unchanged)
# ===========================================================================
s = new_slide()
side_img(s, "s6_Image_0_sm.png")
kicker(s, 0.7, 0.6, "OUR SOLUTION", 1.65)
txt(s, 0.7, 1.08, 7.1, 0.55, [
    {"runs": [("A small rover that farms ", 27, True, INK),
              ("plant by plant.", 27, True, ACC_D)]},
])
txt(s, 0.7, 1.68, 7.0, 0.45, [
    {"runs": [("It drives the crop rows on its own and treats every plant individually — "
               "like a careful farmhand that never gets tired.", 12, False, MUT)], "ls": 1.1},
])
steps = [
    ("1", "CHECK", "A probe dips into the soil and reads its health in seconds — no lab, no waiting."),
    ("2", "SEE", "An AI camera spots weeds and sick plants as the rover drives the rows."),
    ("3", "TREAT", "Drops fertilizer or a tiny spray only where it is needed, plant by plant."),
    ("4", "REPORT", "Sends a simple map to the farmer's phone: what it found, what it did."),
]
positions = [(0.7, 2.35), (4.32, 2.35), (0.7, 3.88), (4.32, 3.88)]
for (num, title, body), (x, y) in zip(steps, positions):
    box(s, x, y, 3.42, 1.38, fill=CARD, line=LINE, round_=0.08)
    chip = box(s, x + 0.18, y + 0.16, 0.34, 0.34, fill=ACC, round_=0.5)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num
    cr.font.name, cr.font.size, cr.font.bold, cr.font.color.rgb = FONT, Pt(12), True, WHITE
    txt(s, x + 0.64, y + 0.2, 2.6, 0.3, [{"runs": [(title, 12.5, True, INK)]}])
    txt(s, x + 0.18, y + 0.62, 3.06, 0.7, [{"runs": [(body, 10.5, False, MUT)], "ls": 1.08}])
box(s, 0.7, 5.48, 7.04, 1.3, fill=CARD2, line=None, round_=0.08)
txt(s, 0.95, 5.62, 6.55, 0.3, [{"runs": [("HOW IT WORKS INSIDE — TWO BRAINS, ALWAYS SAFE", 10.5, True, ACC_D)]}])
txt(s, 0.95, 5.94, 6.55, 0.75, [
    {"runs": [("A thinking brain", 11, True, INK),
              (" sees and decides; an ", 11, False, MUT),
              ("acting brain", 11, True, INK),
              (" drives and doses. If anything fails — camera, AI or signal — the rover "
               "simply stops. It cannot run away or overdose a plant.", 11, False, MUT)], "ls": 1.12},
])
footer(s, 3)

# ===========================================================================
# SLIDE 4 — COMPETITION: XMACHINES & MARUT DRONES
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "COMPETITION", 1.55)
txt(s, 0.7, 1.0, 11.9, 0.55, [
    {"runs": [("India already builds farm robots and drones. ", 25, True, INK),
              ("None reach the small farm.", 25, True, ACC_D)]},
])
# ---- left: the two named competitors ----
comp = [
    ("XMACHINES · Hyderabad", "Autonomous field robots (X100, Neo)",
     "Serious machines for weeding, spraying and mowing — RTK-GPS, AI vision, "
     "300 kg payload. But the Neo is priced at $19,995 (≈ ₹17 lakh) and aimed at "
     "large farms and export markets. No in-soil nutrient testing."),
    ("MARUT DRONES · Hyderabad", "AG365 certified spray drone",
     "Sprays 6 acres an hour from the sky; service model at ~₹400/acre. But the "
     "drone costs ₹7–10 lakh, needs a licensed pilot, and it is blanket spraying — "
     "it cannot test soil or treat one plant differently from the next."),
]
y = 1.78
for name, tag, body in comp:
    box(s, 0.7, y, 5.75, 1.62, fill=CARD, line=LINE, round_=0.07)
    txt(s, 0.95, y + 0.13, 5.25, 0.28, [{"runs": [(name, 11, True, ACC_D)]}])
    txt(s, 0.95, y + 0.42, 5.25, 0.28, [{"runs": [(tag, 11.5, True, INK)]}])
    txt(s, 0.95, y + 0.74, 5.25, 0.8, [{"runs": [(body, 10, False, MUT)], "ls": 1.1}])
    y += 1.78
box(s, 0.7, y, 5.75, 1.28, fill=CARD2, round_=0.07)
txt(s, 0.95, y + 0.13, 5.25, 0.28, [{"runs": [("WHAT THEY PROVE — AND MISS", 10.5, True, ACC_D)]}])
txt(s, 0.95, y + 0.44, 5.25, 0.75, [
    {"runs": [("Both prove Indian farms will pay for automation. Both are priced for the "
               "14% of farms above 2 hectares. The other 86% is open — that's ours.",
               10.5, False, INK)], "ls": 1.12},
])
# ---- right: head-to-head table ----
txt(s, 6.75, 1.78, 5.85, 0.35, [{"runs": [("Head to head", 16, True, INK)]}])
tx, ty = 6.75, 2.28
w_label, w_col, row_h, head_h = 2.55, 1.1, 0.5, 0.38
cols = ["XMACHINES", "MARUT", "AGRIROVER"]
rows = [
    ("Machine price",             ["~₹17 L", "₹7–10 L", "<₹0.5 L"]),
    ("On-the-spot soil test",     ["–", "–", "✓"]),
    ("Per-plant fertilizer dose", ["–", "–", "✓"]),
    ("Spot-sprays single weeds",  ["✓", "–", "✓"]),
    ("Fits a 1-acre veg. plot",   ["–", "–", "✓"]),
    ("No licensed pilot needed",  ["✓", "–", "✓"]),
]
table_h = head_h + row_h * len(rows)
box(s, tx + w_label + 2 * w_col, ty, w_col, table_h, fill=CARD2, round_=None)
for i in range(len(rows)):
    if i % 2 == 0:
        box(s, tx, ty + head_h + i * row_h, w_label + 2 * w_col, row_h, fill=CARD)
for j, c in enumerate(cols):
    color = ACC_D if j == 2 else MUT
    txt(s, tx + w_label + j * w_col, ty + 0.06, w_col, 0.3,
        [{"runs": [(c, 9, True, color)], "align": PP_ALIGN.CENTER}])
for i, (label, marks) in enumerate(rows):
    ry = ty + head_h + i * row_h
    txt(s, tx + 0.12, ry, w_label - 0.12, row_h,
        [{"runs": [(label, 10.5, False, INK)]}], anchor=MSO_ANCHOR.MIDDLE)
    for j, m in enumerate(marks):
        if m == "✓":
            runs = [(m, 13, True, ACC)]
        elif m == "–":
            runs = [(m, 13, True, GRAYX)]
        else:
            runs = [(m, 9.5, True, ACC_D if j == 2 else INK)]
        txt(s, tx + w_label + j * w_col, ry, w_col, row_h,
            [{"runs": runs, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 6.75, ty + table_h + 0.15, 5.85, 0.6, [
    {"runs": [("A drone looks from the sky. A big robot works the big farm. "
               "AgriRover is the first one cheap enough — and small enough — for the "
               "farmer with one acre of vegetables.", 10.5, False, MUT, True)], "ls": 1.1},
])
footer(s, 4)

# ===========================================================================
# SLIDE 5 — PERSONAL STORY (Vivek · Jaunpur, UP)
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "WHY I'M BUILDING THIS", 2.35)
txt(s, 0.7, 1.02, 11.5, 0.55, [
    {"runs": [("I grew up in the fields ", 27, True, INK),
              ("this is built for.", 27, True, ACC_D)]},
])
story = [
    ("My name is Vivek Kumar Gupta. ",
     "I study Mechanical Engineering at IIT Bombay, but I come from a farmer family "
     "in Jaunpur, Uttar Pradesh. My family grows what most of eastern UP grows: "
     "rice in kharif, wheat in rabi, season after season."),
    ("Every season, the same ritual. ",
     "Urea thrown by hand — evenly, everywhere — because no one knows what the soil "
     "actually needs. The nearest soil lab is a bus ride away and the report takes "
     "2–3 weeks. In ten years, I have never seen a neighbour test twice."),
    ("Our wheat and rice can't use this rover — and that's the point. ",
     "Broadcast-sown grain leaves no path for a machine to walk. But all around us, "
     "the money crops — tomato, chilli, cauliflower, potato, brinjal — are planted "
     "in neat rows. Rows a small rover can drive. Those plots take the most "
     "fertilizer and spray per acre, and lose the most when it's wasted."),
]
y = 1.85
for lead, body in story:
    box(s, 0.7, y, 6.9, 1.34, fill=CARD, line=LINE, round_=0.07)
    txt(s, 0.95, y + 0.15, 6.4, 1.05,
        [{"runs": [(lead, 11, True, INK), (body, 11, False, MUT)], "ls": 1.13}])
    y += 1.5
s.shapes.add_picture(f"{ASSETS}/s9_Image_0.png", Inches(7.85), Inches(1.85),
                     width=Inches(4.75))
qy = 4.65
qbox = box(s, 7.85, qy, 4.75, 1.72, fill=ACC_D, round_=0.09)
qtf = qbox.text_frame
qtf.vertical_anchor = MSO_ANCHOR.MIDDLE
qtf.margin_left = qtf.margin_right = Inches(0.22)
qtf.margin_top = qtf.margin_bottom = Inches(0.12)
qp = qtf.paragraphs[0]
qr = qp.add_run()
qr.text = ("“I couldn't build a robot for my father's rice. So I built one for "
           "every crop planted in rows — starting with the vegetable plots "
           "around my own village.”")
qr.font.name, qr.font.size, qr.font.color.rgb = FONT, Pt(12), WHITE
qr.font.italic = True
qp.line_spacing = 1.15
qp2 = qtf.add_paragraph()
qp2.space_before = Pt(6)
qr2 = qp2.add_run(); qr2.text = "— Vivek, AgriRover team"
qr2.font.name, qr2.font.size, qr2.font.bold, qr2.font.color.rgb = FONT, Pt(10.5), True, WHITE
txt(s, 0.7, y + 0.02, 6.9, 0.35, [
    {"runs": [("First pilot fields: the row-crop vegetable plots around Jaunpur — "
               "farmers we already know by name.", 10.5, False, ACC_D, True)]},
])
footer(s, 5)

# ===========================================================================
# SLIDE 6 — MARKET, MODEL & THE ASK  (old slide 5, unchanged content)
# ===========================================================================
s = new_slide()
side_img(s, "s13_Image_0_sm.png")
kicker(s, 0.7, 0.55, "MARKET, MODEL & THE ASK", 2.7)
txt(s, 0.7, 1.0, 7.1, 0.55, [
    {"runs": [("A big market. ", 27, True, INK), ("A simple business.", 27, True, ACC_D)]},
])
stats = [
    ("146M", "farm holdings in India"),
    ("86%", "of farms under 2 hectares — our users"),
    ("$2.5Bn", "India agritech market by 2034"),
]
for i, (num, label) in enumerate(stats):
    x = 0.7 + i * 2.42
    txt(s, x, 1.8, 2.3, 0.45, [{"runs": [(num, 24, True, ACC_D)]}])
    txt(s, x, 2.28, 2.25, 0.55, [{"runs": [(label, 10, False, MUT)], "ls": 1.05}])
box(s, 0.7, 3.05, 7.1, 1.42, fill=CARD, line=LINE, round_=0.07)
txt(s, 0.95, 3.18, 6.6, 0.28, [{"runs": [("HOW WE EARN", 10.5, True, ACC_D)]}])
earn = [
    "1.  Sell the rover with a yearly service plan.",
    "2.  Rent it per acre through farmer groups (FPOs) — no upfront cost for the farmer.",
    "3.  Monthly subscription for soil-health and crop reports.",
]
txt(s, 0.95, 3.5, 6.6, 0.95,
    [{"runs": [(e, 11, False, INK)], "sa": 3} for e in earn])
box(s, 0.7, 4.62, 7.1, 1.05, fill=CARD, line=LINE, round_=0.07)
txt(s, 0.95, 4.74, 6.6, 0.28, [{"runs": [("WHERE WE ARE TODAY", 10.5, True, ACC_D)]}])
txt(s, 0.95, 5.05, 6.6, 0.55, [
    {"runs": [("Full software stack — driving, AI vision, navigation, farmer dashboard — built and "
               "tested in simulation. Hardware fully designed and costed: ₹27,000–48,000 per rover.",
               10.5, False, MUT)], "ls": 1.1},
])
ask = box(s, 0.7, 5.85, 7.1, 0.95, fill=ACC_D, round_=0.09)
atf = ask.text_frame
atf.vertical_anchor = MSO_ANCHOR.MIDDLE
atf.margin_left = atf.margin_right = Inches(0.25)
atf.margin_top = atf.margin_bottom = Inches(0.08)
ap = atf.paragraphs[0]
ar = ap.add_run(); ar.text = "The ask:  "
ar.font.name, ar.font.size, ar.font.bold, ar.font.color.rgb = FONT, Pt(12.5), True, WHITE
ar2 = ap.add_run()
ar2.text = ("seed support to build the first prototype, plus a pilot with one farmer "
            "group to prove 30–50% input savings on real farms.")
ar2.font.name, ar2.font.size, ar2.font.bold, ar2.font.color.rgb = FONT, Pt(12.5), False, WHITE
ap.line_spacing = 1.12
footer(s, 6)

# ===========================================================================
# SLIDE 7 — POC HYPOTHESIS & METRICS  (template requirement)
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "POC HYPOTHESIS & METRICS", 2.85)
txt(s, 0.7, 1.0, 11.9, 0.55, [
    {"runs": [("One season. Three plots. ", 27, True, INK),
              ("Numbers, not promises.", 27, True, ACC_D)]},
])
hyp = box(s, 0.7, 1.85, 11.9, 0.98, fill=ACC_D, round_=0.08)
htf = hyp.text_frame
htf.vertical_anchor = MSO_ANCHOR.MIDDLE
htf.margin_left = htf.margin_right = Inches(0.25)
htf.margin_top = htf.margin_bottom = Inches(0.06)
hp = htf.paragraphs[0]
hr = hp.add_run(); hr.text = "HYPOTHESIS:  "
hr.font.name, hr.font.size, hr.font.bold, hr.font.color.rgb = FONT, Pt(12.5), True, WHITE
hr2 = hp.add_run()
hr2.text = ("if soil is sensed and fertilizer dosed plant-by-plant on row vegetables, "
            "the farmer saves 30–50% of input cost with no loss of yield.")
hr2.font.name, hr2.font.size, hr2.font.color.rgb = FONT, Pt(12.5), WHITE
hp.line_spacing = 1.1
box(s, 0.7, 3.0, 11.9, 0.86, fill=CARD2, round_=0.08)
txt(s, 0.95, 3.12, 11.4, 0.28, [{"runs": [("HOW WE TEST IT", 10.5, True, ACC_D)]}])
txt(s, 0.95, 3.4, 11.4, 0.42, [
    {"runs": [("A/B on 3 vegetable plots (Jaunpur network), one growing season: half of each "
               "plot farmed as usual, half managed by AgriRover. Same seed, same water, same farmer.",
               10.5, False, INK)], "ls": 1.1},
])
metrics = [
    ("−30–50%", "fertilizer used per acre", "measured in kg, A vs B"),
    ("−60–80%", "weedkiller sprayed", "spot-spray vs blanket spray"),
    ("±0", "yield per acre", "must not drop — else we fail"),
    ("30 sec", "soil test, in the field", "vs 2–3 weeks at a lab"),
    ("1 acre/day", "covered per rover", "coverage & uptime logged"),
    ("<2 seasons", "payback for the farmer", "input savings vs rover cost"),
]
positions = [(0.7, 4.1), (4.68, 4.1), (8.66, 4.1), (0.7, 5.45), (4.68, 5.45), (8.66, 5.45)]
for (num, label, note), (x, y) in zip(metrics, positions):
    box(s, x, y, 3.72, 1.2, fill=CARD, line=LINE, round_=0.08)
    txt(s, x + 0.2, y + 0.12, 3.3, 0.4, [{"runs": [(num, 17, True, ACC_D)]}])
    txt(s, x + 0.2, y + 0.55, 3.3, 0.3, [{"runs": [(label, 11, True, INK)]}])
    txt(s, x + 0.2, y + 0.84, 3.3, 0.3, [{"runs": [(note, 9.5, False, MUT)]}])
footer(s, 7)

# ===========================================================================
# SLIDE 8 — EXECUTION & BUSINESS PLAN  (template requirement)
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "EXECUTION & BUSINESS PLAN", 2.9)
txt(s, 0.7, 1.0, 11.9, 0.55, [
    {"runs": [("Build. Test. Iterate. ", 27, True, INK),
              ("Then rent it out.", 27, True, ACC_D)]},
])
miles = [
    ("M1 · BUILD", "Weeks 1–8",
     "Assemble the first rover from our fully-costed BOM (₹27–48k). The design, "
     "wiring and firmware are already done — this is assembly and bench testing, "
     "then autonomous runs on campus."),
    ("M2 · TEST", "Weeks 9–20",
     "One growing season with 3–5 farmers from our own networks (Jaunpur UP, "
     "Nashik MH). Run the A/B protocol from slide 7 and collect the savings "
     "and yield data."),
    ("M3 · ITERATE", "Weeks 21–28",
     "Fix what the field breaks. Add the attachment farmers ask for most. "
     "Start the first paid rental pilot with one FPO."),
]
for i, (name, when, body) in enumerate(miles):
    x = 0.7 + i * 4.0
    box(s, x, 1.85, 3.78, 2.2, fill=CARD, line=LINE, round_=0.07)
    txt(s, x + 0.22, 2.0, 2.2, 0.3, [{"runs": [(name, 12.5, True, ACC_D)]}])
    wpill = box(s, x + 2.42, 2.0, 1.15, 0.3, fill=CARD2, round_=0.5)
    wtf = wpill.text_frame
    wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    wtf.margin_left = wtf.margin_right = Inches(0.05)
    wtf.margin_top = wtf.margin_bottom = 0
    wp = wtf.paragraphs[0]; wp.alignment = PP_ALIGN.CENTER
    wr = wp.add_run(); wr.text = when
    wr.font.name, wr.font.size, wr.font.bold, wr.font.color.rgb = FONT, Pt(9), True, ACC_D
    txt(s, x + 0.22, 2.42, 3.35, 1.55, [{"runs": [(body, 10.5, False, MUT)], "ls": 1.13}])
box(s, 0.7, 4.35, 7.35, 1.55, fill=CARD, line=LINE, round_=0.07)
card_head(s, 0.7, 4.35, 7.35, "PRICING PLAN")
plans = [
    ("Own it:  ", "₹75,000 per rover + ₹6,000/yr service & software plan."),
    ("Rent it:  ", "₹300 per acre-pass through FPOs — cheaper than the drone's ₹400, and it "
                   "tests soil while it drives."),
    ("Know it:  ", "₹199/month soil-health & crop report subscription."),
]
txt(s, 0.95, 4.72, 6.85, 1.1,
    [{"runs": [(lead, 10.5, True, INK), (rest, 10.5, False, MUT)], "sa": 4, "ls": 1.08}
     for lead, rest in plans])
box(s, 8.25, 4.35, 4.35, 1.55, fill=CARD2, round_=0.07)
card_head(s, 8.25, 4.35, 4.35, "COMPETITORS")
txt(s, 8.5, 4.72, 3.85, 1.1, [
    {"runs": [("XMachines (₹17 L robots) and Marut (₹7–10 L drones) — detailed on "
               "slide 4. Nobody serves the under-2-hectare farm at under ₹1 lakh. "
               "We do.", 10.5, False, INK)], "ls": 1.12},
])
txt(s, 0.7, 6.15, 11.9, 0.5, [
    {"runs": [("Every milestone ends in something a farmer can see: a rover that drives (M1), "
               "a savings number (M2), a paying pilot (M3).", 11, False, ACC_D, True)]},
])
footer(s, 8)

# ===========================================================================
# SLIDE 9 — TEAM, BUDGET & IMPACT  (template requirement)
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "TEAM, BUDGET & IMPACT", 2.5)
txt(s, 0.7, 1.0, 11.9, 0.55, [
    {"runs": [("Four builders. ", 27, True, INK),
              ("One of us grew up on the farm.", 27, True, ACC_D)]},
])
team = [
    ("Hitanshu Kapadiya", "Mech. Engg, IITB", "Chassis, CAD & drivetrain"),
    ("Vivek Kumar Gupta", "Mech. Engg, IITB", "Electronics & firmware · farmer-family domain (Jaunpur)"),
    ("Shreyash Wagh", "Mech. Engg, IITB", "AI vision & model training (Nashik veg. belt)"),
    ("Pritish Nandy", "IITB", "Software, dashboard & field data"),
]
for i, (name, dept, role) in enumerate(team):
    x = 0.7 + i * 3.02
    box(s, x, 1.82, 2.86, 1.28, fill=CARD, line=LINE, round_=0.08)
    txt(s, x + 0.18, 1.95, 2.5, 0.3, [{"runs": [(name, 11, True, INK)]}])
    txt(s, x + 0.18, 2.22, 2.5, 0.25, [{"runs": [(dept, 9, False, GRAYX)]}])
    txt(s, x + 0.18, 2.48, 2.5, 0.55, [{"runs": [(role, 9.5, False, MUT)], "ls": 1.05}])
box(s, 0.7, 3.3, 5.9, 1.0, fill=CARD2, round_=0.07)
card_head(s, 0.7, 3.3, 5.9, "WHY US")
txt(s, 0.95, 3.65, 5.4, 0.6, [
    {"runs": [("We already built the full software stack — and we know the user by name. "
               "Our pilot farmers are our own families' neighbours.", 10.5, False, INK)], "ls": 1.1},
])
box(s, 0.7, 4.45, 5.9, 2.15, fill=CARD, line=LINE, round_=0.07)
card_head(s, 0.7, 4.45, 5.9, "USE OF FUNDS — ₹1.0 LAKH, 28 WEEKS")
funds = [
    ("₹49k", "Prototype parts (full BOM, next slide)", "M1 · weeks 1–8"),
    ("₹30k", "Field pilot: travel, spares, consumables", "M2 · weeks 9–20"),
    ("₹21k", "Iteration, attachments, FPO pilot", "M3 · weeks 21–28"),
]
fy = 4.82
for amt, what, when in funds:
    txt(s, 0.95, fy, 0.85, 0.3, [{"runs": [(amt, 11.5, True, ACC_D)]}])
    txt(s, 1.85, fy, 3.55, 0.3, [{"runs": [(what, 10, False, INK)]}])
    txt(s, 5.02, fy, 1.45, 0.3, [{"runs": [(when, 8.5, False, GRAYX)]}])
    fy += 0.36
txt(s, 0.95, fy + 0.05, 5.4, 0.5, [
    {"runs": [("Proposed output: ", 10, True, INK),
              ("one field-proven rover + a one-season savings dataset + a signed FPO "
               "rental pilot.", 10, False, MUT)], "ls": 1.08},
])
box(s, 6.85, 3.3, 5.75, 1.5, fill=CARD, line=LINE, round_=0.07)
card_head(s, 6.85, 3.3, 5.75, "IMPACT — WHAT IT'S WORTH")
imp = [
    ("For one farmer: ", "₹8–15k saved per season on one acre of vegetables — often the "
                         "family's entire margin."),
    ("For the soil & the plate: ", "less chemical in the ground and none of it on the food."),
]
txt(s, 7.1, 3.65, 5.25, 1.05,
    [{"runs": [(lead, 10, True, INK), (rest, 10, False, MUT)], "sa": 4, "ls": 1.08}
     for lead, rest in imp])
box(s, 6.85, 4.95, 5.75, 1.65, fill=ACC_D, round_=0.08)
vtf_box = txt(s, 7.1, 5.12, 5.25, 0.28, [{"runs": [("LONG-TERM VISION", 10.5, True, WHITE)]}])
txt(s, 7.1, 5.42, 5.25, 1.05, [
    {"runs": [("A rover in every village, not every farm: FPO-owned fleets with swappable "
               "attachments (weeder, sower, sprayer) working plots on demand — precision "
               "farming as a ₹300 service, not a ₹17 lakh machine.", 10.5, False, WHITE)], "ls": 1.13},
])
footer(s, 9)

# ===========================================================================
# SLIDE 10 — BUDGET SHEET  (template requirement)
# ===========================================================================
s = new_slide()
kicker(s, 0.7, 0.55, "BUDGET SHEET", 1.6)
txt(s, 0.7, 1.0, 11.9, 0.55, [
    {"runs": [("Every rupee is already itemised — ", 27, True, INK),
              ("110 parts, priced.", 27, True, ACC_D)]},
])
rows = [
    ("Brain & compute (ESP32, Pi 4, cams)", "5,880 – 8,150"),
    ("Power (LiPo, buck, protection, solar)", "2,350 – 4,040"),
    ("Motors & drive (4 gear motors, drivers)", "1,840 – 3,600"),
    ("Soil & field sensors (NPK probe, GPS…)", "3,740 – 8,080"),
    ("AI camera kit (Pi Cam, Coral TPU)", "4,680 – 7,100"),
    ("Dosing & actuation (pump, actuator…)", "1,760 – 3,330"),
    ("Chassis, wiring & passives", "2,210 – 4,130"),
    ("Safety, thermal & tools", "3,630 – 7,040"),
    ("Comms, interface & upgrade parts", "1,740 – 3,440"),
]
tx, ty = 0.7, 1.85
w_item, w_amt, row_h = 4.9, 1.85, 0.44
txt(s, tx + 0.12, ty + 0.02, w_item, 0.3, [{"runs": [("SUBSYSTEM", 9.5, True, MUT)]}])
txt(s, tx + w_item, ty + 0.02, w_amt, 0.3,
    [{"runs": [("₹ (range)", 9.5, True, MUT)], "align": PP_ALIGN.RIGHT}])
for i, (item, amt) in enumerate(rows):
    ry = ty + 0.34 + i * row_h
    if i % 2 == 0:
        box(s, tx, ry, w_item + w_amt, row_h, fill=CARD)
    txt(s, tx + 0.12, ry, w_item - 0.12, row_h,
        [{"runs": [(item, 10, False, INK)]}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, tx + w_item, ry, w_amt - 0.12, row_h,
        [{"runs": [(amt, 10, False, INK)], "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
ty_total = ty + 0.34 + len(rows) * row_h
box(s, tx, ty_total, w_item + w_amt, 0.52, fill=ACC_D, round_=0.12)
txt(s, tx + 0.15, ty_total, w_item - 0.15, 0.52,
    [{"runs": [("FULL AI ROVER — TOTAL", 11, True, WHITE)]}], anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx + w_item, ty_total, w_amt - 0.15, 0.52,
    [{"runs": [("₹27,830 – 48,910", 11, True, WHITE)], "align": PP_ALIGN.RIGHT}],
    anchor=MSO_ANCHOR.MIDDLE)
# right column: tiers + pilot budget + note
box(s, 7.9, 1.85, 4.7, 1.7, fill=CARD, line=LINE, round_=0.07)
card_head(s, 7.9, 1.85, 4.7, "BUILD TIERS")
tiers = [
    ("Core (drives + soil test + dosing)", "₹8–14k"),
    ("+ Navigation (GPS, encoders, safety)", "₹12–20k"),
    ("Full AI (camera + weed/disease AI)", "₹28–49k"),
]
tyy = 2.22
for name, amt in tiers:
    txt(s, 8.15, tyy, 3.3, 0.3, [{"runs": [(name, 9.5, False, INK)]}])
    txt(s, 11.35, tyy, 1.0, 0.3, [{"runs": [(amt, 9.5, True, ACC_D)], "align": PP_ALIGN.RIGHT}])
    tyy += 0.34
txt(s, 8.15, tyy + 0.02, 4.2, 0.35,
    [{"runs": [("We build Full AI for the PoC — it's what we must prove.", 9, False, MUT, True)]}])
box(s, 7.9, 3.75, 4.7, 1.35, fill=CARD, line=LINE, round_=0.07)
card_head(s, 7.9, 3.75, 4.7, "PILOT & CONTINGENCY")
pil = [
    ("Field pilot (travel, spares, consumables)", "₹30k"),
    ("Iteration & attachment parts", "₹21k"),
]
pyy = 4.12
for name, amt in pil:
    txt(s, 8.15, pyy, 3.55, 0.3, [{"runs": [(name, 9.5, False, INK)]}])
    txt(s, 11.6, pyy, 0.75, 0.3, [{"runs": [(amt, 9.5, True, ACC_D)], "align": PP_ALIGN.RIGHT}])
    pyy += 0.34
txt(s, 8.15, pyy + 0.02, 4.2, 0.3,
    [{"runs": [("Biggest swing items: NPK probe, Pi 4, Coral TPU.", 9, False, MUT, True)]}])
askb = box(s, 7.9, 5.3, 4.7, 1.05, fill=ACC_D, round_=0.09)
atf = askb.text_frame
atf.vertical_anchor = MSO_ANCHOR.MIDDLE
atf.margin_left = atf.margin_right = Inches(0.2)
atf.margin_top = atf.margin_bottom = Inches(0.06)
ap = atf.paragraphs[0]
ar = ap.add_run(); ar.text = "Total ask: ₹1.0 lakh "
ar.font.name, ar.font.size, ar.font.bold, ar.font.color.rgb = FONT, Pt(13), True, WHITE
ar2 = ap.add_run()
ar2.text = "— one rover, one season of proof, one paying pilot."
ar2.font.name, ar2.font.size, ar2.font.color.rgb = FONT, Pt(11), WHITE
ap.line_spacing = 1.12
txt(s, 0.7, 6.78, 11.9, 0.3, [
    {"runs": [("Full 110-part priced list with vendors lives in our repo (docs/shopping-list.md) — "
               "available on request.", 9.5, False, GRAYX)]},
])
footer(s, 10)

prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides)} slides")
