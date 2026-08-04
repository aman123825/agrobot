"""Structural verification of the generated deck (no rendering needed)."""
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

path = "AgriRover_GrowwxIITB.pptx"
p = Presentation(path)
n = len(p.slides._sldIdLst)
print("slides:", n, flush=True)
print("size:", round(Emu(p.slide_width).inches, 3), "x",
      round(Emu(p.slide_height).inches, 3), flush=True)
print("-" * 78, flush=True)

total_pics = total_tables = 0
issues = []
for i, s in enumerate(p.slides, 1):
    npic = ntbl = nshape = 0
    title = ""
    for sh in s.shapes:
        nshape += 1
        if sh.shape_type == 13:
            npic += 1
        if sh.has_table:
            ntbl += 1
        if sh.has_text_frame and sh.text_frame.text.strip() and not title:
            title = sh.text_frame.text.strip().split("\n")[0][:44]
    total_pics += npic
    total_tables += ntbl
    has_fade = s._element.find(qn("p:transition")) is not None
    if nshape == 0:
        issues.append(f"slide {i}: no shapes")
    if not has_fade:
        issues.append(f"slide {i}: no transition")
    print(f"{i:2d} | shapes {nshape:2d} | pics {npic} | tbl {ntbl} | "
          f"fade {int(has_fade)} | {title}", flush=True)

print("-" * 78, flush=True)
print("total pictures:", total_pics, "| total tables:", total_tables, flush=True)
for i, s in enumerate(p.slides, 1):
    for sh in s.shapes:
        if sh.has_table:
            t = sh.table
            print(f"table on slide {i}: {len(t.rows)}x{len(t.columns)} "
                  f"header={[t.cell(0,j).text for j in range(len(t.columns))]}",
                  flush=True)
print("RESULT:", "PASS - no structural issues" if not issues else issues, flush=True)
