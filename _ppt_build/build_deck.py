"""Build AgriRover_GrowwxIITB.pptx — polished pitch deck for the Groww x IIT Bombay
final evaluation. python-pptx only. Fade transitions + soft shadows via XML.

Run:  python _ppt_build/build_deck.py
"""
import os
from PIL import Image, ImageChops

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
ASSETS = os.path.join(HERE, "assets")
CROP = os.path.join(HERE, "assets", "cropped")
os.makedirs(CROP, exist_ok=True)
OUT = os.path.join(REPO, "AgriRover_GrowwxIITB.pptx")

# ---------------- Theme ----------------
INK    = "12281C"
FOREST = "0E3B2A"
FOREST2 = "07271B"
GREEN  = "1F9254"
LEAF   = "6FBF44"
AMBER  = "F2A03D"
SKY    = "2B9BD6"
PLUM   = "6A4C93"
CLOUD  = "F5F8F3"
MIST   = "E7EFE7"
PANEL  = "EEF4EC"
GRAY   = "5B6B62"
LGRAY  = "9AA8A0"
WHITE  = "FFFFFF"
RED    = "E4572E"

HEAD_FONT = "Segoe UI Semibold"
HEAD_FONT2 = "Segoe UI"
BODY_FONT = "Calibri"

SW = Inches(13.333)
SH = Inches(7.5)


def C(h):
    return RGBColor.from_string(h)


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------------- low-level helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line)
        sp.line.width = Pt(line_w)
    if not shadow:
        _no_shadow(sp)
    else:
        _soft_shadow(sp)
    return sp


def _no_shadow(sp):
    spPr = sp._element.spPr
    if spPr.find(qn("a:effectLst")) is None:
        spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def _soft_shadow(sp, blur=90000, dist=38000, alpha=62000, direction=5400000):
    spPr = sp._element.spPr
    old = spPr.find(qn("a:effectLst"))
    if old is not None:
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(blur), "dist": str(dist), "dir": str(direction),
        "rotWithShape": "0"})
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "0B241A"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(alpha)})
    clr.append(a)
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def grad(sp, c1, c2, angle=90):
    """Two-stop linear gradient fill on an existing shape."""
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = C(c1)
    stops[1].position = 1.0
    stops[1].color.rgb = C(c2)
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size=16, color=INK, bold=False, font=BODY_FONT,
         align=PP_ALIGN.LEFT, first=False, space_after=6, space_before=0,
         line_spacing=1.08, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = C(color)
    return p, run


def runs(p, parts):
    """parts = list of (text, size, color, bold[, font])."""
    for t in parts:
        r = p.add_run()
        r.text = t[0]
        r.font.size = Pt(t[1])
        r.font.color.rgb = C(t[2])
        r.font.bold = t[3]
        r.font.name = t[4] if len(t) > 4 else BODY_FONT


def fade(s, spd="med"):
    sld = s._element
    for t in sld.findall(qn("p:transition")):
        sld.remove(t)
    tr = sld.makeelement(qn("p:transition"), {"spd": spd})
    tr.append(tr.makeelement(qn("p:fade"), {}))
    cm = sld.find(qn("p:clrMapOvr"))
    if cm is not None:
        cm.addnext(tr)
    else:
        sld.append(tr)


# ---------------- image handling ----------------
_IMG = {}


def _prep(name):
    if name in _IMG:
        return _IMG[name]
    src = os.path.join(ASSETS, name)
    dst = os.path.join(CROP, name)
    im = Image.open(src).convert("RGB")
    diff = ImageChops.difference(im, Image.new("RGB", im.size, (255, 255, 255)))
    bb = diff.getbbox()
    if bb:
        pad = 10
        l, t, r, b = bb
        im = im.crop((max(0, l - pad), max(0, t - pad),
                      min(im.width, r + pad), min(im.height, b + pad)))
    im.save(dst)
    _IMG[name] = (dst, im.width, im.height)
    return _IMG[name]


def place_image(s, name, bx, by, bw, bh, caption=None, shadow=True, frame=True,
                cap_color=GRAY):
    path, iw, ih = _prep(name)
    cap_h = 0.34 if caption else 0.0
    avail_h = bh - cap_h
    scale = min(bw / iw, avail_h / ih)
    w = iw * scale
    h = ih * scale
    x = bx + (bw - w) / 2
    y = by + (avail_h - h) / 2
    if frame:
        pad = 0.06
        bg = rect(s, x - pad, y - pad, w + 2 * pad, h + 2 * pad, fill=WHITE,
                  line=MIST, line_w=1.0, shadow=shadow)
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        _, tf = textbox(s, bx, by + avail_h + 0.02, bw, cap_h,
                        anchor=MSO_ANCHOR.TOP)
        para(tf, caption, size=10.5, color=cap_color, italic=True,
             align=PP_ALIGN.CENTER, first=True, space_after=0)
    return pic


# ---------------- brand + layout ----------------
MX = 0.92
CT = 1.82          # content top
CB = 6.98          # content bottom
CW = 13.333 - 2 * MX


def brand_badge(s, x, y, h=0.5, on_dark=False):
    """A small leaf-tile logo mark (rounded tile + leaf/letter)."""
    tile = rect(s, x, y, h, h, fill=(LEAF if on_dark else GREEN),
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tile.adjustments[0] = 0.28
    made_leaf = False
    try:
        leaf = rect(s, x + h * 0.24, y + h * 0.14, h * 0.52, h * 0.62,
                    fill=WHITE, shape=MSO_SHAPE.TEAR)
        leaf.rotation = 135
        made_leaf = True
    except Exception:
        made_leaf = False
    if not made_leaf:
        _, tf = textbox(s, x, y, h, h, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "A", size=h * 34, color=WHITE, bold=True, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, first=True, space_after=0)
    return x + h + 0.12


def brand_wordmark(s, x, y, h=0.5, on_dark=False):
    nx = brand_badge(s, x, y, h, on_dark)
    _, tf = textbox(s, nx, y - 0.04, 3.0, h + 0.08, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    runs(p, [("Agri", 17, WHITE if on_dark else INK, True, HEAD_FONT),
             ("Rover", 17, LEAF if on_dark else GREEN, True, HEAD_FONT)])


def header(s, kicker, title, n=None, accent=GREEN):
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)               # base bg
    rect(s, 0, 0, 0.22, 7.5, fill=accent)                # left accent spine
    _, tf = textbox(s, MX, 0.46, CW - 2.5, 0.3)
    para(tf, kicker.upper(), size=12, color=accent, bold=True,
         font=HEAD_FONT, first=True, space_after=0)
    _, tf = textbox(s, MX, 0.72, CW, 0.75)
    para(tf, title, size=29, color=INK, bold=True, font=HEAD_FONT,
         first=True, space_after=0)
    rect(s, MX, 1.5, 1.5, 0.055, fill=accent)            # underline accent
    # top-right brand
    _, tf = textbox(s, 13.333 - 3.3, 0.5, 2.55, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    runs(p, [("Agri", 12.5, INK, True, HEAD_FONT),
             ("Rover", 12.5, GREEN, True, HEAD_FONT)])
    dot = rect(s, 13.333 - 0.62, 0.58, 0.2, 0.2, fill=accent,
               shape=MSO_SHAPE.OVAL)
    footer(s, n)


def footer(s, n=None):
    rect(s, MX, 7.06, CW, 0.014, fill=MIST)
    _, tf = textbox(s, MX, 7.12, 8.0, 0.3)
    para(tf, "AgriRover  ·  IITB–Groww INV.ENT  ·  Track A  ·  Team of 4 (IIT Bombay)",
         size=9, color=LGRAY, first=True, space_after=0)
    if n is not None:
        _, tf = textbox(s, 13.333 - MX - 1.2, 7.12, 1.2, 0.3)
        para(tf, str(n), size=9, color=LGRAY, align=PP_ALIGN.RIGHT, first=True,
             space_after=0)


def section_divider(n, kicker, title, subtitle=""):
    s = slide()
    bg = rect(s, 0, 0, 13.333, 7.5, fill=FOREST)
    grad(bg, FOREST, FOREST2, angle=60)
    # decorative accent circles
    c1 = rect(s, 10.4, -1.6, 4.4, 4.4, fill=GREEN, shape=MSO_SHAPE.OVAL)
    c1.fill.fore_color.rgb = C(GREEN)
    _translucent(c1, 22000)
    c2 = rect(s, 11.7, 4.3, 3.4, 3.4, fill=LEAF, shape=MSO_SHAPE.OVAL)
    _translucent(c2, 16000)
    rect(s, 0, 0, 0.28, 7.5, fill=LEAF)
    # big number
    _, tf = textbox(s, MX - 0.06, 1.5, 4.0, 2.2)
    para(tf, f"{n:02d}", size=120, color=LEAF, bold=True, font=HEAD_FONT,
         first=True, space_after=0)
    _, tf = textbox(s, MX, 3.75, 10.5, 0.4)
    para(tf, kicker.upper(), size=14, color=AMBER, bold=True, font=HEAD_FONT,
         first=True, space_after=0)
    _, tf = textbox(s, MX, 4.05, 11.2, 1.4)
    para(tf, title, size=40, color=WHITE, bold=True, font=HEAD_FONT,
         first=True, space_after=0)
    if subtitle:
        _, tf = textbox(s, MX, 5.35, 10.8, 1.0)
        para(tf, subtitle, size=16, color="CFE6D8", first=True, space_after=0,
             line_spacing=1.15)
    fade(s)
    return s


def _translucent(sp, alpha):
    # replace solid fill with alpha
    fe = sp.fill.fore_color._xFill
    srgb = fe.find(qn("a:srgbClr"))
    if srgb is not None:
        a = srgb.makeelement(qn("a:alpha"), {"val": str(alpha)})
        srgb.append(a)


def card(s, x, y, w, h, title, body_lines, accent=GREEN, icon="", num=None,
         title_size=15, body_size=12.5):
    panel = rect(s, x, y, w, h, fill=WHITE, line=MIST, line_w=1.0,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    panel.adjustments[0] = 0.06
    rect(s, x, y, 0.12, h, fill=accent)   # left color rail (clipped visually ok)
    ic = rect(s, x + 0.28, y + 0.26, 0.62, 0.62, fill=accent,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ic.adjustments[0] = 0.3
    _, tf = textbox(s, x + 0.28, y + 0.26, 0.62, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, (num if num else icon), size=(19 if num else 17), color=WHITE,
         bold=True, font=HEAD_FONT, align=PP_ALIGN.CENTER, first=True,
         space_after=0)
    _, tf = textbox(s, x + 1.06, y + 0.22, w - 1.3, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, size=title_size, color=INK, bold=True, font=HEAD_FONT,
         first=True, space_after=0, line_spacing=1.0)
    _, tf = textbox(s, x + 0.3, y + 1.02, w - 0.55, h - 1.15)
    for i, ln in enumerate(body_lines):
        para(tf, ln, size=body_size, color=GRAY, first=(i == 0), space_after=4,
             line_spacing=1.06)
    return panel


def kpi(s, x, y, w, h, big, label, accent=GREEN):
    panel = rect(s, x, y, w, h, fill=WHITE, line=MIST, line_w=1.0,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    panel.adjustments[0] = 0.08
    rect(s, x, y, w, 0.12, fill=accent)
    _, tf = textbox(s, x + 0.15, y + 0.28, w - 0.3, h * 0.5,
                    anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, size=30, color=accent, bold=True, font=HEAD_FONT,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    _, tf = textbox(s, x + 0.15, y + h * 0.58, w - 0.3, h * 0.4,
                    anchor=MSO_ANCHOR.TOP)
    para(tf, label, size=11.5, color=GRAY, align=PP_ALIGN.CENTER, first=True,
         space_after=0, line_spacing=1.02)
    return panel


def chip(s, x, y, text, color=GREEN, tcolor=WHITE, w=None, size=11.5):
    w = w if w else (0.16 * len(text) + 0.3)
    c = rect(s, x, y, w, 0.36, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.5
    _, tf = textbox(s, x, y, w, 0.36, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, size=size, color=tcolor, bold=True, align=PP_ALIGN.CENTER,
         font=HEAD_FONT2, first=True, space_after=0)
    return x + w + 0.14


def bullets(s, x, y, w, h, items, size=15, marker="▸", mcolor=GREEN,
            gap=9, tcolor=INK):
    _, tf = textbox(s, x, y, w, h)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            head, sub = it
        else:
            head, sub = it, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.06
        r = p.add_run(); r.text = marker + "  "
        r.font.size = Pt(size); r.font.bold = True; r.font.name = HEAD_FONT
        r.font.color.rgb = C(mcolor)
        r = p.add_run(); r.text = head
        r.font.size = Pt(size); r.font.name = BODY_FONT; r.font.color.rgb = C(tcolor)
        r.font.bold = True
        if sub:
            r = p.add_run(); r.text = "  —  " + sub
            r.font.size = Pt(size - 1.5); r.font.name = BODY_FONT
            r.font.color.rgb = C(GRAY); r.font.bold = False
    return tf




# =====================================================================
# SLIDES
# =====================================================================
def s_title():
    s = slide()
    bg = rect(s, 0, 0, 13.333, 7.5, fill=FOREST)
    grad(bg, FOREST, FOREST2, angle=55)
    rect(s, 0, 0, 0.3, 7.5, fill=LEAF)
    c1 = rect(s, 9.7, -2.2, 5.6, 5.6, fill=GREEN, shape=MSO_SHAPE.OVAL)
    _translucent(c1, 20000)
    c2 = rect(s, 11.6, 3.6, 4.2, 4.2, fill=LEAF, shape=MSO_SHAPE.OVAL)
    _translucent(c2, 13000)
    c3 = rect(s, 8.7, 4.9, 2.2, 2.2, fill=AMBER, shape=MSO_SHAPE.OVAL)
    _translucent(c3, 12000)

    brand_wordmark(s, MX, 0.62, h=0.52, on_dark=True)

    _, tf = textbox(s, MX, 2.05, 9.6, 0.4)
    para(tf, "AUTONOMOUS  PRECISION-AGRICULTURE  ROVER", size=15, color=AMBER,
         bold=True, font=HEAD_FONT, first=True, space_after=0)
    _, tf = textbox(s, MX - 0.03, 2.42, 10.5, 1.5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    runs(p, [("Agri", 68, WHITE, True, HEAD_FONT), ("Rover", 68, LEAF, True, HEAD_FONT)])
    _, tf = textbox(s, MX, 3.86, 10.8, 0.7)
    para(tf, "Sense the soil.  See the crop.  Target the action.", size=24,
         color="D9EAD9", bold=True, font=HEAD_FONT2, first=True, space_after=0)
    _, tf = textbox(s, MX, 4.62, 10.3, 0.9)
    para(tf, "A dual-controller ESP32 + Raspberry Pi rover that pairs 7-in-1 soil "
             "sensing with on-board AI vision to micro-dose fertiliser and spot-spray "
             "weeds — cutting input waste on every pass.",
         size=14.5, color="BFD8C8", first=True, space_after=0, line_spacing=1.2)

    chip(s, MX, 5.72, "Groww × IIT Bombay  ·  Final Evaluation", color=AMBER,
         tcolor=INK, w=4.6, size=12.5)

    rect(s, MX, 6.42, 11.5, 0.02, fill="27503C")
    _, tf = textbox(s, MX, 6.56, 11.6, 0.5)
    p = tf.paragraphs[0]
    runs(p, [("Kapadiya Hitanshu Mukeshbhai", 12.5, WHITE, True),
             ("  25B2228      ", 11.5, LEAF, False),
             ("Vivek Kumar Gupta", 12.5, WHITE, True),
             ("  25B2269      ", 11.5, LEAF, False),
             ("Shreyash Wagh", 12.5, WHITE, True),
             ("  25B2227", 11.5, LEAF, False)])
    _, tf = textbox(s, MX, 6.92, 11.6, 0.4)
    para(tf, "B.Tech Mechanical Engineering  ·  Indian Institute of Technology Bombay",
         size=11.5, color="9FC2AC", first=True, space_after=0)
    fade(s)


def s_problem(n):
    s = slide()
    header(s, "The Problem", "Farming still runs on guesswork", n, accent=RED)
    _, tf = textbox(s, MX, CT, CW, 0.7)
    para(tf, "Agriculture supports nearly half of India's workforce — yet most "
             "input decisions are still made without field data, wasting money and "
             "degrading soil.", size=15.5, color=GRAY, first=True, line_spacing=1.15)
    cy = 2.78
    cw2 = (CW - 0.4) / 2
    ch = 1.85
    card(s, MX, cy, cw2, ch, "Blanket, blind inputs",
         ["Fertiliser & pesticide sprayed uniformly with no",
          "soil test — over-applied here, starved there."],
         accent=RED, icon="✕")
    card(s, MX + cw2 + 0.4, cy, cw2, ch, "Rising cost & labour",
         ["Farm-labour shortages and climbing input prices",
          "squeeze already-thin margins."],
         accent=AMBER, icon="₹")
    card(s, MX, cy + ch + 0.32, cw2, ch, "Soil & yield damage",
         ["NPK imbalance and chemical overuse erode long-",
          "term fertility and pollute groundwater."],
         accent="B4632B", icon="!")
    card(s, MX + cw2 + 0.4, cy + ch + 0.32, cw2, ch, "Late, un-targeted detection",
         ["Weeds, pests and disease are spotted late and",
          "treated field-wide instead of plant-by-plant."],
         accent=PLUM, icon="?")
    fade(s)


def s_solution(n):
    s = slide()
    header(s, "Our Solution", "One rover: sense → decide → act", n)
    _, tf = textbox(s, MX, CT, CW, 0.7)
    para(tf, "AgriRover turns every pass across the field into data-driven, "
             "input-saving action — an affordable robot built from off-the-shelf parts.",
         size=15.5, color=GRAY, first=True, line_spacing=1.15)
    cy = 2.72
    w4 = (CW - 3 * 0.32) / 4
    ch = 3.6
    data = [
        ("SENSE", GREEN, "1",
         ["7-in-1 NPK probe", "soil moisture · TDS · pH", "GPS + IMU + encoders",
          "ultrasonic + ToF"]),
        ("DECIDE", SKY, "2",
         ["On-board AI vision", "obstacle detection", "weed detection",
          "crop-disease ID"]),
        ("ACT", AMBER, "3",
         ["Targeted micro-dosing", "vision-aimed spot-spray", "grass-cutter · seeder",
          "weeder · sprayer"]),
        ("LEARN", PLUM, "4",
         ["Geo-tagged field logs", "prescription maps", "live dashboard + alerts",
          "per-plant history"]),
    ]
    for i, (t, col, num, lines) in enumerate(data):
        card(s, MX + i * (w4 + 0.32), cy, w4, ch, t, lines, accent=col, num=num,
             title_size=16, body_size=12.5)
    fade(s)


def s_architecture(n):
    s = slide()
    header(s, "System Architecture", "Two brains, one mission", n)
    _, tf = textbox(s, MX, CT - 0.02, CW, 0.62)
    para(tf, "A real-time ESP32 (FreeRTOS) handles drive, sensing and dosing; a "
             "Raspberry Pi 4 runs AI vision and analytics — linked by an "
             "authenticated command channel.", size=14, color=GRAY, first=True,
         line_spacing=1.12)
    place_image(s, "architecture.png", MX, 2.5, CW, 4.35,
                caption="Dual-controller architecture — real-time control on the ESP32, "
                        "AI + analytics on the Pi, joined by an HMAC-signed link.")
    fade(s)


def s_loop(n):
    s = slide()
    header(s, "How It Works", "A closed sense–decide–act loop", n)
    steps = [
        ("SENSE", GREEN, ["Soil NPK, moisture,", "GPS, camera frame"]),
        ("DECIDE", SKY, ["AI + rules pick the", "right action & spot"]),
        ("ACT", AMBER, ["Micro-dose, spot-spray,", "steer or stop"]),
        ("LEARN", PLUM, ["Geo-tag, log, alert,", "update the map"]),
    ]
    n_steps = len(steps)
    bw = 2.35
    gap = (CW - n_steps * bw) / (n_steps - 1)
    y = 3.0
    h = 2.0
    for i, (t, col, lines) in enumerate(steps):
        x = MX + i * (bw + gap)
        panel = rect(s, x, y, bw, h, fill=WHITE, line=col, line_w=1.6,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        panel.adjustments[0] = 0.08
        rect(s, x, y, bw, 0.62, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y + 0.32, bw, 0.30, fill=col)
        _, tf = textbox(s, x, y + 0.02, bw, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, f"{i+1}.  {t}", size=16, color=WHITE, bold=True, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, first=True, space_after=0)
        _, tf = textbox(s, x + 0.2, y + 0.78, bw - 0.4, h - 0.9,
                        anchor=MSO_ANCHOR.MIDDLE)
        for j, ln in enumerate(lines):
            para(tf, ln, size=12.5, color=GRAY, align=PP_ALIGN.CENTER,
                 first=(j == 0), space_after=2, line_spacing=1.05)
        if i < n_steps - 1:
            ar = rect(s, x + bw + 0.04, y + h / 2 - 0.22, gap - 0.08, 0.44,
                      fill=LEAF, shape=MSO_SHAPE.RIGHT_ARROW)
    # loop-back caption
    lb = rect(s, MX, y + h + 0.5, CW, 0.62, fill=PANEL,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lb.adjustments[0] = 0.5
    _, tf = textbox(s, MX + 0.3, y + h + 0.5, CW - 0.6, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    runs(p, [("↻  Continuous loop at the field edge  —  ", 14, GREEN, True, HEAD_FONT),
             ("50 Hz real-time safety on the ESP32, ~vision cadence on the Pi, "
              "telemetry streamed live.", 13.5, GRAY, False)])
    fade(s)


def s_sensing(n):
    s = slide()
    header(s, "Precision Soil Sensing", "Know the soil before you feed it", n)
    bullets(s, MX, CT + 0.05, 6.7, 4.0, [
        ("7-in-1 RS485 probe", "N, P, K, pH, EC, moisture & soil temp over Modbus RTU"),
        ("Calibrated & robust", "16× oversampled ADC, multi-point moisture curve, temp-compensated TDS"),
        ("Geo-tagged readings", "every sample carries a sub-meter GPS + vision-fused position"),
        ("Prescription maps", "IDW / kriging interpolation → variable-rate CSV & GeoJSON"),
        ("Anomaly alerts", "rolling-window stream flags out-of-range nutrients in real time"),
    ], size=14.5)
    # right panel: the 7 parameters as chips + note
    px = MX + 7.0
    pw = CW - 7.0
    panel = rect(s, px, CT + 0.05, pw, 4.35, fill=PANEL, line=MIST, line_w=1,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    panel.adjustments[0] = 0.05
    _, tf = textbox(s, px + 0.3, CT + 0.28, pw - 0.6, 0.4)
    para(tf, "WHAT THE PROBE REPORTS", size=12, color=GREEN, bold=True,
         font=HEAD_FONT, first=True, space_after=0)
    params = ["Nitrogen (N)", "Phosphorus (P)", "Potassium (K)", "pH",
              "Conductivity (EC)", "Moisture", "Soil temp"]
    cx = px + 0.3
    cyy = CT + 0.78
    for i, pchip in enumerate(params):
        w = 0.135 * len(pchip) + 0.34
        if cx + w > px + pw - 0.25:
            cx = px + 0.3
            cyy += 0.52
        cx = chip(s, cx, cyy, pchip, color=GREEN, size=11)
    _, tf = textbox(s, px + 0.3, cyy + 0.7, pw - 0.6, 1.4)
    para(tf, "Rover halts and freezes the drive while the probe is in the soil "
             "(EVT_DOSING), then micro-doses only what each spot needs.",
         size=12.5, color=GRAY, first=True, line_spacing=1.15)
    fade(s)


def s_ai(n):
    s = slide()
    header(s, "On-board AI Vision", "The rover sees what it treats", n, accent=SKY)
    bullets(s, MX, CT + 0.05, 6.2, 4.2, [
        ("Three models, one camera", "obstacle · weed · crop-disease, run on captured frames"),
        ("YOLOv8n detection", "trained on DeepWeeds (9 species) & a custom obstacle set"),
        ("Disease classifier", "MobileNetV2 on PlantVillage — 38 classes across 14 crops"),
        ("Coral Edge TPU", "~10× faster inference than CPU for real-time safety"),
        ("Fails safe", "no model? detection degrades gracefully — it never crashes"),
    ], size=14, mcolor=SKY)
    place_image(s, "ai_performance.png", MX + 6.6, CT + 0.15, CW - 6.6, 3.4)
    # KPI row
    ky = 5.55
    kw = (CW - 3 * 0.3) / 4
    kpi(s, MX, ky, kw, 1.2, "3", "vision models", SKY)
    kpi(s, MX + (kw + 0.3), ky, kw, 1.2, "38", "disease classes", GREEN)
    kpi(s, MX + 2 * (kw + 0.3), ky, kw, 1.2, "~30", "FPS on Coral", AMBER)
    kpi(s, MX + 3 * (kw + 0.3), ky, kw, 1.2, "≥15", "FPS safety floor", PLUM)
    fade(s)


def s_action(n):
    s = slide()
    header(s, "Targeted Action", "Treat the plant, not the field", n, accent=AMBER)
    bullets(s, MX, CT + 0.05, 6.6, 4.2, [
        ("Sequential micro-dosing", "pre-soak → probe extend → dose → retract, never two relays at once"),
        ("Vision-aimed spot-spray", "pan/tilt nozzle points at each detected weed (FC-01)"),
        ("Modular attachments", "grass-cutter, seed-sower, weeder & sprayer swap in"),
        ("Drive-frozen while dosing", "the rover cannot move while the probe is in the soil"),
        ("Less chemical, less cost", "spray only where weeds are — not the whole bed"),
    ], size=14, mcolor=AMBER)
    place_image(s, "savings.png", MX + 6.9, CT + 0.15, CW - 6.9, 4.3)
    fade(s)


def s_nav(n):
    s = slide()
    header(s, "Autonomous Navigation", "Straight rows, repeatable coverage", n, accent=PLUM)
    bullets(s, MX, CT + 0.05, CW, 2.7, [
        ("Affordable GPS, sharpened", "Neo-6M + SBAS/GAGAN + stationary averaging → sub-meter fixes"),
        ("Sensor fusion (EKF)", "GPS + wheel odometry + IMU yaw → smooth pose between fixes"),
        ("Vision plant geo-tagging", "camera geometry localises each plant to ~10–20 cm"),
        ("Boustrophedon coverage", "snake path-planner with cross-track guidance covers the whole field"),
        ("Closed-loop velocity PID", "encoder feedback → straight rows & repeatable seed spacing"),
    ], size=14, mcolor=PLUM)
    ky = 5.5
    kw = (CW - 2 * 0.3) / 3
    kpi(s, MX, ky, kw, 1.25, "Sub-meter", "absolute GPS fix", PLUM)
    kpi(s, MX + (kw + 0.3), ky, kw, 1.25, "~10–20 cm", "per-plant (vision)", GREEN)
    kpi(s, MX + 2 * (kw + 0.3), ky, kw, 1.25, "Full-field", "snake coverage", SKY)
    fade(s)


def s_safety(n):
    s = slide()
    header(s, "Safety & Security", "Engineered to fail safe", n, accent=GREEN)
    place_image(s, "security.png", MX, CT + 0.1, 6.5, 4.4, frame=False,
                shadow=False)
    bullets(s, MX + 6.8, CT + 0.15, CW - 6.8, 4.4, [
        ("Hardware E-stop", "cuts everything via the ESP32 EN pin — no software in the loop"),
        ("Dead-man heartbeat", "drive halts if no signed command arrives in 1.5 s"),
        ("Thermal guardian", "pack NTC + CPU + ambient → throttle, stop or shutdown (FC-02)"),
        ("Stall detection", "motor current-sense stops a jammed wheel or cutter"),
        ("Fail-safe boot", "relays forced OFF first; watchdog resets a hung task"),
    ], size=13.5, mcolor=GREEN)
    fade(s)


def s_field(n):
    s = slide()
    header(s, "Field-Ready Engineering", "Ten real-world challenges, solved", n, accent=AMBER)
    rows = [
        ("ID", "Field challenge", "Engineered solution", "Status"),
        ("FC-01", "Weeds grow above the pathway", "Vision-guided pan/tilt aimed spray", "Code-ready"),
        ("FC-02", "Field heat → battery fire risk", "Thermal guardian + pack NTC + LiFePO4 option", "Code-ready"),
        ("FC-03", "Actuator retract type unknown", "Dual compile-flag branches (spring / DC)", "Code-ready"),
        ("FC-04", "Per-plant GPS precision", "SBAS + averaging + EKF + vision tag", "Implemented"),
        ("FC-05", "Analog sensors over-volt the ADC", "Sensors powered from the 3.3 V rail", "Implemented"),
        ("FC-06", "Relays twitch ON at boot", "10 kΩ pull-ups + firmware fail-safe OFF", "Implemented"),
        ("FC-07", "Comms link drops mid-drive", "HMAC heartbeat dead-man halt", "Implemented"),
        ("FC-08", "Dust, dew & spray on electronics", "Conformal coat + glands + gaskets", "Implemented"),
        ("FC-09", "Wheel / cutter jam (stall)", "ADS1115 current-sense stall stop", "Implemented"),
        ("FC-10", "Motor driver under-spec (L298N)", "2× BTS7960 MOSFET drivers", "Code-ready"),
    ]
    x, y, w = MX, CT + 0.02, CW
    h = 4.9
    tbl_shape = s.shapes.add_table(len(rows), 4, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    widths = [1.15, 4.3, 4.55, 1.5]
    for j, wd in enumerate(widths):
        tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(h / len(rows))
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C(FOREST)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C(WHITE if i % 2 else PANEL)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in (0, 3) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = HEAD_FONT if (i == 0 or j == 0) else BODY_FONT
            r.font.size = Pt(11 if i == 0 else 10.5)
            r.font.bold = (i == 0 or j == 0 or j == 3)
            if i == 0:
                r.font.color.rgb = C(WHITE)
            elif j == 0:
                r.font.color.rgb = C(GREEN)
            elif j == 3:
                r.font.color.rgb = C(GREEN if val == "Implemented" else AMBER)
            else:
                r.font.color.rgb = C(INK)
    fade(s)


def s_mechanical(n):
    s = slide()
    header(s, "Mechanical Design", "A stable, service-friendly platform", n)
    place_image(s, "chassis_layout.png", MX, CT, 8.4, 4.8,
                caption="To-scale top view rendered from the verified component coordinates.")
    bullets(s, MX + 8.7, CT + 0.25, CW - 8.7, 4.4, [
        ("Double-decker deck", "320 × 450 mm, two acrylic layers"),
        ("Low centre of gravity", "battery + tank sit low & centred"),
        ("Fluids below electronics", "leaks drain away from the Pi"),
        ("Heat kept apart", "motor drivers at the rear by the fan"),
        ("Field-hardened", "grommets, glands & conformal coat"),
    ], size=13, mcolor=GREEN)
    fade(s)


def s_wiring(n):
    s = slide()
    header(s, "Engineering Depth", "Every wire accounted for", n, accent=SKY)
    _, tf = textbox(s, MX, CT - 0.02, CW, 0.5)
    para(tf, "A fully verified electrical design — 110 BOM components + 15 gap-audit "
             "items, every bus, divider and safety part placed.", size=13.5,
         color=GRAY, first=True, line_spacing=1.1)
    place_image(s, "wiring.png", MX, 2.42, CW, 4.45,
                caption="Consolidated v2 wiring schematic — red = safety-critical paths.")
    fade(s)


def s_drive(n):
    s = slide()
    header(s, "Drive System", "MOSFET power for real soil loads", n, accent=SKY)
    place_image(s, "drive_schematic.png", MX, CT, 8.5, 4.85,
                caption="Pin-level 2× BTS7960 (IBT-2) drive schematic.")
    bullets(s, MX + 8.8, CT + 0.25, CW - 8.8, 4.4, [
        ("2× BTS7960 (IBT-2)", "replaced the L298N (FC-10)"),
        ("~43 A peak per side", "handles soil-loaded gear motors"),
        ("Dual-PWM control", "forward / reverse / coast"),
        ("1N5819 flyback", "clamps back-EMF on every motor"),
        ("Current-sense ready", "IS output → stall detection"),
    ], size=13, mcolor=SKY)
    fade(s)


def s_stack(n):
    s = slide()
    header(s, "Technology Stack", "Proven, open, off-the-shelf", n, accent=PLUM)
    groups = [
        ("FIRMWARE", GREEN, ["ESP32", "FreeRTOS", "PlatformIO", "C++", "mbedTLS"]),
        ("VISION / AI", SKY, ["Python", "OpenCV", "YOLOv8n", "TFLite", "Coral TPU"]),
        ("CONNECTIVITY", AMBER, ["MQTT / Mosquitto", "TLS", "LoRa", "PySerial", "HMAC link"]),
        ("DATA / UI", PLUM, ["Pathway", "Streamlit", "Folium", "InfluxDB", "WebSocket UI"]),
        ("ROBOTICS", "2E7D5B", ["ROS 2", "EKF", "Velocity PID", "Boustrophedon", "Simulator"]),
    ]
    y = CT + 0.1
    for label, col, items in groups:
        _, tf = textbox(s, MX, y + 0.02, 2.1, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, label, size=12.5, color=col, bold=True, font=HEAD_FONT, first=True,
             space_after=0)
        cx = MX + 2.2
        for it in items:
            cx = chip(s, cx, y + 0.06, it, color=col, size=12)
        y += 0.92
    fade(s)


def s_software(n):
    s = slide()
    header(s, "Software & Operations", "Built like a product, not a demo", n)
    cards = [
        ("Live dashboard", SKY, ["Streamlit metrics, NPK", "trends, GPS map & alerts"]),
        ("Phone remote control", GREEN, ["Stdlib web UI + WebSocket", "joystick & mission control"]),
        ("ROS 2 bridge", PLUM, ["Nav2 / RViz-ready topics", "& mission services"]),
        ("Black-box recorder", AMBER, ["Every command & event", "logged for replay"]),
        ("Secure OTA + sim", "2E7D5B", ["Password OTA updates &", "a full hardware-free sim"]),
        ("CI-verified", RED, ["C++ tests, py-compile,", "lint & notebook checks"]),
    ]
    w3 = (CW - 2 * 0.32) / 3
    ch = 1.9
    for i, (t, col, lines) in enumerate(cards):
        cx = MX + (i % 3) * (w3 + 0.32)
        cyy = CT + 0.15 + (i // 3) * (ch + 0.3)
        card(s, cx, cyy, w3, ch, t, lines, accent=col, icon="●",
             title_size=14, body_size=12)
    fade(s)


def s_cost(n):
    s = slide()
    header(s, "Affordable & Modular", "Start small, scale to full autonomy", n, accent=AMBER)
    place_image(s, "cost_tiers.png", MX, CT + 0.1, 7.4, 4.2, frame=True)
    bullets(s, MX + 7.7, CT + 0.3, CW - 7.7, 4.2, [
        ("Off-the-shelf parts", "no exotic hardware — all India-market components"),
        ("Three build tiers", "demo → navigation → full AI, add capability as you grow"),
        ("Software is free", "open-source stack end-to-end"),
        ("Repairable in the field", "modular decks, labelled harness, spares are cheap"),
    ], size=13.5, mcolor=AMBER)
    fade(s)


def s_impact(n):
    s = slide()
    header(s, "Impact", "Cheaper inputs, healthier soil, better data", n)
    ky = CT + 0.15
    kw = (CW - 2 * 0.3) / 3
    kpi(s, MX, ky, kw, 1.5, "↓ up to ~75%", "chemical use (spot-spray)", GREEN)
    kpi(s, MX + (kw + 0.3), ky, kw, 1.5, "Per-plant", "data & health history", SKY)
    kpi(s, MX + 2 * (kw + 0.3), ky, kw, 1.5, "Less labour", "autonomous field passes", AMBER)
    bullets(s, MX, ky + 1.95, CW, 2.4, [
        ("For farmers", "lower input bills, earlier problem detection, higher effective yield"),
        ("For the soil", "right-rate nutrients protect long-term fertility & groundwater"),
        ("For scale", "fleet-ready rover IDs & topics — one dashboard, many rovers"),
        ("Aligned with sustainability", "precision agriculture cuts waste and runoff"),
    ], size=14.5)
    fade(s)


def s_roadmap(n):
    s = slide()
    header(s, "Roadmap", "From working scaffold to field fleet", n, accent=SKY)
    milestones = [
        ("NOW", GREEN, "Working dual-controller scaffold", "sim + CI green, full firmware & Pi stack"),
        ("NEXT", SKY, "Field trials + trained models", "deploy Edge-TPU models, calibrate in-field"),
        ("THEN", AMBER, "RTK cm-GPS + crop health", "ZED-F9P + NDVI / multispectral imaging"),
        ("FUTURE", PLUM, "SLAM · Jetson · fleet", "on-board mapping & multi-rover operation"),
    ]
    y = 3.55
    rect(s, MX + 0.2, y, CW - 0.4, 0.03, fill=MIST)
    n_m = len(milestones)
    step = (CW - 0.4) / (n_m - 1)
    for i, (tag, col, title, sub) in enumerate(milestones):
        cx = MX + 0.2 + i * step
        dot = rect(s, cx - 0.16, y - 0.16, 0.34, 0.34, fill=col, shape=MSO_SHAPE.OVAL,
                   shadow=True)
        above = (i % 2 == 0)
        bx = cx - 1.35
        if i == 0:
            bx = cx - 0.2
        if i == n_m - 1:
            bx = cx - 2.5
        by = y - 1.85 if above else y + 0.45
        chip(s, bx if bx > MX else MX, by, tag, color=col, w=1.1, size=12)
        _, tf = textbox(s, (bx if bx > MX else MX), by + 0.44, 2.9, 1.2)
        para(tf, title, size=13.5, color=INK, bold=True, font=HEAD_FONT, first=True,
             space_after=2, line_spacing=1.0)
        para(tf, sub, size=11.5, color=GRAY, space_after=0, line_spacing=1.05)
    fade(s)


def s_team(n):
    s = slide()
    header(s, "The Team", "Built by first-year IIT Bombay engineers", n)
    members = [
        ("Kapadiya Hitanshu Mukeshbhai", "25B2228", GREEN),
        ("Vivek Kumar Gupta", "25B2269", SKY),
        ("Shreyash Wagh", "25B2227", AMBER),
    ]
    w3 = (CW - 2 * 0.4) / 3
    y = CT + 0.35
    h = 3.0
    for i, (name, roll, col) in enumerate(members):
        x = MX + i * (w3 + 0.4)
        panel = rect(s, x, y, w3, h, fill=WHITE, line=MIST, line_w=1,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        panel.adjustments[0] = 0.05
        rect(s, x, y, w3, 0.7, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y + 0.4, w3, 0.3, fill=col)
        av = rect(s, x + w3 / 2 - 0.55, y + 0.5, 1.1, 1.1, fill=col,
                  shape=MSO_SHAPE.OVAL, shadow=True)
        _, tf = textbox(s, x + w3 / 2 - 0.55, y + 0.5, 1.1, 1.1, anchor=MSO_ANCHOR.MIDDLE)
        initials = "".join([p[0] for p in name.split()[:2]]).upper()
        para(tf, initials, size=30, color=WHITE, bold=True, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, first=True, space_after=0)
        _, tf = textbox(s, x + 0.2, y + 1.75, w3 - 0.4, 0.8, anchor=MSO_ANCHOR.TOP)
        para(tf, name, size=14.5, color=INK, bold=True, font=HEAD_FONT,
             align=PP_ALIGN.CENTER, first=True, space_after=3, line_spacing=1.0)
        para(tf, f"Roll {roll}", size=12.5, color=col, bold=True,
             align=PP_ALIGN.CENTER, space_after=0)
    _, tf = textbox(s, MX, y + h + 0.35, CW, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "B.Tech Mechanical Engineering  ·  Indian Institute of Technology Bombay  ·  Class of 2029",
         size=14, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)
    fade(s)


def s_thanks(n):
    s = slide()
    bg = rect(s, 0, 0, 13.333, 7.5, fill=FOREST)
    grad(bg, FOREST, FOREST2, angle=55)
    rect(s, 0, 0, 0.3, 7.5, fill=LEAF)
    c1 = rect(s, 9.9, 3.0, 5.6, 5.6, fill=GREEN, shape=MSO_SHAPE.OVAL)
    _translucent(c1, 18000)
    c2 = rect(s, -1.6, -1.6, 3.6, 3.6, fill=LEAF, shape=MSO_SHAPE.OVAL)
    _translucent(c2, 12000)
    brand_wordmark(s, MX, 0.7, h=0.5, on_dark=True)
    _, tf = textbox(s, MX, 2.5, 11.0, 1.3)
    para(tf, "Thank you.", size=58, color=WHITE, bold=True, font=HEAD_FONT,
         first=True, space_after=0)
    _, tf = textbox(s, MX, 3.85, 11.0, 0.7)
    para(tf, "We'd love your questions — and to show you a live run.", size=20,
         color="D9EAD9", first=True, space_after=0)
    chip(s, MX, 4.95, "Groww × IIT Bombay  ·  Final Evaluation", color=AMBER,
         tcolor=INK, w=4.6, size=12.5)
    rect(s, MX, 5.75, 11.5, 0.02, fill="27503C")
    _, tf = textbox(s, MX, 5.9, 11.6, 1.0)
    para(tf, "Kapadiya Hitanshu Mukeshbhai · Vivek Kumar Gupta · Shreyash Wagh",
         size=13, color=WHITE, bold=True, first=True, space_after=4)
    para(tf, "25b2228@iitb.ac.in   ·   25b2269@iitb.ac.in   ·   25b2227@iitb.ac.in",
         size=12, color="9FC2AC", space_after=0)
    fade(s)


# =====================================================================
# ASSEMBLE
# =====================================================================
def main():
    s_title()
    s_problem(2)
    s_solution(3)
    section_divider(1, "Part One", "How AgriRover Works",
                    "Sense the soil, see the crop, act with precision.")
    s_architecture(5)
    s_loop(6)
    s_sensing(7)
    s_ai(8)
    s_action(9)
    s_nav(10)
    section_divider(2, "Part Two", "Under the Hood",
                    "The engineering that makes it field-ready.")
    s_safety(12)
    s_field(13)
    s_mechanical(14)
    s_wiring(15)
    s_drive(16)
    s_stack(17)
    s_software(18)
    section_divider(3, "Part Three", "Impact & Road Ahead",
                    "Affordable today, scalable tomorrow.")
    s_cost(20)
    s_impact(21)
    s_roadmap(22)
    s_team(23)
    s_thanks(24)

    prs.save(OUT)
    print("Saved:", OUT)
    print("Slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
